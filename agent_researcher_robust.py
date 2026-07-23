import sys
import asyncio

# --- Windows asyncio fix -----------------------------------------------------
# On Windows, Python's default ProactorEventLoop is incompatible with how the
# `websockets` library (used under the hood by Streamlit's dev server) handles
# ping/pong keepalives. During long-running operations — like this app's long
# code-generation streams — the mismatch surfaces as:
#   "keepalive ping failed... AssertionError: waiter is None or waiter.cancelled()"
# and kills the websocket connection to the browser tab mid-response.
# Forcing the SelectorEventLoop policy (the standard workaround for this bug)
# must happen before Streamlit/Tornado create their event loop, hence this is
# the very first thing the file does.
if sys.platform == "win32":
    asyncio.set_event_loop_policy(asyncio.WindowsSelectorEventLoopPolicy())

import streamlit as st
import os
import io
import re
import json
import queue
import threading
import base64
import zipfile
import subprocess
import platform
import shutil
import requests
from datetime import datetime
from crewai import Agent, Task, Crew, LLM

from streamlit.runtime.scriptrunner import add_script_run_ctx, get_script_run_ctx

import ssl
ssl._create_default_https_context = ssl._create_unverified_context
import urllib3
urllib3.disable_warnings(urllib3.exceptions.InsecureRequestWarning)

import arxiv
from pypdf import PdfReader
import shlex
import html
import hashlib
import secrets

BASE_DIR = os.path.dirname(os.path.abspath(__file__))


# =============================================================================
# USER ACCOUNTS — lightweight username/email + password auth for a shared
# workstation deployment. Backed by a single JSON file (fine at the scale of a
# workstation with a handful of users; not meant to replace a real auth system
# at larger scale). Each account maps to its own private data folder under
# /user_data/<username>/, isolating one user's papers/sessions/memory from
# another's — see the reassignment of *_DIR constants right after the login
# gate, further down this file.
# =============================================================================
USERS_DIR = os.path.join(BASE_DIR, "user_accounts")
USERS_FILE = os.path.join(USERS_DIR, "users.json")
_USERS_LOCK = threading.Lock()


def _load_users():
    try:
        with open(USERS_FILE, "r", encoding="utf-8") as f:
            return json.load(f)
    except Exception:
        return {"by_username": {}, "email_to_username": {}}


def _save_users(data):
    os.makedirs(USERS_DIR, exist_ok=True)
    _atomic_write_json(USERS_FILE, data)


def _hash_password(password, salt_hex=None):
    salt = bytes.fromhex(salt_hex) if salt_hex else secrets.token_bytes(16)
    digest = hashlib.pbkdf2_hmac("sha256", password.encode("utf-8"), salt, 200_000)
    return salt.hex(), digest.hex()


def _safe_username_for_path(username):
    """Turns a username into a filesystem-safe folder name for that user's private data."""
    safe = re.sub(r"[^A-Za-z0-9_-]+", "_", (username or "").strip().lower()) or "user"
    return safe


def _atomic_write_json(path, payload):
    """Writes JSON atomically: to a temp file first, then renamed into place. This means a
    crash or serialization error mid-write (e.g. an unexpected non-JSON-serializable field
    slipping into the payload) can never leave a truncated/corrupt file behind — if
    json.dump() raises, the exception propagates before the rename, so the previous
    on-disk version (if any) is untouched. Every project/session save function below uses
    this instead of writing directly to the real path."""
    tmp_path = f"{path}.tmp"
    with open(tmp_path, "w", encoding="utf-8") as f:
        json.dump(payload, f, ensure_ascii=False, indent=2)
    os.replace(tmp_path, path)


def create_account(username, email, password):
    """Returns (ok, error_message_or_None). The very first account ever created on this
    deployment is automatically made an admin, so there's always at least one admin
    without needing to hand-edit the user store."""
    username = (username or "").strip()
    email = (email or "").strip().lower()
    if not (3 <= len(username) <= 32) or not re.match(r"^[A-Za-z0-9_.-]+$", username):
        return False, "Username must be 3-32 characters: letters, numbers, underscore, dot, or hyphen only."
    if not re.match(r"^[^@\s]+@[^@\s]+\.[^@\s]+$", email):
        return False, "Please enter a valid email address."
    if len(password) < 8:
        return False, "Password must be at least 8 characters."

    with _USERS_LOCK:
        users = _load_users()
        uname_key = username.lower()
        if uname_key in users["by_username"]:
            return False, "That username is already taken."
        if email in users["email_to_username"]:
            return False, "An account with that email already exists."
        salt_hex, hash_hex = _hash_password(password)
        users["by_username"][uname_key] = {
            "username": username, "email": email, "salt": salt_hex, "hash": hash_hex,
            "created_at": datetime.now().strftime("%Y-%m-%d %H:%M"),
            "last_login": None, "is_admin": len(users["by_username"]) == 0,
        }
        users["email_to_username"][email] = uname_key
        _save_users(users)
    return True, None


def authenticate(identifier, password):
    """identifier can be a username or an email. Returns (username_or_None, error_message_or_None)."""
    identifier = (identifier or "").strip().lower()
    if not identifier or not password:
        return None, "Enter your username/email and password."
    users = _load_users()
    uname_key = users["email_to_username"].get(identifier, identifier)
    record = users["by_username"].get(uname_key)
    if not record:
        return None, "No account found with that username or email."
    _, check_hash = _hash_password(password, record["salt"])
    if not secrets.compare_digest(check_hash, record["hash"]):
        return None, "Incorrect password."
    return record["username"], None


def record_login(username):
    """Stamps last_login on the account and logs a 'login' usage event."""
    with _USERS_LOCK:
        users = _load_users()
        rec = users["by_username"].get(username.lower())
        if rec:
            rec["last_login"] = datetime.now().strftime("%Y-%m-%d %H:%M")
            _save_users(users)
    log_event(username, "login")


def is_admin_user(username):
    users = _load_users()
    rec = users["by_username"].get((username or "").lower())
    return bool(rec and rec.get("is_admin"))


def set_admin_flag(username, flag):
    with _USERS_LOCK:
        users = _load_users()
        rec = users["by_username"].get((username or "").lower())
        if rec:
            rec["is_admin"] = bool(flag)
            _save_users(users)


def list_all_users():
    """Returns the {username_lower: record} dict of every account, for the admin dashboard."""
    return _load_users()["by_username"]


def dir_size_bytes(path):
    total = 0
    try:
        for dirpath, _dirnames, filenames in os.walk(path):
            for fn in filenames:
                try:
                    total += os.path.getsize(os.path.join(dirpath, fn))
                except OSError:
                    pass
    except Exception:
        pass
    return total


# =============================================================================
# USAGE LOGGING — one JSON line per event (login, or a tool being used), so an
# admin account can see who's using this deployment and what they're using it
# for. Shared across all users (unlike each user's own private data folder) so
# an admin can see everyone's activity from one place.
# =============================================================================
EVENTS_FILE = os.path.join(BASE_DIR, "usage_logs", "events.jsonl")
_EVENTS_LOCK = threading.Lock()


def log_event(username, event, detail=""):
    """Best-effort: a logging failure should never break the app."""
    try:
        os.makedirs(os.path.dirname(EVENTS_FILE), exist_ok=True)
        line = json.dumps({
            "ts": datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
            "user": username, "event": event, "detail": detail,
        }, ensure_ascii=False)
        with _EVENTS_LOCK:
            with open(EVENTS_FILE, "a", encoding="utf-8") as f:
                f.write(line + "\n")
    except Exception:
        pass


def read_events(limit=5000):
    """Returns the most recent `limit` logged events, oldest first."""
    try:
        with open(EVENTS_FILE, "r", encoding="utf-8") as f:
            lines = f.readlines()[-limit:]
        events = []
        for ln in lines:
            try:
                events.append(json.loads(ln))
            except Exception:
                continue
        return events
    except Exception:
        return []


def format_claude_code_event(evt):
    """Turns one parsed `claude -p --output-format stream-json --verbose` event into a
    human-readable snippet, mirroring the ⏺ tool-call / ⎿ result style Claude Code itself
    prints in a real terminal, so it's clear what it's doing at each step (not just the
    final answer)."""
    etype = evt.get("type")
    out = ""
    if etype == "assistant":
        for block in ((evt.get("message") or {}).get("content") or []):
            btype = block.get("type")
            if btype == "text":
                txt = block.get("text", "")
                if txt.strip():
                    out += txt
            elif btype == "tool_use":
                name = block.get("name", "tool")
                inp = block.get("input", {}) or {}
                if "command" in inp:
                    detail = inp["command"]
                elif "file_path" in inp:
                    detail = inp["file_path"]
                elif "path" in inp:
                    detail = inp["path"]
                else:
                    detail = json.dumps(inp, ensure_ascii=False)[:200]
                out += f"\n⏺ {name}({detail})\n"
            elif btype == "thinking":
                out += "\n💭 (thinking…)\n"
    elif etype == "user":
        for block in ((evt.get("message") or {}).get("content") or []):
            if block.get("type") == "tool_result":
                content = block.get("content")
                if isinstance(content, list):
                    text = "".join(c.get("text", "") for c in content if isinstance(c, dict))
                else:
                    text = str(content or "")
                preview = text.strip().replace("\n", " ")[:300]
                if preview:
                    out += f"  ⎿ {preview}\n"
    elif etype == "result":
        if evt.get("subtype") == "success":
            extras = []
            cost = evt.get("total_cost_usd")
            dur = evt.get("duration_ms")
            if cost is not None:
                extras.append(f"${cost:.4f}")
            if dur is not None:
                extras.append(f"{dur / 1000:.1f}s")
            out += "\n✅ Done" + (f" ({', '.join(extras)})" if extras else "") + "\n"
        else:
            out += f"\n⚠️ {evt.get('result') or evt.get('subtype') or 'error'}\n"
    # "system" (session init) events carry nothing worth displaying — skipped.
    return out


def browse_for_folder(initial_dir=None):
    """Opens a native OS 'choose folder' dialog on the machine running this app and
    returns the chosen absolute path, or "" if the user cancelled or no display/tkinter
    is available (e.g. running on a headless server) — callers should fall back to the
    manual text input in that case."""
    try:
        import tkinter as tk
        from tkinter import filedialog
        root = tk.Tk()
        root.withdraw()
        root.attributes("-topmost", True)  # bring the dialog to the front of other windows
        chosen = filedialog.askdirectory(
            initialdir=initial_dir if initial_dir and os.path.isdir(initial_dir) else None,
            title="Choose project folder")
        root.destroy()
        return chosen or ""
    except Exception:
        return ""


def path_in_sandbox(path, sandbox_root):
    """True if `path` resolves to a location inside `sandbox_root` (or is the root
    itself). Resolves symlinks/'..'/relative segments before comparing, so it can't
    be tricked by path traversal — used to confine the Agent/Claude Code
    project-folder pickers to a signed-in user's own sandbox on a shared workstation."""
    try:
        real_path = os.path.realpath(os.path.abspath(path))
        real_root = os.path.realpath(os.path.abspath(sandbox_root))
        return real_path == real_root or real_path.startswith(real_root + os.sep)
    except Exception:
        return False


BROWSER_HEADERS = {
    "User-Agent": ("Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 "
                   "(KHTML, like Gecko) Chrome/125.0 Safari/537.36"),
    "Accept": "text/html,application/xhtml+xml;q=0.9,*/*;q=0.8",
    "Accept-Language": "en-US,en;q=0.9",
}


# =============================================================================
# PERSISTENT CHAT MEMORY — like ChatGPT's memory feature: the model notices
# durable facts about the user as they chat and recalls them in later sessions,
# even after the app restarts. Backed by a small JSON file per "profile" (so
# multiple people sharing one deployment each get their own memory), stored on
# disk rather than in st.session_state — session_state resets when a browser
# session ends, which is exactly what this is meant to survive.
# =============================================================================
MEMORY_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "chat_memory")


def _memory_path(profile):
    safe = re.sub(r"[^A-Za-z0-9_-]+", "_", (profile or "default").strip()) or "default"
    return os.path.join(MEMORY_DIR, f"{safe}.json")


def load_memory(profile):
    """Returns a list of {"fact": str, "ts": str} dicts, oldest first. Empty list if the
    profile has no memory file yet or it can't be read."""
    try:
        with open(_memory_path(profile), "r", encoding="utf-8") as f:
            data = json.load(f)
        return data if isinstance(data, list) else []
    except Exception:
        return []


def save_memory(profile, facts):
    try:
        os.makedirs(MEMORY_DIR, exist_ok=True)
        with open(_memory_path(profile), "w", encoding="utf-8") as f:
            json.dump(facts, f, ensure_ascii=False, indent=2)
        return True
    except Exception:
        return False


def add_memory_fact(profile, fact):
    """Append a new fact, skipping if it's a near-duplicate (case-insensitive exact match) of
    something already stored."""
    facts = load_memory(profile)
    norm = fact.strip().lower()
    if any(f.get("fact", "").strip().lower() == norm for f in facts):
        return facts
    facts.append({"fact": fact.strip(), "ts": datetime.now().strftime("%Y-%m-%d %H:%M")})
    save_memory(profile, facts)
    return facts


def delete_memory_fact(profile, index):
    facts = load_memory(profile)
    if 0 <= index < len(facts):
        facts.pop(index)
        save_memory(profile, facts)
    return facts


def clear_memory(profile):
    save_memory(profile, [])


def get_user_memory_context(username, max_facts=30):
    """Formats a user's stored memory facts into a block ready to drop into a system
    prompt. Returns "" if there's nothing stored yet, so callers can skip injecting it."""
    facts = load_memory(username)
    if not facts:
        return ""
    facts = facts[-max_facts:]  # most recently learned facts, in case the list has grown large
    lines = "\n".join(f"- {f.get('fact', '')}" for f in facts if f.get("fact"))
    if not lines:
        return ""
    return (
        "KNOWN CONTEXT ABOUT THIS USER (learned automatically from their past activity in "
        "this app — their research interests, the topics they've written about, languages "
        "or tools they use, etc.). Use it to personalize your response where it's actually "
        "relevant — e.g. defaulting to their usual programming language, or connecting a new "
        "request to a topic they've worked on before — but don't force a mention of it in "
        "when it doesn't apply to the current request:\n" + lines
    )


MEMORY_EXTRACT_SYSTEM = (
    "You are deciding whether a chat exchange revealed anything worth remembering long-term "
    "about the user, for future conversations — like ChatGPT's memory feature. Only durable "
    "facts belong here: their name, role, ongoing projects, stated preferences, or recurring "
    "context. NOT one-off details, small talk, or anything already listed as already "
    "remembered below.\n\n"
    "Respond in EXACTLY this two-line format, FACT on the FIRST line always:\n"
    "FACT: <one short factual sentence to remember, or NONE if nothing new and durable came up>\n"
    "REASON: <one short reason for your decision>"
)


def extract_memory_fact(user_message, assistant_response, existing_facts, model,
                        api_base=None, api_key=None, num_ctx=None):
    """Ask the LLM whether this exchange contained a new durable fact worth remembering.
    Returns the fact string, or None if there's nothing new (or the call fails)."""
    existing_block = "\n".join(f"- {f}" for f in existing_facts) or "(none yet)"
    user_msg = (
        f"Already remembered about the user:\n{existing_block}\n\n"
        f"User said: {user_message}\n\nAssistant replied: {assistant_response[:800]}\n\n"
        f"Decide now."
    )
    msgs = [{"role": "system", "content": MEMORY_EXTRACT_SYSTEM},
            {"role": "user", "content": user_msg}]
    try:
        chunks = list(stream_chat(msgs, model, api_base=api_base, api_key=api_key,
                                  temperature=0.0, num_ctx=num_ctx, max_tokens=150))
        raw = "".join(chunks)
    except Exception:
        return None
    m_fact = re.search(r"FACT:\s*(.+)", raw, re.I)
    if not m_fact:
        return None
    fact = m_fact.group(1).strip().splitlines()[0].strip()
    if not fact or fact.upper().startswith("NONE"):
        return None
    return fact[:300]


# =============================================================================
# WRITE PAPER PROJECTS — save/load the entire pipeline state to disk (not just
# st.session_state, which is lost when the browser session ends) so a paper can
# be started, saved, closed, and picked back up later — including the gathered
# sources, novelty report, approved idea, outline, and the written draft with
# its figures. Multiple named projects can exist side by side.
# =============================================================================
WP_PROJECTS_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "write_paper_projects")
WP_PDFS_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "write_paper_project_pdfs")

# Every session_state key that makes up a Write Paper project's state, in the order the
# pipeline produces them. Kept as one list so save/load/new-project logic can't drift apart.
WP_STATE_KEYS = [
    "wp_idea", "wp_web", "wp_web_full", "wp_max_rounds", "wp_max_chars_per_paper",
    "wp_candidates", "wp_removed", "wp_sources", "wp_novelty_report",
    "wp_approved_idea", "wp_title", "wp_title_input", "wp_sections", "wp_paper",
]


def _wp_project_path(name):
    safe = re.sub(r"[^A-Za-z0-9_-]+", "_", (name or "").strip()) or "untitled"
    return os.path.join(WP_PROJECTS_DIR, f"{safe}.json")


def list_wp_projects():
    try:
        os.makedirs(WP_PROJECTS_DIR, exist_ok=True)
        return sorted(f[:-5] for f in os.listdir(WP_PROJECTS_DIR) if f.endswith(".json"))
    except Exception:
        return []


def _encode_figs_for_storage(figs):
    """Figures hold raw PNG bytes, which JSON can't store directly — base64-encode them."""
    out = []
    for f in (figs or []):
        out.append({
            "title": f.get("title", ""), "caption": f.get("caption", ""),
            "aspect": f.get("aspect", 0.6),
            "png_b64": base64.b64encode(f["png"]).decode("ascii") if f.get("png") else "",
        })
    return out


def _decode_figs_from_storage(figs_enc):
    out = []
    for f in (figs_enc or []):
        png = base64.b64decode(f["png_b64"]) if f.get("png_b64") else b""
        out.append({"title": f.get("title", ""), "caption": f.get("caption", ""),
                    "aspect": f.get("aspect", 0.6), "png": png})
    return out


def save_wp_project(name, state):
    """`state` is a dict keyed by (a subset of) WP_STATE_KEYS, pulled from session_state.
    Returns (success, error_message_or_None)."""
    try:
        os.makedirs(WP_PROJECTS_DIR, exist_ok=True)
        payload = dict(state)
        if payload.get("wp_paper"):
            wp = dict(payload["wp_paper"])
            wp["figs"] = _encode_figs_for_storage(wp.get("figs"))
            payload["wp_paper"] = wp
        payload["_saved_at"] = datetime.now().strftime("%Y-%m-%d %H:%M")
        _atomic_write_json(_wp_project_path(name), payload)
        return True, None
    except Exception as e:
        return False, str(e)


def load_wp_project(name):
    """Returns (payload_dict_or_None, error_message_or_None). The payload's wp_paper (if any)
    has its figures already decoded back to raw PNG bytes, ready to drop into session_state."""
    try:
        with open(_wp_project_path(name), "r", encoding="utf-8") as f:
            payload = json.load(f)
        if payload.get("wp_paper"):
            wp = dict(payload["wp_paper"])
            wp["figs"] = _decode_figs_from_storage(wp.get("figs"))
            payload["wp_paper"] = wp
        return payload, None
    except Exception as e:
        return None, str(e)


def delete_wp_project(name):
    try:
        os.remove(_wp_project_path(name))
        return True
    except Exception:
        return False


# =============================================================================
# SURVEY PROJECTS — same save/load/delete pattern as Write Paper, so a survey's
# setup (topic, chosen sources, generated survey text) survives closing the app.
# =============================================================================
SURVEY_PROJECTS_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "survey_projects")
SURVEY_PDFS_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "survey_project_pdfs")

SURVEY_STATE_KEYS = ["sv_topic", "sv_candidates", "sv_removed", "survey_result"]


def _survey_project_path(name):
    safe = re.sub(r"[^A-Za-z0-9_-]+", "_", (name or "").strip()) or "untitled"
    return os.path.join(SURVEY_PROJECTS_DIR, f"{safe}.json")


def list_survey_projects():
    try:
        os.makedirs(SURVEY_PROJECTS_DIR, exist_ok=True)
        return sorted(f[:-5] for f in os.listdir(SURVEY_PROJECTS_DIR) if f.endswith(".json"))
    except Exception:
        return []


def _encode_survey_result_for_storage(sv):
    """survey_result holds raw PNG bytes (figs) and a raw zip (tex_zip) — JSON can't store
    either directly, so base64-encode them for storage."""
    if not sv:
        return sv
    sv = dict(sv)
    if sv.get("figs"):
        sv["figs"] = _encode_figs_for_storage(sv["figs"])
    if sv.get("tex_zip"):
        sv["tex_zip_b64"] = base64.b64encode(sv["tex_zip"]).decode("ascii")
        sv["tex_zip"] = None
    return sv


def _decode_survey_result_from_storage(sv):
    if not sv:
        return sv
    sv = dict(sv)
    if sv.get("figs"):
        sv["figs"] = _decode_figs_from_storage(sv["figs"])
    if sv.get("tex_zip_b64"):
        sv["tex_zip"] = base64.b64decode(sv["tex_zip_b64"])
        sv.pop("tex_zip_b64", None)
    return sv


def save_survey_project(name, state):
    try:
        os.makedirs(SURVEY_PROJECTS_DIR, exist_ok=True)
        payload = dict(state)
        if payload.get("survey_result"):
            payload["survey_result"] = _encode_survey_result_for_storage(payload["survey_result"])
        payload["_saved_at"] = datetime.now().strftime("%Y-%m-%d %H:%M")
        _atomic_write_json(_survey_project_path(name), payload)
        return True, None
    except Exception as e:
        return False, str(e)


def load_survey_project(name):
    try:
        with open(_survey_project_path(name), "r", encoding="utf-8") as f:
            payload = json.load(f)
        if payload.get("survey_result"):
            payload["survey_result"] = _decode_survey_result_from_storage(payload["survey_result"])
        return payload, None
    except Exception as e:
        return None, str(e)


def delete_survey_project(name):
    try:
        os.remove(_survey_project_path(name))
        return True
    except Exception:
        return False


# =============================================================================
# RESEARCH CREW PROJECTS — same pattern, for the Research Crew tab's query,
# chosen papers, and the crew's output.
# =============================================================================
CREW_PROJECTS_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "research_crew_projects")
CREW_PDFS_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "research_crew_project_pdfs")


def _crew_project_path(name):
    safe = re.sub(r"[^A-Za-z0-9_-]+", "_", (name or "").strip()) or "untitled"
    return os.path.join(CREW_PROJECTS_DIR, f"{safe}.json")


def list_crew_projects():
    try:
        os.makedirs(CREW_PROJECTS_DIR, exist_ok=True)
        return sorted(f[:-5] for f in os.listdir(CREW_PROJECTS_DIR) if f.endswith(".json"))
    except Exception:
        return []


def save_crew_project(name, state):
    try:
        os.makedirs(CREW_PROJECTS_DIR, exist_ok=True)
        payload = dict(state)
        payload["_saved_at"] = datetime.now().strftime("%Y-%m-%d %H:%M")
        _atomic_write_json(_crew_project_path(name), payload)
        return True, None
    except Exception as e:
        return False, str(e)


def load_crew_project(name):
    try:
        with open(_crew_project_path(name), "r", encoding="utf-8") as f:
            return json.load(f), None
    except Exception as e:
        return None, str(e)


def delete_crew_project(name):
    try:
        os.remove(_crew_project_path(name))
        return True
    except Exception:
        return False


def save_project_source_pdfs(pdfs_root, project_name, sources):
    """Downloads and saves the actual PDF (not just extracted text) for each source that
    has a pdf_url, into <pdfs_root>/<project_name>/ — so the papers used are kept
    alongside the project and can be browsed/downloaded later. Best-effort: sources with
    no pdf_url or that fail to fetch are silently skipped. Returns the list of filenames
    actually saved."""
    safe_proj = re.sub(r"[^A-Za-z0-9_-]+", "_", (project_name or "").strip()) or "untitled"
    dest = os.path.join(pdfs_root, safe_proj)
    os.makedirs(dest, exist_ok=True)
    saved = []
    for p in sources or []:
        url = p.get("pdf_url")
        if not url:
            continue
        try:
            pdf_bytes, err = fetch_pdf_bytes(url)
            if not pdf_bytes:
                continue
            fname = re.sub(r"[^A-Za-z0-9_ -]+", "_", (p.get("title") or "paper").strip())[:80] + ".pdf"
            with open(os.path.join(dest, fname), "wb") as f:
                f.write(pdf_bytes)
            saved.append(fname)
        except Exception:
            continue
    return saved


def list_project_pdfs(pdfs_root, project_name):
    safe_proj = re.sub(r"[^A-Za-z0-9_-]+", "_", (project_name or "").strip()) or "untitled"
    dest = os.path.join(pdfs_root, safe_proj)
    try:
        return sorted(f for f in os.listdir(dest) if f.lower().endswith(".pdf"))
    except Exception:
        return []


# =============================================================================
# CODING AGENT SESSIONS — save/load the agent's full conversation (history +
# transcript + project folder) to disk, so a coding session can be closed and
# resumed later with complete context intact, like Claude Code's session
# persistence. Plain JSON — no binary data to encode here, unlike Write Paper's
# figures.
# =============================================================================
CODING_SESSIONS_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "coding_agent_sessions")


def _coding_session_path(name):
    safe = re.sub(r"[^A-Za-z0-9_-]+", "_", (name or "").strip()) or "untitled"
    return os.path.join(CODING_SESSIONS_DIR, f"{safe}.json")


def list_coding_sessions():
    try:
        os.makedirs(CODING_SESSIONS_DIR, exist_ok=True)
        return sorted(f[:-5] for f in os.listdir(CODING_SESSIONS_DIR) if f.endswith(".json"))
    except Exception:
        return []


def save_coding_session(name, agent_state):
    try:
        os.makedirs(CODING_SESSIONS_DIR, exist_ok=True)
        payload = {
            "project_dir": agent_state.get("project_dir", ""),
            "history": agent_state.get("history", []),
            "transcript": agent_state.get("transcript", []),
            "steps": agent_state.get("steps", 0),
            "_saved_at": datetime.now().strftime("%Y-%m-%d %H:%M"),
        }
        _atomic_write_json(_coding_session_path(name), payload)
        return True, None
    except Exception as e:
        return False, str(e)


def load_coding_session(name):
    try:
        with open(_coding_session_path(name), "r", encoding="utf-8") as f:
            return json.load(f), None
    except Exception as e:
        return None, str(e)


def delete_coding_session(name):
    try:
        os.remove(_coding_session_path(name))
        return True
    except Exception:
        return False


def autoname_session(seed_text, existing_names, maxlen=40):
    """Turns the first bit of a task/prompt into a short session name, and disambiguates
    it against names that already exist (so two sessions never silently overwrite)."""
    seed_text = re.sub(r"\s+", " ", seed_text or "").strip()
    base = (seed_text[:maxlen].rstrip() + "…") if len(seed_text) > maxlen else (seed_text or "Untitled session")
    existing = set(existing_names)
    if base not in existing:
        return base
    n = 2
    while f"{base} ({n})" in existing:
        n += 1
    return f"{base} ({n})"


# =============================================================================
# CLAUDE CODE (NATIVE) SESSIONS — save/load the transcript from the "Claude Code
# (native CLI, any LLM)" mode to disk, so it can be closed and picked back up
# later. Note this only persists THIS APP's own display of the conversation
# (the turns shown in the UI) — Claude Code's own --continue history lives in
# its ~/.claude session store on disk and is separate from this.
# =============================================================================
CC_SESSIONS_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "claude_code_sessions")
CHAT_CONVOS_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "chat_conversations")


def _cc_session_path(name):
    safe = re.sub(r"[^A-Za-z0-9_-]+", "_", (name or "").strip()) or "untitled"
    return os.path.join(CC_SESSIONS_DIR, f"{safe}.json")


def list_cc_sessions():
    try:
        os.makedirs(CC_SESSIONS_DIR, exist_ok=True)
        return sorted(f[:-5] for f in os.listdir(CC_SESSIONS_DIR) if f.endswith(".json"))
    except Exception:
        return []


def save_cc_session(name, ccn_state, project_dir):
    try:
        os.makedirs(CC_SESSIONS_DIR, exist_ok=True)
        payload = {
            "project_dir": project_dir,
            "turns": ccn_state.get("turns", []),
            "_saved_at": datetime.now().strftime("%Y-%m-%d %H:%M"),
        }
        _atomic_write_json(_cc_session_path(name), payload)
        return True, None
    except Exception as e:
        return False, str(e)


def load_cc_session(name):
    try:
        with open(_cc_session_path(name), "r", encoding="utf-8") as f:
            return json.load(f), None
    except Exception as e:
        return None, str(e)


def delete_cc_session(name):
    try:
        os.remove(_cc_session_path(name))
        return True
    except Exception:
        return False


# =============================================================================
# CHAT CONVERSATION PERSISTENCE — save/load chat conversations to disk per-user,
# so chats survive browser closes and app restarts (like ChatGPT's sidebar).
# =============================================================================

def _chat_conv_path(cid):
    safe = re.sub(r"[^A-Za-z0-9_-]+", "_", (cid or "").strip()) or "untitled"
    return os.path.join(CHAT_CONVOS_DIR, f"{safe}.json")

_CHAT_INDEX_PATH = os.path.join(CHAT_CONVOS_DIR, "_index.json")


def _load_chat_index():
    """Returns the active conversation ID from the index file, or None."""
    try:
        with open(_CHAT_INDEX_PATH, "r", encoding="utf-8") as f:
            return json.load(f).get("active")
    except Exception:
        return None


def _save_chat_index(active_cid):
    """Persists which conversation is currently active."""
    try:
        _atomic_write_json(_CHAT_INDEX_PATH, {"active": active_cid})
    except Exception:
        pass


def list_chat_conversations():
    """Returns a sorted list of conversation IDs saved on disk for this user."""
    try:
        os.makedirs(CHAT_CONVOS_DIR, exist_ok=True)
        return sorted(
            f[:-5] for f in os.listdir(CHAT_CONVOS_DIR)
            if f.endswith(".json") and not f.startswith("_")
        )
    except Exception:
        return []


def save_chat_conversation(cid, conv_data):
    """Persists a single conversation to disk. conv_data = {title, messages}."""
    try:
        os.makedirs(CHAT_CONVOS_DIR, exist_ok=True)
        payload = {
            "title": conv_data.get("title", "New chat"),
            "messages": conv_data.get("messages", []),
            "_saved_at": datetime.now().strftime("%Y-%m-%d %H:%M"),
        }
        _atomic_write_json(_chat_conv_path(cid), payload)
        _save_chat_index(cid)
        return True
    except Exception:
        return False


def load_chat_conversation(cid):
    """Loads a single conversation from disk. Returns (conv_data, error)."""
    try:
        with open(_chat_conv_path(cid), "r", encoding="utf-8") as f:
            return json.load(f), None
    except Exception as e:
        return None, str(e)


def delete_chat_conversation(cid):
    """Deletes a conversation file from disk."""
    try:
        path = _chat_conv_path(cid)
        if os.path.isfile(path):
            os.remove(path)
        return True
    except Exception:
        return False


def load_all_chat_conversations():
    """Loads every saved conversation into a dict keyed by conversation ID.
    Returns (conversations_dict, active_cid_or_None)."""
    saved_ids = list_chat_conversations()
    convos = {}
    for cid in saved_ids:
        data, err = load_chat_conversation(cid)
        if err is None and data is not None:
            convos[cid] = {
                "title": data.get("title", "New chat"),
                "messages": data.get("messages", []),
            }
    active = _load_chat_index()
    if active not in convos:
        active = saved_ids[0] if saved_ids else None
    return convos, active



# =============================================================================
# Ollama & Google helpers
# =============================================================================
@st.cache_data(ttl=300, show_spinner=False)
def list_ollama_models(base_url, api_key=None):
    try:
        headers = {}
        if api_key:
            headers["Authorization"] = f"Bearer {api_key}"
        r = requests.get(base_url.rstrip("/") + "/api/tags", timeout=8, headers=headers)
        r.raise_for_status()
        return sorted(m["name"] for m in r.json().get("models", []))
    except Exception:
        return []


@st.cache_data(ttl=300, show_spinner=False)
def list_anthropic_compatible_models(base_url, api_key=None):
    """List available models from an Anthropic Messages-format endpoint.

    Tries multiple auth flavours, endpoint paths, AND parent URLs because custom
    proxies vary wildly:
    - DeepSeek: base = …/anthropic (translation layer, no /models there) but the
      parent domain exposes /models (OpenAI-style)
    - Anthropic official / Qwen Anthropic: x-api-key, GET /v1/models
    - OpenRouter / LiteLLM / local proxies: Bearer token, GET /v1/models or /models
    """
    if not base_url:
        return []
    base = base_url.rstrip("/")

    # Build candidate base URLs — the literal one AND the parent (strip last path
    # segment) so that e.g. https://api.deepseek.com/anthropic also probes
    # https://api.deepseek.com/v1/models and https://api.deepseek.com/models
    bases_to_try = [base]
    scheme_sep = base.find("://")
    path_start = base.find("/", scheme_sep + 3) if scheme_sep != -1 else -1
    if path_start != -1:
        parent = base[:path_start]
        if parent not in bases_to_try:
            bases_to_try.append(parent)

    # Auth headers — some proxies want x-api-key, others want Bearer
    auth_headers_list = []
    if api_key:
        auth_headers_list.append({"x-api-key": api_key, "anthropic-version": "2023-06-01"})
        auth_headers_list.append({"Authorization": f"Bearer {api_key}"})
        auth_headers_list.append({"Authorization": f"Bearer {api_key}", "anthropic-version": "2023-06-01"})
    else:
        auth_headers_list.append({"anthropic-version": "2023-06-01"})
        auth_headers_list.append({})

    for try_base in bases_to_try:
        for path in ("/v1/models", "/models"):
            for headers in auth_headers_list:
                try:
                    r = requests.get(try_base + path, timeout=10, headers=headers)
                    r.raise_for_status()
                    data = r.json().get("data", [])
                    models = sorted(m["id"] for m in data if m.get("id"))
                    if models:
                        return models
                except Exception:
                    continue
    return []


@st.cache_data(ttl=300, show_spinner=False)
def list_openai_compatible_models(base_url, api_key=None):
    """List available models from any OpenAI-compatible endpoint via GET {base}/models —
    works for Qwen/DashScope, Together, Groq, Fireworks, a local vLLM server, etc."""
    if not base_url:
        return []
    try:
        headers = {}
        if api_key:
            headers["Authorization"] = f"Bearer {api_key}"
        r = requests.get(base_url.rstrip("/") + "/models", timeout=10, headers=headers)
        r.raise_for_status()
        data = r.json().get("data", [])
        return sorted(m["id"] for m in data if m.get("id"))
    except Exception:
        return []


@st.cache_data(ttl=300, show_spinner=False)
def list_google_models(api_key):
    """Dynamically fetches models and ensures 'gemini-flash-lite-latest' is the default."""
    default_models = ["gemini-flash-lite-latest", "gemini-1.5-flash", "gemini-1.5-pro", "gemini-2.0-flash"]
    if not api_key:
        return default_models
        
    try:
        url = f"https://generativelanguage.googleapis.com/v1beta/models?key={api_key}"
        r = requests.get(url, timeout=10)
        r.raise_for_status()
        
        found_models = []
        for m in r.json().get("models", []):
            if "generateContent" in m.get("supportedGenerationMethods", []):
                name = m.get("name", "").replace("models/", "")
                # Avoid duplicates
                if name != "gemini-flash-lite-latest" and name not in found_models:
                    found_models.append(name)
        
        # Sort the others alphabetically
        found_models.sort()
        
        # Force 'gemini-flash-lite-latest' to be the first item
        return ["gemini-flash-lite-latest"] + found_models
    except Exception:
        return default_models


def default_model_index(models, provider="ollama"):
    """Prefer an Ollama Cloud model or a fast Google Flash model by default."""
    if provider == "google":
        for i, m in enumerate(models):
            if "flash" in m.lower() and "8b" not in m.lower():
                return i
        return 0
    else:
        for i, m in enumerate(models):
            if "cloud" in m.lower():
                return i
        for pref in ("qwen2.5:7b-instruct", "qwen2.5:7b", "llama3.1:8b", "llama3.2"):
            if pref in models:
                return models.index(pref)
        return 0


# =============================================================================
# MULTI-SOURCE SEARCH  (metadata only — fast; full text fetched later)
# Unified candidate schema: source,title,authors,published,pdf_url,
#                           landing_url,abstract,doi,full_text
# =============================================================================
def search_arxiv_meta(query, n=5, sort="relevance"):
    sort_map = {"relevance": arxiv.SortCriterion.Relevance,
                "newest": arxiv.SortCriterion.SubmittedDate}
    client = arxiv.Client(page_size=n, delay_seconds=3, num_retries=3)
    search = arxiv.Search(query=query, max_results=n,
                          sort_by=sort_map.get(sort, arxiv.SortCriterion.Relevance))
    out = []
    for r in client.results(search):
        out.append({
            "source": "arXiv",
            "title": r.title,
            "authors": ", ".join(a.name for a in r.authors[:6]),
            "published": r.published.strftime("%Y-%m-%d") if r.published else "n/a",
            "pdf_url": r.pdf_url,
            "landing_url": r.entry_id,
            "abstract": " ".join(r.summary.split()),
            "doi": r.doi or "",
            "full_text": "",
        })
    return out


def _parse_s2(data):
    out = []
    for p in data.get("data", []):
        authors = ", ".join(a.get("name", "") for a in (p.get("authors") or [])[:6])
        pdf = (p.get("openAccessPdf") or {}).get("url") or ""
        out.append({
            "source": "Semantic Scholar",
            "title": p.get("title") or "",
            "authors": authors,
            "published": str(p.get("year") or "n/a"),
            "pdf_url": pdf,
            "landing_url": p.get("url") or "",
            "abstract": p.get("abstract") or "(no abstract available)",
            "doi": (p.get("externalIds") or {}).get("DOI", ""),
            "full_text": "",
        })
    return out


def search_semantic_scholar(query, n=5):
    """Free API, no key. Covers ALL fields (medical, bio, social science, etc.).
    The free tier is ~1 request/second and returns 429 when exceeded, so we retry."""
    import time
    url = "https://api.semanticscholar.org/graph/v1/paper/search"
    params = {"query": query, "limit": n,
              "fields": "title,abstract,authors,year,openAccessPdf,externalIds,url,venue"}
    last_err = None
    for attempt in range(3):
        try:
            r = requests.get(url, params=params, timeout=25, headers=BROWSER_HEADERS)
            if r.status_code == 429:
                time.sleep(3 * (attempt + 1))   # back off and retry
                last_err = "rate limited (429)"
                continue
            r.raise_for_status()
            return _parse_s2(r.json())
        except Exception as e:
            last_err = str(e)
            time.sleep(2)
    raise RuntimeError(f"Semantic Scholar unavailable after retries ({last_err})")


def _openalex_abstract(inv):
    """OpenAlex returns abstracts as an inverted index {word: [positions]}; rebuild it."""
    if not inv:
        return "(no abstract available)"
    pos = {}
    for word, idxs in inv.items():
        for i in idxs:
            pos[i] = word
    return " ".join(pos[i] for i in sorted(pos)) or "(no abstract available)"


def _parse_openalex(data):
    out = []
    for w in data.get("results", []):
        auth = ", ".join((a.get("author") or {}).get("display_name", "")
                         for a in (w.get("authorships") or [])[:6])
        loc = w.get("primary_location") or {}
        src = loc.get("source") or {}
        boa = w.get("best_oa_location") or {}
        oa = w.get("open_access") or {}
        pdf = boa.get("pdf_url") or oa.get("oa_url") or loc.get("pdf_url") or ""
        doi = (w.get("doi") or "").replace("https://doi.org/", "")
        out.append({
            "source": "OpenAlex",
            "venue": src.get("display_name") or "",
            "publisher": src.get("host_organization_name") or "",
            "title": w.get("display_name") or w.get("title") or "",
            "authors": auth,
            "published": str(w.get("publication_year") or "n/a"),
            "pdf_url": pdf,
            "landing_url": (loc.get("landing_page_url")
                            or (("https://doi.org/" + doi) if doi else "")
                            or w.get("id", "")),
            "abstract": _openalex_abstract(w.get("abstract_inverted_index")),
            "doi": doi,
            "full_text": "",
        })
    return out


def search_openalex(query, n=5):
    """Free, no key. Indexes IEEE, Elsevier, Springer, ACM, Nature, etc.
    Returns title + abstract for all, and an open-access PDF link when one exists."""
    url = "https://api.openalex.org/works"
    params = {"search": query, "per-page": n,
              "mailto": "research-app@example.com"}   # polite pool (recommended by OpenAlex)
    r = requests.get(url, params=params, timeout=25, headers=BROWSER_HEADERS)
    r.raise_for_status()
    return _parse_openalex(r.json())


# --- Local / uploaded PDFs as a paper source -------------------------------
def _pdf_bytes_to_candidate(name, data, chars=12000):
    """Turn raw PDF bytes into a candidate dict with full text already extracted."""
    title = re.sub(r"\.pdf$", "", os.path.basename(name), flags=re.I)
    full_text, abstract = "", ""
    try:
        reader = PdfReader(io.BytesIO(data))
        raw = "\n".join((pg.extract_text() or "") for pg in reader.pages)
        full_text = " ".join(raw.split())
        abstract = (full_text[:600] + "…") if len(full_text) > 600 else full_text
        # try to use the first line as a nicer title
        first_line = next((ln.strip() for ln in raw.splitlines() if len(ln.strip()) > 15), "")
        if first_line:
            title = first_line[:160]
        if len(full_text) > chars:
            full_text = full_text[:chars] + " ...[truncated]"
    except Exception as e:
        abstract = f"[Could not read PDF: {e}]"
    return {
        "source": "Uploaded", "title": title, "authors": "", "published": "local",
        "pdf_url": "", "landing_url": os.path.basename(name),
        "abstract": abstract or "(no text extracted — may be a scanned PDF)",
        "doi": "", "full_text": full_text,
        "venue": "", "publisher": "",
    }


def uploaded_pdfs_to_candidates(uploaded_files, chars=12000):
    out = []
    for uf in uploaded_files or []:
        try:
            data = uf.getvalue()
        except Exception:
            data = uf.read()
        out.append(_pdf_bytes_to_candidate(getattr(uf, "name", "upload.pdf"), data, chars))
    return out


def folder_pdfs_to_candidates(folder, chars=12000, max_files=100):
    """Recursively scan a local folder (and subfolders) for PDFs."""
    out = []
    if not folder or not os.path.isdir(folder):
        return out, f"Folder not found: {folder}"
    found = []
    for root, _dirs, files in os.walk(folder):
        for fn in files:
            if fn.lower().endswith(".pdf"):
                found.append(os.path.join(root, fn))
    found = sorted(found)[:max_files]
    for path in found:
        try:
            with open(path, "rb") as f:
                data = f.read()
            out.append(_pdf_bytes_to_candidate(path, data, chars))
        except Exception as e:
            out.append({"source": "Uploaded", "title": os.path.basename(path),
                        "authors": "", "published": "local", "pdf_url": "",
                        "landing_url": path, "abstract": f"[Could not read: {e}]",
                        "doi": "", "full_text": "", "venue": "", "publisher": ""})
    return out, None


def ensure_fulltext(p, chars=6000, log=None):
    """Download + extract full text for one paper (only called for SELECTED papers). Handles
    both academic PDFs (pdf_url) and general web results (fetches the actual page content).
    Falls back to the abstract if there's nothing to fetch or extraction fails."""
    if p.get("full_text"):
        return p
    text = ""
    if p.get("source") == "Web" and not p.get("pdf_url"):
        page_url = p.get("landing_url") or p.get("pdf_url")
        if page_url:
            if log:
                log(f"🌐 Reading web page: **{p['title'][:70]}**\n")
            page = fetch_page_text(page_url, chars=chars)
            if page and not page.startswith("["):
                text = page
            elif log:
                log(f"   ⚠️ could not read page ({page}); using snippet instead\n")
    else:
        url = p.get("pdf_url")
        if url:
            if log:
                log(f"📥 Reading: **{p['title'][:70]}**\n")
            try:
                resp = requests.get(url, timeout=60, verify=False, headers=BROWSER_HEADERS)
                resp.raise_for_status()
                reader = PdfReader(io.BytesIO(resp.content))
                raw = "\n".join((pg.extract_text() or "") for pg in reader.pages)
                text = " ".join(raw.split())
                if len(text) > chars:
                    text = text[:chars] + " ...[truncated]"
            except Exception as e:
                text = ""
                if log:
                    log(f"   ⚠️ could not read PDF ({e}); using abstract instead\n")
    if not text:
        text = f"(Full text unavailable — using abstract only)\n{p.get('abstract','')}"
    p["full_text"] = text
    return p


def fetch_pdf_bytes(url, timeout=60):
    """Download a paper's raw PDF bytes for local saving. Returns (bytes_or_None, error_or_None)."""
    if not url:
        return None, "No PDF URL available for this source."
    try:
        resp = requests.get(url, timeout=timeout, verify=False, headers=BROWSER_HEADERS)
        resp.raise_for_status()
        return resp.content, None
    except Exception as e:
        return None, str(e)


def render_paper_card_extra(p, key_prefix):
    """Renders 'Read full text' and 'Download PDF' controls for one paper, meant to sit inside
    that paper's expander in a Gatekeeper list. Mutates `p` in place when full text is fetched
    (persists automatically since `p` is a reference into the session_state candidates list)."""
    c1, c2 = st.columns(2)
    with c1:
        if st.button("📖 Read full text", key=f"{key_prefix}_readfull"):
            with st.spinner("Fetching and extracting full text…"):
                ensure_fulltext(p, chars=8000)
            st.rerun()
    with c2:
        if p.get("pdf_url"):
            if st.button("📥 Fetch PDF for download", key=f"{key_prefix}_fetchpdf"):
                with st.spinner("Downloading PDF…"):
                    pdf_bytes, err = fetch_pdf_bytes(p["pdf_url"])
                if err:
                    st.warning(f"Could not download: {err}")
                else:
                    st.session_state[f"{key_prefix}_pdfbytes"] = pdf_bytes
                    st.rerun()
        else:
            st.caption("No direct PDF link available for this source.")

    if p.get("full_text"):
        st.text_area("Full text (extracted):", value=p["full_text"], height=200,
                     key=f"{key_prefix}_fulltext_view", disabled=True)

    cached_pdf = st.session_state.get(f"{key_prefix}_pdfbytes")
    if cached_pdf:
        safe_name = re.sub(r"[^A-Za-z0-9]+", "_", p["title"])[:60] or "paper"
        st.download_button("💾 Download PDF", data=cached_pdf, file_name=f"{safe_name}.pdf",
                           mime="application/pdf", key=f"{key_prefix}_dlbtn")


def papers_to_context_block(papers):
    parts = []
    for i, p in enumerate(papers, 1):
        parts.append(
            f"=== PAPER {i} ({p['source']}) ===\n"
            f"Title: {p['title']}\nAuthors: {p['authors']}\nPublished: {p['published']}\n"
            f"Link: {p.get('pdf_url') or p.get('landing_url')}\n"
            f"Abstract: {p['abstract']}\n\n"
            f"FULL TEXT:\n{p['full_text']}\n"
        )
    return "\n\n".join(parts)


# =============================================================================
# BibTeX
# =============================================================================
def to_bibtex(papers):
    entries, used = [], {}
    for p in papers:
        au = (p.get("authors") or "").strip()
        first = (au.split(",")[0].strip().split()[-1] if au else "anon")
        first = re.sub(r"[^A-Za-z]", "", first).lower() or "anon"
        year = (re.sub(r"\D", "", str(p.get("published", ""))) or "0000")[:4]
        tw = (p.get("title", "").split() or ["paper"])[0]
        tw = re.sub(r"[^A-Za-z]", "", tw).lower() or "paper"
        key = f"{first}{year}{tw}"
        n = used.get(key, 0); used[key] = n + 1
        if n:
            key += chr(ord('a') + n - 1)
        fields = [f"  title={{{p.get('title','')}}}",
                  f"  author={{{au}}}", f"  year={{{year}}}"]
        if p.get("source"):
            fields.append(f"  note={{Source: {p['source']}}}")
        if p.get("doi"):
            fields.append(f"  doi={{{p['doi']}}}")
        link = p.get("pdf_url") or p.get("landing_url")
        if link:
            fields.append(f"  url={{{link}}}")
        entries.append("@article{" + key + ",\n" + ",\n".join(fields) + "\n}")
    return "\n\n".join(entries)


def est_tokens(papers, chars_per_paper):
    """Rough token estimate (~4 chars/token) for the selected papers' injected text."""
    total = 0
    for p in papers:
        total += min(len(p.get("abstract", "")), chars_per_paper) + 300
        total += chars_per_paper  # reserve for full text
    return int(total / 4)


# =============================================================================
# SURVEY GENERATION — build an evidence digest, then write section by section
# =============================================================================
def build_source_digest(papers, per_chars=1500):
    """Numbered evidence base the model cites as [1], [2], …"""
    lines = []
    for i, p in enumerate(papers, 1):
        body = (p.get("full_text") or p.get("abstract") or "").strip()
        body = body[:per_chars] if body else "(no text available)"
        venue = p.get("venue") or p.get("source", "")
        lines.append(f"[{i}] {p.get('title','')} ({p.get('published','n/a')}; {venue})\n{body}")
    return "\n\n".join(lines)


def numbered_references(papers):
    out = []
    for i, p in enumerate(papers, 1):
        link = p.get("pdf_url") or p.get("landing_url") or ""
        out.append(f"[{i}] {p.get('title','')} — {p.get('authors','') or 'n/a'} "
                   f"({p.get('published','n/a')}). {link}")
    return "\n".join(out)


def build_novelty_digest(papers, web_results, per_chars=1500):
    """Like build_source_digest, but numbers papers AND general web hits (products, patents,
    blog posts) together in one evidence base, since a novelty check needs both."""
    lines = []
    n = 0
    for p in papers:
        n += 1
        body = (p.get("full_text") or p.get("abstract") or "").strip()
        body = body[:per_chars] if body else "(no text available)"
        venue = p.get("venue") or p.get("source", "")
        lines.append(f"[{n}] (Paper — {venue}) {p.get('title', '')} ({p.get('published', 'n/a')})\n{body}")
    for r in web_results:
        n += 1
        snippet = (r.get("page_text") or r.get("snippet") or "").strip()[:per_chars]
        lines.append(f"[{n}] (Web) {r.get('title', '')}\nURL: {r.get('url', '')}\n{snippet}")
    return "\n\n".join(lines), n


def numbered_novelty_references(papers, web_results):
    out = []
    n = 0
    for p in papers:
        n += 1
        link = p.get("pdf_url") or p.get("landing_url") or ""
        out.append(f"[{n}] {p.get('title', '')} — {p.get('authors', '') or 'n/a'} "
                   f"({p.get('published', 'n/a')}). {link}")
    for r in web_results:
        n += 1
        out.append(f"[{n}] {r.get('title', '')} — {r.get('url', '')}")
    return "\n".join(out)


NOVELTY_SYSTEM = (
    "You are a rigorous research/IP analyst assessing the novelty of an idea. You are given the "
    "user's idea description and a numbered set of prior-art sources (academic papers and general "
    "web results — products, patents, blog posts) gathered specifically to check for overlap. "
    "Using ONLY the provided sources, cite them as [n] wherever relevant. Write a structured "
    "report with EXACTLY these Markdown sections, in this order:\n\n"
    "## Idea Summary\n"
    "Restate the idea clearly and concisely in your own words (2-4 sentences).\n\n"
    "## Novelty Verdict\n"
    "State one of: **Novel**, **Not Novel**, or **Partially Novel / Incremental** as the first "
    "line, then explain why in 3-6 sentences, citing specific overlapping or closely related "
    "sources by [n]. Be honest and specific — do not default to 'novel' just to be encouraging, "
    "and do not default to 'not novel' just because *something* related exists.\n\n"
    "## Related Work\n"
    "A numbered list matching the source numbers. For each: one line stating what it does and "
    "how it relates to (overlaps with, or clearly differs from) the idea.\n\n"
    "## Enhancement Suggestions\n"
    "3-5 concrete, specific suggestions for how to strengthen, differentiate, or extend the idea "
    "to increase its novelty and real-world impact, grounded in the gaps visible in the related "
    "work — not generic advice.\n\n"
    "## Suggested Angle / Research Gap\n"
    "1-2 sentences on the most promising specific angle or gap this idea could target to be "
    "clearly novel and publishable/patentable.\n\n"
    "RULES: Base every claim ONLY on the provided sources or the idea text — never invent papers, "
    "products, companies, or results. If sources are sparse or weak, say so honestly in the "
    "verdict rather than overclaiming novelty or lack of it."
)


_PDF_UNSAFE_CHARS = {
    "\u2010": "-", "\u2011": "-", "\u2012": "-", "\u2013": "-", "\u2014": "--",
    "\u2015": "--", "\u2212": "-",                                     # hyphen/dash variants, minus sign
    "\u2018": "'", "\u2019": "'", "\u201a": "'", "\u201b": "'",         # smart single quotes
    "\u201c": '"', "\u201d": '"', "\u201e": '"', "\u201f": '"',         # smart double quotes
    "\u2026": "...",                                                    # ellipsis
    "\u00a0": " ", "\u2000": " ", "\u2001": " ", "\u2002": " ",
    "\u2003": " ", "\u2004": " ", "\u2005": " ", "\u2006": " ",
    "\u2007": " ", "\u2008": " ", "\u2009": " ", "\u200a": " ",         # various Unicode spaces
    "\u200b": "", "\u200c": "", "\u200d": "", "\ufeff": "",             # zero-width characters
    "\u2022": "-", "\u25cf": "-", "\u25aa": "-",                        # bullet-ish glyphs used inline
}


def _sanitize_pdf_text(t):
    """Map Unicode punctuation that ReportLab's default Helvetica font can't render (it shows
    a missing-glyph black square instead of erroring) to plain-ASCII equivalents it can."""
    if not t:
        return t
    return "".join(_PDF_UNSAFE_CHARS.get(ch, ch) for ch in t)


def md_to_pdf_bytes(md_text, title="Document", figs=None):
    """Render simple markdown (headings, paragraphs, bullets, a References section, and
    [[FIGURE:n]] markers referencing entries in `figs`) into a clean single-column PDF.
    Returns bytes, or None if reportlab is missing."""
    try:
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.units import mm
        from reportlab.lib import colors
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.platypus import (SimpleDocTemplate, Paragraph, Spacer, HRFlowable, Image)
    except ImportError:
        return None
    import html as _html
    buf = io.BytesIO()
    doc = SimpleDocTemplate(buf, pagesize=A4, leftMargin=20*mm, rightMargin=20*mm,
                            topMargin=18*mm, bottomMargin=18*mm, title=title)
    ss = getSampleStyleSheet()
    body = ParagraphStyle("body", parent=ss["BodyText"], fontSize=10.5, leading=15,
                          spaceAfter=6, alignment=4)
    h1 = ParagraphStyle("h1", parent=ss["Heading1"], fontSize=17, spaceBefore=8, spaceAfter=6)
    h2 = ParagraphStyle("h2", parent=ss["Heading2"], fontSize=13, spaceBefore=10, spaceAfter=4)
    h3 = ParagraphStyle("h3", parent=ss["Heading3"], fontSize=11.5, spaceBefore=6, spaceAfter=3)
    refst = ParagraphStyle("ref", parent=body, fontSize=9, leading=12.5, leftIndent=12,
                           firstLineIndent=-12, spaceAfter=3)
    ital = ParagraphStyle("ital", parent=body, fontName="Helvetica-Oblique",
                          textColor=colors.grey)

    def inline(t):
        # ReportLab's default Helvetica (Base14) font only covers WinAnsiEncoding — anything
        # outside it (e.g. a non-breaking hyphen U+2011, en/em dash, smart quotes, minus sign)
        # silently renders as a missing-glyph black square instead of erroring, which is what
        # was showing up as "low■power" in place of "low-power". Normalize the common ones to
        # their plain-ASCII equivalents before the font ever has to render them.
        t = _sanitize_pdf_text(t)
        t = _html.escape(t)
        t = re.sub(r"\*\*(.+?)\*\*", r"<b>\1</b>", t)
        t = re.sub(r"(?<!\*)\*(?!\*)(.+?)(?<!\*)\*(?!\*)", r"<i>\1</i>", t)
        t = re.sub(r"`(.+?)`", r'<font face="Courier">\1</font>', t)
        return t

    flow, in_refs = [], False
    for raw in md_text.split("\n"):
        line = raw.rstrip()
        if not line.strip():
            continue
        fig_m = re.match(r"^\[\[FIGURE:(\d+)\]\]$", line)
        if line.startswith("# "):
            flow.append(Paragraph(inline(line[2:]), h1))
            flow.append(HRFlowable(width="100%", thickness=0.6, color=colors.grey, spaceAfter=6))
        elif line.startswith("## "):
            txt = line[3:]
            in_refs = txt.strip().lower().startswith("references") or txt.strip().lower().startswith("sources")
            flow.append(Paragraph(inline(txt), h2))
        elif line.startswith("### "):
            flow.append(Paragraph(inline(line[4:]), h3))
        elif fig_m and figs is not None:
            idx = int(fig_m.group(1))
            if 0 <= idx < len(figs):
                fdata = figs[idx]
                try:
                    width_mm = 150
                    height_mm = width_mm * fdata.get("aspect", 0.55)
                    flow.append(Spacer(1, 6))
                    flow.append(Image(io.BytesIO(fdata["png"]), width=width_mm * mm, height=height_mm * mm))
                except Exception:
                    pass
        elif line.startswith("---"):
            flow.append(Spacer(1, 4))
        elif re.match(r"^\*[^*].*\*$", line):     # *italic caption line*
            flow.append(Paragraph(inline(line.strip("*")), ital))
        elif line.startswith("- ") or line.startswith("• "):
            flow.append(Paragraph("&bull;&nbsp;&nbsp;" + inline(line[2:]), body))
        else:
            flow.append(Paragraph(inline(line), refst if in_refs else body))
    doc.build(flow)
    return buf.getvalue()


SURVEY_SYSTEM = (
    "You are an expert academic writing ONE section of a rigorous SURVEY (literature review) "
    "paper. Write formal, in-depth academic prose. Cite sources by their number, e.g. [1], [3] — "
    "EXCEPT in the Abstract, which must stand alone with NO citation markers at all (write it so "
    "it's fully understandable without the numbered source list). "
    "Use ONLY facts present in the provided numbered sources — never invent results, numbers, or "
    "citations. Do not repeat material that belongs to other sections. Output only the requested "
    "section's body text (no section number/heading — that is added for you)."
)

# (title, guidance, base_word_target)
SURVEY_SECTIONS = [
    ("Abstract", "Summarize the scope of the survey, what it covers, and the key takeaways. "
     "Do NOT use any citation markers like [1] here — the abstract must read standalone.", 180),
    ("1. Introduction", "Motivate the topic and its importance, state the scope of this survey, "
     "and outline the paper's structure.", 500),
    ("2. Background and Preliminaries", "Explain the key concepts, definitions, and foundational "
     "methods a reader needs to understand the field.", 500),
    ("3. Taxonomy and Themes", "Organize the surveyed works into coherent themes or categories and "
     "compare how the sources relate to each other.", 650),
    ("4. Methods and Approaches", "Discuss the methodologies used across the works, comparing their "
     "strengths and weaknesses, with citations.", 650),
    ("5. Key Findings and Results", "Synthesize the main findings across the sources, highlighting "
     "agreements, contradictions, and trends.", 600),
    ("6. Open Challenges and Limitations", "Discuss the gaps, limitations, and unresolved problems "
     "evident from the literature.", 450),
    ("7. Future Directions", "Propose promising, well-grounded research directions that address the "
     "identified gaps.", 450),
    ("8. Conclusion", "Summarize the survey and its main messages.", 250),
]

# =============================================================================
# MAKE PRESENTATION — turn any source text (an idea, a written paper, code, or a
# past conversation) into a slide-by-slide outline via the LLM, then render an
# actual .pptx from it with python-pptx.
# =============================================================================
PRESENTATION_OUTLINE_SYSTEM = (
    "You turn source material (an idea, a paper, code, or a conversation) into a clear, "
    "well-structured slide deck outline for a presentation.\n\n"
    'Respond with ONLY a JSON object of this form:\n'
    '{"title": "Deck title", "subtitle": "Optional one-line subtitle", "slides": '
    '[{"title": "Slide title", "bullets": ["point 1", "point 2"], "notes": "optional speaker notes"}]}\n\n'
    "Rules: produce exactly the requested number of content slides (you may add one final "
    "'Thank You / Questions' slide beyond that count). Each slide should have 3-6 short, punchy "
    "bullets — a slide is a prompt for the speaker, not a document; no full paragraphs. Keep the "
    "flow logical: motivation/context first, then the core content, then results, impact, or a "
    "wrap-up. If the source is code, focus on what it does and why, not a line-by-line walkthrough. "
    "Output ONLY the JSON object — no explanation, no markdown fences."
)


def generate_slide_outline(source_text, n_slides, model, api_base=None, api_key=None,
                           num_ctx=None, guidance="", max_tokens=4000):
    """Ask the LLM to turn source_text into a slide-by-slide outline.
    Returns (outline_dict_or_None, raw_reply)."""
    user_msg = f"Number of content slides: {n_slides}\n"
    if guidance.strip():
        user_msg += f"Extra guidance: {guidance.strip()}\n"
    user_msg += f"\nSource material:\n\n{source_text}"
    msgs = [{"role": "system", "content": PRESENTATION_OUTLINE_SYSTEM}, {"role": "user", "content": user_msg}]
    try:
        chunks = list(stream_chat(msgs, model, api_base=api_base, api_key=api_key,
                                  temperature=0.4, num_ctx=num_ctx, max_tokens=max_tokens))
        raw = "".join(chunks).strip()
    except Exception as e:
        return None, f"(LLM call failed: {e})"
    cleaned = re.sub(r"^```(json)?\s*|\s*```$", "", raw, flags=re.I).strip()
    try:
        return json.loads(cleaned), raw
    except Exception:
        m = re.search(r"\{.*\}", cleaned, re.S)
        if m:
            try:
                return json.loads(m.group(0)), raw
            except Exception:
                pass
    return None, raw


def _blend_hex(hex_str, toward, factor):
    """Blends `hex_str` toward `toward` ('FFFFFF' or '000000', etc.) by `factor` (0=no
    change, 1=fully `toward`). Used to derive a dark variant (for title/closing slide
    backgrounds) and a light tint (for subtitle text) from one accent color."""
    h = hex_str.lstrip("#")
    t = toward.lstrip("#")
    r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
    tr, tg, tb = int(t[0:2], 16), int(t[2:4], 16), int(t[4:6], 16)
    r = int(r + (tr - r) * factor)
    g = int(g + (tg - g) * factor)
    b = int(b + (tb - b) * factor)
    return f"{r:02X}{g:02X}{b:02X}"


def build_pptx_from_outline(outline, theme_color="1F4E79"):
    """Builds a .pptx from an outline dict {"title", "subtitle", "slides": [{"title",
    "bullets", "notes"}, ...]}. Returns bytes, or None if python-pptx isn't installed.

    Built from a blank layout with everything hand-positioned, rather than the default
    template's bare "Title and Content" layout — that one renders as plain black text on
    white with no color at all. Colour comes from full-bleed background fills and
    typography (dark bookend slides, a colored kicker numeral, colored bullet markers),
    deliberately avoiding thin decorative bars/stripes under titles or along edges, which
    read as an obvious AI-generated tell rather than an actual design choice."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    except ImportError:
        return None

    theme_color = re.sub(r"[^0-9A-Fa-f]", "", theme_color or "1F4E79")[:6].upper() or "1F4E79"
    dark_bg = _blend_hex(theme_color, "000000", 0.55)      # bookend slide backgrounds
    light_tint = _blend_hex(theme_color, "FFFFFF", 0.55)   # subtitle text on the dark bg
    body_text, muted_text = "2B2B2B", "8A8A8A"

    SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)
    prs = Presentation()
    prs.slide_width, prs.slide_height = SLIDE_W, SLIDE_H
    blank_layout = prs.slide_layouts[6]

    def set_bg(slide, hex_color):
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor.from_string(hex_color)

    def add_text(slide, text, left, top, width, height, size, color, bold=False,
                align=PP_ALIGN.LEFT, anchor=None):
        box = slide.shapes.add_textbox(left, top, width, height)
        tf = box.text_frame
        tf.word_wrap = True
        if anchor is not None:
            tf.vertical_anchor = anchor
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size, run.font.bold, run.font.name = Pt(size), bold, "Calibri"
        run.font.color.rgb = RGBColor.from_string(color)
        return box

    total = len(outline.get("slides", [])) + 2  # + title + closing

    # ---- Title slide (full-bleed dark background) ---------------------------
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide, dark_bg)
    add_text(slide, outline.get("title") or "Presentation", Inches(1), Inches(2.6),
             Inches(11.3), Inches(1.7), 40, "FFFFFF", bold=True, anchor=MSO_ANCHOR.BOTTOM)
    if outline.get("subtitle"):
        add_text(slide, outline["subtitle"], Inches(1), Inches(4.35), Inches(11.3), Inches(0.7),
                 18, light_tint)

    # ---- Content slides (white background, colored kicker + title + bullets)
    for i, sl in enumerate(outline.get("slides", []), start=1):
        slide = prs.slides.add_slide(blank_layout)
        set_bg(slide, "FFFFFF")

        add_text(slide, f"{i:02d}", Inches(0.9), Inches(0.55), Inches(1.5), Inches(0.5),
                 16, theme_color, bold=True)
        add_text(slide, sl.get("title") or "", Inches(0.9), Inches(0.95), Inches(11.5), Inches(1.0),
                 30, theme_color, bold=True)

        bullets = [b for b in (sl.get("bullets") or []) if str(b).strip()] or [""]
        n = len(bullets)
        size = 20 if n <= 5 else (17 if n <= 7 else 15)
        body_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.05), Inches(11.3), Inches(4.7))
        tf = body_box.text_frame
        tf.word_wrap = True
        for j, b in enumerate(bullets):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.space_after = Pt(max(10, 22 - n))
            dot = p.add_run()
            dot.text = "●  "
            dot.font.size, dot.font.color.rgb, dot.font.name = Pt(size), RGBColor.from_string(theme_color), "Calibri"
            txt = p.add_run()
            txt.text = str(b)
            txt.font.size, txt.font.color.rgb, txt.font.name = Pt(size), RGBColor.from_string(body_text), "Calibri"

        add_text(slide, f"{i + 1} / {total}", Inches(11.9), Inches(7.05), Inches(1.1), Inches(0.35),
                 10, muted_text, align=PP_ALIGN.RIGHT)

        notes = sl.get("notes") or ""
        if notes:
            slide.notes_slide.notes_text_frame.text = notes

    # ---- Closing slide (bookends the title slide) ---------------------------
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide, dark_bg)
    add_text(slide, "Thank You", Inches(1), Inches(3.1), Inches(11.3), Inches(1.1),
             40, "FFFFFF", bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, "Questions?", Inches(1), Inches(4.25), Inches(11.3), Inches(0.7),
             18, light_tint, align=PP_ALIGN.CENTER)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# =============================================================================
# RELEVANCE FILTER — after gathering candidate papers, ask the LLM which ones are
# actually off-topic and should be dropped before selection/writing. Fails open
# (keeps everything) if the model's output can't be parsed, so a formatting hiccup
# never silently wipes out the candidate list.
# =============================================================================
RELEVANCE_FILTER_SYSTEM = (
    "You are screening search results for a literature survey/review. You are given the "
    "survey topic and a numbered list of gathered papers (title + short snippet). Identify "
    "which papers are NOT actually relevant to the topic — off-topic, tangential, or clearly "
    "about something else — so they can be excluded before the survey is written.\n\n"
    'Respond with ONLY a JSON object mapping each number to REMOVE to a short reason, e.g. '
    '{"2": "different sensor modality, unrelated to this topic", "5": "general survey of an '
    'unrelated field"}. If every paper is relevant, respond with an empty object: {}. Output '
    "ONLY the JSON object — no explanation, no markdown fences."
)


def filter_relevant_papers(topic, papers, model, api_base=None, api_key=None, num_ctx=None):
    """Ask the LLM which gathered papers are off-topic. Returns (kept, removed) where each
    removed entry has an added '_removal_reason' key. Fails open: any parsing problem, or the
    model trying to remove everything, just returns the original list untouched."""
    if not papers:
        return papers, []
    listing = "\n".join(
        f"{i}. [{p.get('source', '?')}] {p.get('title', '')} — {(p.get('abstract') or '')[:150]}"
        for i, p in enumerate(papers, 1))
    msgs = [{"role": "system", "content": RELEVANCE_FILTER_SYSTEM},
            {"role": "user", "content": f"Topic: {topic}\n\nPapers:\n{listing}"}]
    try:
        chunks = list(stream_chat(msgs, model, api_base=api_base, api_key=api_key,
                                  temperature=0.0, num_ctx=num_ctx, max_tokens=500))
        raw = re.sub(r"^```(json)?\s*|\s*```$", "", "".join(chunks).strip(), flags=re.I)
        mapping = json.loads(raw) or {}
    except Exception:
        return papers, []

    reasons = {}
    for k, v in mapping.items():
        try:
            reasons[int(k)] = v
        except (TypeError, ValueError):
            continue
    if not reasons:
        return papers, []

    kept, removed = [], []
    for i, p in enumerate(papers, 1):
        if i in reasons:
            removed.append({**p, "_removal_reason": reasons[i]})
        else:
            kept.append(p)
    return (papers, []) if not kept else (kept, removed)


# =============================================================================
# DYNAMIC LLM-AUTHORED FIGURES — instead of a fixed chart template, the model
# writes actual matplotlib Python code for each section (block diagram, flowchart,
# taxonomy tree, timeline, or a data chart when real data exists — its choice),
# which is executed in a restricted sandbox. Retries once with the error fed back
# if the generated code fails.
# =============================================================================
def _fig_to_png_bytes(fig):
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=150, bbox_inches="tight")
    try:
        import matplotlib.pyplot as plt
        plt.close(fig)
    except Exception:
        pass
    return buf.getvalue()


def _get_plt():
    """Lazily import matplotlib with a non-interactive backend. Returns the pyplot module,
    or None if matplotlib isn't installed (figures are then silently skipped)."""
    try:
        import matplotlib
        matplotlib.use("Agg")
        import matplotlib.pyplot as plt
        return plt
    except ImportError:
        return None


FIGURE_CODE_SYSTEM = (
    "You are writing Python code that implements an ALREADY-DECIDED figure plan exactly — you "
    "are not choosing what to draw here, only translating the given plan into working code.\n\n"
    "For DIAGRAMS (block diagram, flowchart, taxonomy tree, conceptual framework, timeline) you "
    "have a small toolkit ALREADY AVAILABLE in the sandbox — use it, do NOT use raw "
    "matplotlib.patches, it is much less error-prone and gives consistent, clean styling:\n"
    "  fig, ax = new_diagram_fig(title='...')            # clean canvas on a 0-10 x 0-10 grid\n"
    "  b1 = add_box(ax, x, y, w, h, 'Label text')          # a styled box; x,y,w,h are in the "
    "0-10 grid (e.g. a box at x=1,y=4,w=2.2,h=1.4); returns a HANDLE — save it in a variable\n"
    "  add_arrow(ax, b1, b2, label='opt.')                 # connects two box HANDLES from "
    "add_box with a clean arrow between their nearest edges, correctly directed automatically\n"
    "CRITICAL: add_arrow takes the box HANDLES returned by add_box (e.g. `b1`, `b2`) — NEVER "
    "raw x/y coordinates. Passing coordinates instead of handles is exactly what causes arrows "
    "to cut through boxes at the wrong angle, so always do `b1 = add_box(...)` then "
    "`add_arrow(ax, b1, b2)`.\n"
    "Text is automatically wrapped and shrunk to fit inside each box, so don't worry about "
    "exact sizing — but keep labels SHORT (2-4 words, or a short phrase) and size box width "
    "roughly as (label length in characters) / 5.5, e.g. an 18-character label wants w≈3.3. "
    "Avoid long parenthetical asides in labels; put detail in the caption instead.\n"
    "Place every box so it fits within the 0-10 by 0-10 grid (if it doesn't quite fit, the view "
    "auto-expands afterward, but staying within 0-10 keeps proportions sensible). Leave a real "
    "gap of at least 0.4-0.5 units between adjacent boxes — boxes packed edge-to-edge leave no "
    "room for arrow labels and make them overlap. If two arrows run close and parallel to each "
    "other (e.g. two stacked boxes both connecting to a third), stagger them or only label one "
    "to avoid the labels colliding. A simple left-to-right or top-to-bottom layout works well.\n\n"
    "For a genuine DATA chart (bar/line/pie) — ONLY if real numbers are actually given to you, "
    "never invented — use plain `plt.subplots()` / `ax.bar()` / `ax.plot()` etc. directly "
    "instead of the diagram toolkit.\n\n"
    "RULES:\n"
    "- The FIRST line of your output must be exactly: # CAPTION: <one short sentence caption>\n"
    "- No imports needed for diagrams (the toolkit + `plt` are already available). For data "
    "charts you may also use `np`. No other imports — the code runs in a restricted sandbox.\n"
    "- `ax` must always come from `new_diagram_fig()` or `plt.subplots()` — never reference an "
    "axes variable under a different name than the one that function actually returned.\n"
    "- The code MUST end with a variable named `fig` holding the Figure. Do not call plt.show().\n"
    "- No file I/O, no network access.\n"
    "- EVERY box created with add_box MUST appear in at least one add_arrow call — a box with no "
    "arrow in or out will be rejected. If a box genuinely has nothing to connect to, remove it "
    "from the diagram instead of leaving it floating.\n"
    "- Implement EXACTLY the elements in the given plan — do not add, remove, or embellish "
    "beyond it. Keep the code well under 30 lines; a simple, complete diagram beats an "
    "elaborate, unfinished one.\n"
    "- Output ONLY the caption comment line followed by a single fenced Python code block — "
    "no explanation before or after it."
)


def _extract_figure_caption(code):
    for line in code.splitlines()[:4]:
        m = re.match(r"^\s*#\s*CAPTION:\s*(.+)$", line, re.I)
        if m:
            return m.group(1).strip()
    return ""


# Restricted execution environment for LLM-authored figure code: only a whitelisted set of
# builtins, and imports limited to matplotlib/numpy-family modules. This is a pragmatic sandbox
# against accidental mistakes in model-generated code, not a hardened security boundary against
# a deliberately adversarial payload — appropriate here since the "attacker" would be the model
# itself trying (and failing) to draw a chart, not an untrusted third party.
_FIGURE_ALLOWED_MODULE_ROOTS = {"matplotlib", "numpy", "math", "itertools", "textwrap"}


def _figure_restricted_import(name, globals=None, locals=None, fromlist=(), level=0):
    if name.split(".")[0] not in _FIGURE_ALLOWED_MODULE_ROOTS:
        raise ImportError(f"import of '{name}' is not allowed in generated figure code")
    return __import__(name, globals, locals, fromlist, level)


def _figure_safe_builtins():
    import builtins as _b
    allowed = ["range", "len", "enumerate", "zip", "min", "max", "sum", "abs", "round", "pow",
              "list", "dict", "tuple", "set", "frozenset", "str", "int", "float", "bool",
              "complex", "sorted", "reversed", "map", "filter", "all", "any", "isinstance",
              "issubclass", "type", "chr", "ord", "divmod", "format", "True", "False", "None",
              "Exception", "ValueError", "TypeError", "KeyError", "IndexError", "StopIteration"]
    safe = {name: getattr(_b, name) for name in allowed if hasattr(_b, name)}
    safe["__import__"] = _figure_restricted_import
    safe["print"] = lambda *a, **k: None
    return safe


# A small, deliberately opinionated "diagram toolkit" handed to the model instead of raw
# matplotlib.patches. This closes off the two bug classes that kept showing up in practice —
# an axes variable referenced under the wrong name, and hand-rolled FancyBboxPatch calls
# colliding on kwargs like `linewidth` — by giving the model a handful of simple, fixed-signature
# functions to call instead. It also guarantees every generated figure shares the same clean,
# consistent visual style rather than whatever ad hoc styling the model improvises each time.
_DIAGRAM_PALETTE = ["#4C72B0", "#55A868", "#C44E52", "#8172B2", "#CCB974", "#64B5CD"]


def _make_diagram_helpers(plt):
    import matplotlib.patches as mpatches
    import textwrap as _textwrap

    def new_diagram_fig(figsize=(7, 4.5), title=None):
        fig, ax = plt.subplots(figsize=figsize)
        ax.set_xlim(0, 10)
        ax.set_ylim(0, 10)
        ax.axis("off")
        # Tracks every box's extent so the view can be auto-expanded afterward if the model
        # places anything outside this default 0-10 grid — see the auto-fit step in
        # execute_figure_code. Without this, boxes placed beyond the fixed limits used to get
        # silently clipped off the edge of the rendered image.
        ax._db_extents = []
        # Tracks box IDs and which pairs got connected via add_arrow, so execute_figure_code can
        # detect a box that was drawn but never connected to anything — a diagram with visibly
        # disconnected, floating boxes reads as broken/incomplete.
        ax._db_box_ids = []
        ax._db_connections = []
        if title:
            ax.set_title(title, fontsize=13, fontweight="bold", pad=12)
        return fig, ax

    _counter = {"n": 0, "box_id": 0}

    def add_box(ax, x, y, w, h, label, color=None, fontsize=10, text_color="white"):
        chosen = color or _DIAGRAM_PALETTE[_counter["n"] % len(_DIAGRAM_PALETTE)]
        _counter["n"] += 1
        box_id = _counter["box_id"]
        _counter["box_id"] += 1
        box = mpatches.FancyBboxPatch(
            (x, y), w, h, boxstyle="round,pad=0.08,rounding_size=0.12",
            linewidth=1.4, edgecolor="#333333", facecolor=chosen)
        ax.add_patch(box)
        if hasattr(ax, "_db_extents"):
            ax._db_extents.append((x, y, x + w, y + h))
        if hasattr(ax, "_db_box_ids"):
            ax._db_box_ids.append((box_id, label))

        # matplotlib's `wrap=True` wraps to the FIGURE edge, not to this box — it does nothing
        # useful here, which is why text used to spill out past the box boundary. A fixed
        # "characters/points per data-unit" formula turned out unreliable across font/DPI
        # combinations, so instead: wrap with a reasonable starting guess, then actually
        # MEASURE the rendered text with matplotlib's own renderer and iteratively shrink the
        # font and/or rewrap narrower until it provably fits inside the box (with a small margin).
        use_fontsize = fontsize
        chars_per_line = max(int(w * 6.5), 6)
        wrapped = "\n".join(_textwrap.wrap(label, width=chars_per_line)) if label else ""
        txt = ax.text(x + w / 2, y + h / 2, wrapped, ha="center", va="center",
                     fontsize=use_fontsize, color=text_color, fontweight="medium")

        fig = ax.figure
        pad_w, pad_h = w * 0.88, h * 0.82
        for _ in range(6):
            if not label:
                break
            fig.canvas.draw()
            bbox = txt.get_window_extent(renderer=fig.canvas.get_renderer())
            bbox = bbox.transformed(ax.transData.inverted())
            text_w, text_h = bbox.x1 - bbox.x0, bbox.y1 - bbox.y0
            if text_w <= pad_w and text_h <= pad_h:
                break
            if use_fontsize <= 6:
                break  # floor reached — accept slight overflow rather than shrinking to nothing
            use_fontsize = max(use_fontsize - 1, 6)
            txt.set_fontsize(use_fontsize)
            if text_w > pad_w:
                chars_per_line = max(chars_per_line - 2, 5)
                txt.set_text("\n".join(_textwrap.wrap(label, width=chars_per_line)))

        # Handle describing this box's geometry, for add_arrow to connect to automatically.
        return {"_id": box_id, "x": x, "y": y, "w": w, "h": h, "cx": x + w / 2, "cy": y + h / 2,
               "left": (x, y + h / 2), "right": (x + w, y + h / 2),
               "top": (x + w / 2, y + h), "bottom": (x + w / 2, y)}

    def add_arrow(ax, box_a, box_b, label=None, color="#333333"):
        """Connect two boxes (the handles returned by add_box) with a clean arrow between
        their nearest facing edges. Always pass box handles here — never raw coordinates —
        so the arrow direction and anchor points are derived from real box geometry instead
        of guessed numbers, which is what used to send arrows through box interiors at
        arbitrary angles."""
        if hasattr(ax, "_db_connections") and "_id" in box_a and "_id" in box_b:
            ax._db_connections.append((box_a["_id"], box_b["_id"]))
        dx, dy = box_b["cx"] - box_a["cx"], box_b["cy"] - box_a["cy"]
        horizontal = abs(dx) >= abs(dy)
        if horizontal:
            start, end = (box_a["right"], box_b["left"]) if dx > 0 else (box_a["left"], box_b["right"])
        else:
            start, end = (box_a["top"], box_b["bottom"]) if dy > 0 else (box_a["bottom"], box_b["top"])
        ax.annotate("", xy=end, xytext=start,
                   arrowprops=dict(arrowstyle="-|>", color=color, lw=1.6, shrinkA=2, shrinkB=2))
        if label:
            # Placing the label AT the arrow midpoint used to put it right on top of (or inside)
            # a box whenever the model left little or no gap between the two boxes — exactly
            # what happened in practice. Instead, place it OUTSIDE the pair's combined bounding
            # box entirely: above both boxes for a horizontal connection, beside both for a
            # vertical one — this guarantees no overlap with either box's interior regardless
            # of how close together they are.
            mx = (start[0] + end[0]) / 2
            my = (start[1] + end[1]) / 2
            label_pad = 0.28
            if horizontal:
                mx = (start[0] + end[0]) / 2
                my = max(box_a["y"] + box_a["h"], box_b["y"] + box_b["h"]) + label_pad
            else:
                mx = max(box_a["x"] + box_a["w"], box_b["x"] + box_b["w"]) + label_pad
                my = (start[1] + end[1]) / 2
            txt = ax.text(mx, my, label, ha="center", va="center", fontsize=8.5, color=color,
                         style="italic", zorder=5,
                         bbox=dict(facecolor="white", edgecolor="none", alpha=0.85,
                                  boxstyle="round,pad=0.15"))
            if hasattr(ax, "_db_extents"):
                # A rough box around the label so the auto-fit view expands to include it too.
                half_w, half_h = 0.35 + 0.06 * len(label), 0.18
                ax._db_extents.append((mx - half_w, my - half_h, mx + half_w, my + half_h))

    return {"new_diagram_fig": new_diagram_fig, "add_box": add_box, "add_arrow": add_arrow}


def execute_figure_code(code):
    """Execute LLM-generated matplotlib code in a restricted namespace. Returns
    ({"png":..., "aspect":...}, None) on success, or (None, error_message) on failure."""
    plt = _get_plt()
    if plt is None:
        return None, "matplotlib not installed"

    safe_globals = {"__builtins__": _figure_safe_builtins(), "plt": plt}
    try:
        import numpy as np
        safe_globals["np"] = np
    except ImportError:
        pass
    try:
        safe_globals.update(_make_diagram_helpers(plt))
    except Exception:
        pass  # diagram helpers unavailable (e.g. patches submodule missing) — raw plt still works

    plt.close("all")
    local_ns = {}
    try:
        exec(code, safe_globals, local_ns)
    except Exception as e:
        return None, f"{type(e).__name__}: {e}"

    fig = local_ns.get("fig", safe_globals.get("fig"))
    if fig is None:
        nums = plt.get_fignums()
        if nums:
            fig = plt.figure(nums[-1])
    if fig is None:
        return None, "code did not produce a 'fig' variable"

    # Reject diagrams with disconnected boxes: a box that was drawn but never appears in any
    # add_arrow call reads as visually broken/incomplete (a floating box nothing points to or
    # from). This feeds back into the existing retry loop with the specific box name(s), rather
    # than silently accepting an obviously-incomplete diagram.
    for a in getattr(fig, "axes", []):
        box_ids = getattr(a, "_db_box_ids", None)
        connections = getattr(a, "_db_connections", [])
        if box_ids and len(box_ids) > 1:
            connected = {bid for pair in connections for bid in pair}
            orphans = [label for bid, label in box_ids if bid not in connected]
            if orphans:
                names = ", ".join(f"'{o}'" for o in orphans)
                return None, (f"disconnected box(es) with no arrow in or out: {names} — every "
                             f"box must be connected via add_arrow")

    # Auto-fit safety net: if any box (tracked via _db_extents) falls outside the diagram's
    # current view limits — the model drew beyond the suggested 0-10 grid — expand the view to
    # include it, with a small margin. This only ever widens the view, never shrinks it, so
    # diagrams that stayed within bounds render exactly as before. This is what used to cause
    # boxes to get silently clipped off the edge of the image.
    for a in getattr(fig, "axes", []):
        extents = getattr(a, "_db_extents", None)
        if not extents:
            continue
        xs0, ys0 = min(e[0] for e in extents), min(e[1] for e in extents)
        xs1, ys1 = max(e[2] for e in extents), max(e[3] for e in extents)
        pad_x = max((xs1 - xs0) * 0.06, 0.3)
        pad_y = max((ys1 - ys0) * 0.06, 0.3)
        cur_xlim, cur_ylim = a.get_xlim(), a.get_ylim()
        a.set_xlim(min(cur_xlim[0], xs0 - pad_x), max(cur_xlim[1], xs1 + pad_x))
        a.set_ylim(min(cur_ylim[0], ys0 - pad_y), max(cur_ylim[1], ys1 + pad_y))

    try:
        png = _fig_to_png_bytes(fig)
        w, h = fig.get_size_inches()
        aspect = (h / w) if w else 0.6
    except Exception as e:
        return None, f"render error: {e}"
    return {"png": png, "aspect": aspect}, None


FIGURE_PLAN_SYSTEM = (
    "You are planning ONE illustrative figure for a specific section of a literature survey. "
    "Think it through carefully before anything gets built — you are NOT writing code here, "
    "only deciding and describing the figure in plain text.\n\n"
    "Work through:\n"
    "1. TYPE — what kind of figure best fits this section: a block diagram, a flowchart "
    "(boxes and arrows), a conceptual framework diagram, a taxonomy tree, a timeline, or — "
    "only if real numeric data is actually available — a data chart.\n"
    "2. ELEMENTS — list each box/node/bar with its exact label text (3-6 elements; simple and "
    "clear beats elaborate and cluttered). EVERY element must connect to at least one other via "
    "an arrow — never plan a box that just sits there disconnected from the rest.\n"
    "3. RELATIONSHIPS — how the elements connect: arrows, grouping, hierarchy, or sequence. "
    "Plan a layout where connecting lines don't cross each other — e.g. if two branches "
    "reconverge, route them into a single box directly below/between them rather than crossing "
    "diagonals into swapped positions.\n"
    "4. LAYOUT — for diagrams, this will be drawn on a 0-10 by 0-10 grid, so describe roughly "
    "where each element sits in that space (e.g. 'Sensor at top-left around (1,7), flowing "
    "right to Preprocessing around (4,7), down to Model around (4,3)') with enough spacing "
    "that nothing overlaps.\n\n"
    "Never invent numeric data that isn't actually given to you — prefer a conceptual diagram "
    "over a fabricated chart. Output your plan as short labeled bullet points (TYPE / ELEMENTS "
    "/ RELATIONSHIPS / LAYOUT). Do not write any code."
)


def plan_section_figure_concept(topic, section_title, section_text, is_graphical_abstract=False,
                                model=None, api_base=None, api_key=None, num_ctx=None):
    """Reasoning step BEFORE any code is written: decide figure type, exact elements, their
    relationships, and layout. Returns the plan text, or None on failure (caller falls back to
    letting the code-writing step decide inline)."""
    framing = (
        "This is the GRAPHICAL ABSTRACT for the survey as a whole — plan a big-picture visual "
        "overview (e.g. a pipeline or scope diagram) summarizing what the survey covers, not a "
        "narrow diagram of one detail."
        if is_graphical_abstract else
        "Base the plan on what this section actually discusses."
    )
    user_msg = (
        f"Survey topic: {topic}\n"
        f"Section: {section_title}\n"
        f"{framing}\n\n"
        f"Section text:\n{section_text[:2500]}\n\n"
        f"Plan the figure now."
    )
    msgs = [{"role": "system", "content": FIGURE_PLAN_SYSTEM},
            {"role": "user", "content": user_msg}]
    try:
        chunks = list(stream_chat(msgs, model, api_base=api_base, api_key=api_key,
                                  temperature=0.5, num_ctx=num_ctx, max_tokens=500))
        plan = "".join(chunks).strip()
        return plan or None
    except Exception:
        return None


def plan_section_figure_code(topic, section_title, section_text, model, figure_plan=None,
                             is_graphical_abstract=False, api_base=None, api_key=None, num_ctx=None):
    framing = (
        "This is the GRAPHICAL ABSTRACT for the survey as a whole."
        if is_graphical_abstract else
        "Base the figure on what this section actually discusses."
    )
    plan_block = (
        f"FIGURE PLAN (implement exactly this — do not add extra elements beyond it):\n{figure_plan}\n\n"
        if figure_plan else
        "(No separate plan was made — decide the figure type and elements yourself, keeping it "
        "simple: 3-6 elements.)\n\n"
    )
    user_msg = (
        f"Survey topic: {topic}\n"
        f"Section: {section_title}\n"
        f"{framing}\n\n"
        f"{plan_block}"
        f"Section text (for reference only):\n{section_text[:1500]}\n\n"
        f"Write the matplotlib code for this figure now, implementing the plan precisely. "
        f"Before finishing, mentally re-check that every variable you use (ax, np, mpatches, "
        f"etc.) was actually assigned earlier in the code."
    )
    msgs = [{"role": "system", "content": FIGURE_CODE_SYSTEM},
            {"role": "user", "content": user_msg}]
    try:
        chunks = list(stream_chat(msgs, model, api_base=api_base, api_key=api_key,
                                  temperature=0.3, num_ctx=num_ctx, max_tokens=2200))
        return extract_code_block("".join(chunks))
    except Exception:
        return None


def _looks_like_truncation(error_msg):
    """Heuristic: these SyntaxError messages almost always mean the code got cut off by the
    token limit mid-expression, not a genuine logic bug — so the fix should be 'write less',
    not 'fix the same amount of code'."""
    if not error_msg:
        return False
    signals = ("was never closed", "unexpected eof", "eol while scanning",
              "unexpected end of file", "invalid syntax")
    low = error_msg.lower()
    return "syntaxerror" in low and any(s in low for s in signals)


def generate_section_figure(topic, section_title, section_text, model, is_graphical_abstract=False,
                            api_base=None, api_key=None, num_ctx=None, max_attempts=3, status_cb=None):
    """Plan the figure (type, elements, relationships, layout — a reasoning step with no code),
    then write matplotlib code implementing that plan, execute it, and retry (feeding the error
    back) if it fails. Returns a fig dict, or None if all attempts fail."""
    if status_cb:
        status_cb(f"🧠 Thinking about the best figure for “{section_title}”…")
    figure_plan = plan_section_figure_concept(
        topic, section_title, section_text, is_graphical_abstract=is_graphical_abstract,
        model=model, api_base=api_base, api_key=api_key, num_ctx=num_ctx)
    if status_cb and figure_plan:
        status_cb(f"🧠 Plan: {figure_plan.splitlines()[0][:120]}…")

    code, last_error = None, None
    for attempt in range(1, max_attempts + 1):
        if attempt == 1:
            code = plan_section_figure_code(
                topic, section_title, section_text, model, figure_plan=figure_plan,
                is_graphical_abstract=is_graphical_abstract,
                api_base=api_base, api_key=api_key, num_ctx=num_ctx)
        else:
            if _looks_like_truncation(last_error):
                fix_msg = (
                    f"Your previous code was cut off before it finished (error: {last_error}) — "
                    f"it was too long to complete within the response budget. Write a MUCH "
                    f"SIMPLER version of the SAME idea: fewer elements (3-5 boxes/nodes max), "
                    f"shorter labels, less code overall, so it completes fully this time. Output "
                    f"the caption comment line plus ONLY the new fenced Python code block."
                )
            else:
                fix_msg = (
                    f"Your previous code failed with this error:\n{last_error}\n\n"
                    f"Previous code:\n```python\n{code}\n```\n\n"
                    f"If this used raw matplotlib.patches or a manually-created axes variable, "
                    f"switch to the diagram toolkit instead (new_diagram_fig / add_box / "
                    f"add_arrow) — it avoids this whole class of mistake. Otherwise carefully "
                    f"trace through the code and fix the specific error. Output the caption "
                    f"comment line plus ONLY the corrected fenced Python code block."
                )
            try:
                chunks = list(stream_chat(
                    [{"role": "system", "content": FIGURE_CODE_SYSTEM},
                     {"role": "user", "content": fix_msg}],
                    model, api_base=api_base, api_key=api_key, temperature=0.3,
                    num_ctx=num_ctx, max_tokens=2200))
                code = extract_code_block("".join(chunks))
            except Exception as e:
                last_error = str(e)
                continue

        if not code:
            last_error = "model returned no code"
            continue

        if status_cb:
            status_cb(f"🎨 Rendering figure for “{section_title}”"
                      + (f" (attempt {attempt})" if attempt > 1 else "") + "…")
        result, err = execute_figure_code(code)
        if result:
            result["title"] = section_title
            result["caption"] = _extract_figure_caption(code) or f"Illustration for {section_title}."
            return result
        last_error = err
        if status_cb:
            note = "retrying" if attempt < max_attempts else "skipping this figure"
            status_cb(f"⚠️ Figure code failed ({err}) — {note}…")
    return None



def materialize_markdown_images(md_with_markers, figs):
    """Replace [[FIGURE:n]] markers with real embedded (base64 data-URI) images, for the
    live preview and the downloadable .md file. The marker-based source itself (unmodified)
    is what feeds the PDF renderer and the LaTeX converter instead."""
    def _sub(m):
        idx = int(m.group(1))
        if 0 <= idx < len(figs):
            b64 = base64.b64encode(figs[idx]["png"]).decode("ascii")
            return f"![Figure {idx + 1}: {figs[idx]['title']}](data:image/png;base64,{b64})"
        return ""
    return re.sub(r"\[\[FIGURE:(\d+)\]\]", _sub, md_with_markers)


# =============================================================================
# LaTeX EXPORT — converts the same marker-based survey markdown into a compilable
# LaTeX article, bundled as a .zip with the figure PNGs and a .bib file (a bare
# .tex can't carry embedded images, so the zip is the actual usable deliverable).
# =============================================================================
def _latex_escape(s):
    out = []
    for ch in s:
        out.append({
            "\\": r"\textbackslash{}", "&": r"\&", "%": r"\%", "$": r"\$", "#": r"\#",
            "_": r"\_", "{": r"\{", "}": r"\}", "~": r"\textasciitilde{}",
            "^": r"\textasciicircum{}",
        }.get(ch, ch))
    return "".join(out)


def _inline_to_latex(text):
    t = _latex_escape(text)
    t = re.sub(r"\*\*(.+?)\*\*", r"\\textbf{\1}", t)
    t = re.sub(r"(?<!\*)\*(?!\*)(.+?)(?<!\*)\*(?!\*)", r"\\textit{\1}", t)
    t = re.sub(r"`(.+?)`", r"\\texttt{\1}", t)
    return t


def survey_to_latex(md_with_markers, title, figs, fig_filenames):
    """Convert the marker-based survey markdown into a LaTeX article body. `fig_filenames` maps
    figure index -> filename as bundled in the zip (e.g. {0: 'fig0.png'})."""
    lines = md_with_markers.split("\n")
    body, in_itemize, in_refs, refs_started = [], False, False, False
    n = len(lines)
    i = 0

    def close_itemize():
        nonlocal in_itemize
        if in_itemize:
            body.append(r"\end{itemize}")
            in_itemize = False

    while i < n:
        line = lines[i].rstrip()
        if not line.strip():
            i += 1
            continue
        fig_m = re.match(r"^\[\[FIGURE:(\d+)\]\]$", line)
        if line.startswith("# "):
            close_itemize()
            i += 1
            continue  # main title is emitted separately via \maketitle
        elif line.startswith("## "):
            close_itemize()
            txt = line[3:]
            in_refs = txt.strip().lower().startswith("references") or txt.strip().lower().startswith("sources")
            body.append(r"\section{" + _inline_to_latex(txt) + "}")
            if in_refs:
                body.append(r"\begin{thebibliography}{99}")
                refs_started = True
        elif line.startswith("### "):
            close_itemize()
            body.append(r"\subsection{" + _inline_to_latex(line[4:]) + "}")
        elif fig_m:
            close_itemize()
            idx = int(fig_m.group(1))
            caption = ""
            j = i + 1
            while j < n and not lines[j].strip():
                j += 1
            if j < n:
                cap_m = re.match(r"^\*(.+)\*$", lines[j].strip())
                if cap_m:
                    caption = cap_m.group(1)
                    i = j
            fname = fig_filenames.get(idx, f"fig{idx}.png")
            body.append(r"\begin{figure}[h]\centering")
            body.append(r"\includegraphics[width=0.8\textwidth]{" + fname + "}")
            if caption:
                body.append(r"\caption{" + _inline_to_latex(caption) + "}")
            body.append(r"\end{figure}")
        elif line.startswith("---"):
            body.append(r"\bigskip")
        elif re.match(r"^\*[^*].*\*$", line):
            close_itemize()
            body.append(r"\textit{" + _inline_to_latex(line.strip("*")) + "}")
        elif line.startswith("- ") or line.startswith("• "):
            if not in_itemize:
                body.append(r"\begin{itemize}")
                in_itemize = True
            body.append(r"\item " + _inline_to_latex(line[2:]))
        elif in_refs and re.match(r"^\[\d+\]", line):
            m = re.match(r"^\[(\d+)\]\s*(.*)$", line)
            num, rest = m.group(1), m.group(2)
            url_m = re.search(r"(https?://\S+)$", rest)
            if url_m:
                text_part, url = rest[:url_m.start()].strip(), url_m.group(1)
                body.append(r"\bibitem{ref" + num + "} " + _inline_to_latex(text_part) +
                           r" \url{" + url + "}")
            else:
                body.append(r"\bibitem{ref" + num + "} " + _inline_to_latex(rest))
        else:
            close_itemize()
            body.append(_inline_to_latex(line))
        i += 1

    close_itemize()
    if refs_started:
        body.append(r"\end{thebibliography}")

    preamble = (
        "\\documentclass[11pt]{article}\n"
        "\\usepackage[margin=1in]{geometry}\n"
        "\\usepackage{graphicx}\n"
        "\\usepackage{hyperref}\n"
        "\\usepackage{parskip}\n"
        "\\title{" + _inline_to_latex(title) + "}\n"
        "\\date{\\today}\n"
        "\\begin{document}\n"
        "\\maketitle\n"
    )
    return preamble + "\n".join(body) + "\n\\end{document}\n"


def build_latex_zip(tex_str, figs, fig_filenames, bib_str=None, tex_filename="document.tex"):
    """Bundle the .tex source with its figure PNGs (and optional .bib) into a downloadable
    zip — a bare .tex file can't carry embedded images on its own."""
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zf:
        zf.writestr(tex_filename, tex_str)
        for idx, fig in enumerate(figs):
            zf.writestr(fig_filenames.get(idx, f"fig{idx}.png"), fig["png"])
        if bib_str:
            zf.writestr("references.bib", bib_str)
    return buf.getvalue()


# =============================================================================
# WRITE PAPER — a guided pipeline: idea -> literature search -> source curation ->
# novelty check -> outline -> section-by-section writing (with figures) -> AI
# fine-tuning -> export. Reuses the search/novelty/figure/export machinery above.
# =============================================================================
DEFAULT_PAPER_SECTIONS = [
    {"title": "Abstract", "description": "A concise, standalone summary of the problem, approach, and contribution.", "figure": False},
    {"title": "1. Introduction", "description": "Motivate the problem, state the contribution, and outline the paper's structure.", "figure": True},
    {"title": "2. Related Work", "description": "Survey the closest prior work from the gathered literature and position this paper relative to it.", "figure": False},
    {"title": "3. Proposed Approach", "description": "Describe the proposed idea/method in detail.", "figure": True},
    {"title": "4. Methodology", "description": "Describe the experimental setup, data, and evaluation plan.", "figure": True},
    {"title": "5. Results and Discussion", "description": "Present and discuss the expected or obtained results.", "figure": True},
    {"title": "6. Conclusion", "description": "Summarize the contribution and outline future work.", "figure": False},
]

PAPER_OUTLINE_SYSTEM = (
    "You are an expert academic writer proposing a structure for a research paper. Given the "
    "author's idea and a digest of the related literature gathered, suggest a strong paper "
    "TITLE and a section-by-section outline.\n\n"
    'Respond with ONLY a JSON object in this exact shape:\n'
    '{"title": "...", "sections": [{"title": "Abstract", "description": "1-2 sentences on what '
    'this section should cover"}, ...]}\n\n'
    "Choose sections appropriate to the idea and field — a typical paper might include Abstract, "
    "Introduction, Related Work, Proposed Method, Experimental Setup, Results, Discussion, "
    "Conclusion — but adapt this to what actually fits the idea; don't force sections that "
    "don't make sense for it. Output ONLY the JSON object — no markdown fences, no explanation."
)


def suggest_paper_outline(idea_text, digest, model, api_base=None, api_key=None, num_ctx=None):
    """Ask the LLM for a title + section outline. Returns (title, sections) where sections is a
    list of {"title", "description"} dicts; falls back to DEFAULT_PAPER_SECTIONS on failure."""
    user_msg = (f"Idea:\n{idea_text}\n\nRelated literature (numbered sources):\n{digest}\n\n"
               f"Propose the title and outline now.")
    msgs = [{"role": "system", "content": PAPER_OUTLINE_SYSTEM},
            {"role": "user", "content": user_msg}]
    try:
        chunks = list(stream_chat(msgs, model, api_base=api_base, api_key=api_key,
                                  temperature=0.5, num_ctx=num_ctx, max_tokens=1200))
        raw = re.sub(r"^```(json)?\s*|\s*```$", "", "".join(chunks).strip(), flags=re.I)
        data = json.loads(raw)
        title = (data.get("title") or "Untitled Paper").strip()
        sections = []
        for s in data.get("sections", []):
            if isinstance(s, dict) and s.get("title"):
                sections.append({"title": s["title"].strip(), "description": (s.get("description") or "").strip()})
        if not sections:
            raise ValueError("no sections returned")
        return title, sections
    except Exception:
        return "Untitled Paper", [dict(s) for s in DEFAULT_PAPER_SECTIONS]


PAPER_SECTION_SYSTEM = (
    "You are an expert academic writing ONE section of a research paper, following the given "
    "outline description for this section. Write formal, in-depth academic prose. Cite sources "
    "by their number, e.g. [1], [3] — EXCEPT in the Abstract, which must stand alone with NO "
    "citation markers at all. Use ONLY facts present in the provided numbered sources, or "
    "reasonable statements about the author's own idea/approach — never invent results, "
    "numbers, or citations that aren't grounded in what's given. Output only the requested "
    "section's body text (no heading — that is added for you)."
)


PAPER_FINETUNE_SYSTEM = (
    "You are helping revise an academic paper draft. You are given the CURRENT full paper "
    "(Markdown, which may include [[FIGURE:n]] markers) and a REQUEST describing what to "
    "change. Preserve every existing [[FIGURE:n]] marker exactly as-is, in place, unless the "
    "request specifically asks to remove that section or figure — do not renumber or invent "
    "new ones. Output the COMPLETE, revised paper as a SINGLE fenced Markdown code block — the "
    "entire paper, never a snippet, diff, or 'rest unchanged' placeholder. Do not put any "
    "explanation outside the code block."
)


def finetune_paper(current_md, instruction, model, api_base=None, api_key=None, num_ctx=None, max_tokens=4000):
    """Ask the LLM to revise the full paper per a free-text instruction. Returns the revised
    markdown (with [[FIGURE:n]] markers preserved), or None on failure."""
    user_msg = f"CURRENT PAPER:\n```markdown\n{current_md}\n```\n\nREQUEST: {instruction}"
    msgs = [{"role": "system", "content": PAPER_FINETUNE_SYSTEM},
            {"role": "user", "content": user_msg}]
    try:
        chunks = list(stream_chat(msgs, model, api_base=api_base, api_key=api_key,
                                  temperature=0.4, num_ctx=num_ctx, max_tokens=max_tokens))
        revised = extract_code_block("".join(chunks))
        return revised or None
    except Exception:
        return None


SECTION_FINETUNE_SYSTEM = (
    "You are revising ONE section of an academic paper. You are given the CURRENT text of just "
    "this section and a REQUEST describing what to change. Output ONLY the complete revised "
    "section text — no heading, no explanation, no markdown code fences, nothing else."
)


def finetune_section(section_text, instruction, model, api_base=None, api_key=None, num_ctx=None, max_tokens=1500):
    """Ask the LLM to revise just one section's text per a free-text instruction. Returns the
    revised text, or None on failure. Scoped to one section, so it's faster/cheaper than
    rewriting the whole paper for a small, localized change."""
    user_msg = f"CURRENT SECTION TEXT:\n{section_text}\n\nREQUEST: {instruction}"
    msgs = [{"role": "system", "content": SECTION_FINETUNE_SYSTEM},
            {"role": "user", "content": user_msg}]
    try:
        chunks = list(stream_chat(msgs, model, api_base=api_base, api_key=api_key,
                                  temperature=0.4, num_ctx=num_ctx, max_tokens=max_tokens))
        revised = "".join(chunks).strip()
        return revised or None
    except Exception:
        return None


def rebuild_paper_markdown(title, sections_data, refs, generated_date, source_count):
    """Reassemble the full paper's marker-based markdown from structured section data. Used
    both at initial generation and every time a section is individually edited/fine-tuned
    afterward, so the two stay in sync."""
    parts = [f"# {title}\n", f"*Generated on {generated_date}, based on {source_count} source(s).*\n"]
    for sec in sections_data:
        parts.append(f"## {sec['title']}\n\n{sec['text']}\n")
        if sec.get("figure_idx") is not None:
            idx = sec["figure_idx"]
            parts.append(f"[[FIGURE:{idx}]]\n\n*Figure {idx + 1}. {sec.get('figure_caption', '')}*\n")
    parts.append("## References\n\n" + refs)
    return "\n".join(parts)


# =============================================================================
# WEB SEARCH (DuckDuckGo, no API key) + live-time API + page reader
#   pip install ddgs
# =============================================================================
def _strip_html(html):
    html = re.sub(r"(?is)<(script|style|noscript|header|footer|nav).*?</\1>", " ", html)
    text = re.sub(r"(?s)<[^>]+>", " ", html)
    text = re.sub(r"&[a-zA-Z#0-9]+;", " ", text)
    return " ".join(text.split())


def fetch_page_text(url, chars=2500):
    try:
        resp = requests.get(url, timeout=15, verify=False, headers=BROWSER_HEADERS)
        if resp.status_code != 200:
            return f"[page returned HTTP {resp.status_code} — likely blocking bots]"
        text = _strip_html(resp.text)
        return text[:chars] if text.strip() else "[no readable text — page content is JS-rendered]"
    except Exception as e:
        return f"[could not read page: {e}]"


def get_live_time(query):
    """For 'what time in X' questions: read a time API (real value is in the body,
    unlike JS clock sites) or compute locally. Returns None if not a time question."""
    t = query.lower()
    if not any(w in t for w in ("what time", "current time", "time now", "time in",
                                "whats the time", "what's the time", "local time")):
        return None
    zones = {
        "egypt": "Africa/Cairo", "cairo": "Africa/Cairo", "alexandria": "Africa/Cairo",
        "assiut": "Africa/Cairo", "asyut": "Africa/Cairo",
        "saudi": "Asia/Riyadh", "riyadh": "Asia/Riyadh", "jeddah": "Asia/Riyadh",
        "mecca": "Asia/Riyadh", "kaust": "Asia/Riyadh",
        "uae": "Asia/Dubai", "dubai": "Asia/Dubai",
        "uk": "Europe/London", "london": "Europe/London",
        "paris": "Europe/Paris", "france": "Europe/Paris",
        "new york": "America/New_York", "nyc": "America/New_York",
        "los angeles": "America/Los_Angeles", "california": "America/Los_Angeles",
        "tokyo": "Asia/Tokyo", "japan": "Asia/Tokyo", "utc": "UTC", "gmt": "UTC",
    }
    tz = place = None
    for key, zone in zones.items():
        if key in t:
            tz, place = zone, key.title(); break
    if not tz:
        return None
    for api in (f"https://timeapi.io/api/Time/current/zone?timeZone={tz}",
                f"https://worldtimeapi.org/api/timezone/{tz}"):
        try:
            r = requests.get(api, timeout=10, headers=BROWSER_HEADERS)
            if r.status_code == 200:
                data = r.json()
                if "dateTime" in data:
                    return f"The current time in {place} ({tz}) is {data['dateTime'].replace('T',' ')[:19]}. (timeapi.io)"
                if "datetime" in data:
                    return f"The current time in {place} ({tz}) is {data['datetime'].replace('T',' ')[:19]}. (worldtimeapi.org)"
        except Exception:
            continue
    try:
        from zoneinfo import ZoneInfo
        now = datetime.now(ZoneInfo(tz))
        return f"The current time in {place} ({tz}) is {now:%H:%M:%S on %A, %d %B %Y} (computed locally)."
    except Exception:
        return None


def web_search(query, max_results=5, read_pages=False, chars_per_page=2500):
    """Return (results, error). Each result: {title,url,snippet,page_text}."""
    try:
        from ddgs import DDGS
    except ImportError:
        try:
            from duckduckgo_search import DDGS
        except ImportError:
            return None, "Install the search package first:  pip install ddgs"
    results = []
    try:
        with DDGS() as ddgs:
            for r in ddgs.text(query, max_results=max_results):
                results.append({
                    "title": r.get("title") or r.get("heading") or "",
                    "url": r.get("href") or r.get("url") or r.get("link") or "",
                    "snippet": r.get("body") or r.get("snippet") or r.get("description") or "",
                    "page_text": "",
                })
    except Exception as e:
        return None, f"Search failed: {e}"
    if read_pages:
        for item in results:
            if item["url"]:
                item["page_text"] = fetch_page_text(item["url"], chars=chars_per_page)
    return results, None


def results_to_context(results, use_pages=False):
    blocks = []
    for i, r in enumerate(results, 1):
        pt = r["page_text"]
        usable = use_pages and pt and not pt.startswith("[")
        body = pt if usable else r["snippet"]
        blocks.append(f"[{i}] {r['title']}\nURL: {r['url']}\n{body}")
    return "\n\n".join(blocks)


SEARCH_PLANNER_SYSTEM = (
    "You are a search-query planner for an AI chat assistant that has web search available. "
    "You are given the recent conversation and the user's latest message. Decide whether a web "
    "search would actually help answer it (e.g. current events, facts you can't be sure of, "
    "prices, specific people/products/versions, anything time-sensitive).\n\n"
    "Respond in EXACTLY this two-line format, nothing else, QUERY on the FIRST line always:\n"
    "QUERY: <the search query (2-6 words, resolve any pronouns/references using the "
    "conversation), or NONE if a search would not help — e.g. small talk, opinions, coding "
    "help, math, writing help, general knowledge that doesn't change over time, or the answer "
    "is already in the conversation>\n"
    "REASON: <one short sentence explaining your thinking>"
)


def _parse_reason_query(text):
    """Parse a 'QUERY: ...\\nREASON: ...' reply. Returns (query_or_None, reason_str).
    QUERY must come first in the prompt format specifically so that if the model's output gets
    cut off by the token limit, it's the (less important) REASON line that gets truncated away —
    never the query itself, which used to cause searches to silently get skipped."""
    reason, query = "", None
    m_r = re.search(r"REASON:\s*(.+)", text, re.I)
    m_q = re.search(r"QUERY:\s*(.+)", text, re.I)
    if m_r:
        reason = m_r.group(1).strip().splitlines()[0].strip()
    if m_q:
        q = m_q.group(1).strip().splitlines()[0].strip().strip('"').strip("'")
        if q and not q.upper().startswith("NONE"):
            query = q[:200]
    if not m_r and not m_q:
        # Model ignored the format and just replied with a bare query/NONE — fall back gracefully.
        bare = text.strip().strip('"').strip("'").splitlines()[0].strip() if text.strip() else ""
        if bare and not bare.upper().startswith("NONE"):
            query = bare[:200]
    return query, reason


def plan_search_query(user_message, recent_history, model, api_base=None, api_key=None, num_ctx=None):
    """Ask the LLM what to search for given the user's message and recent context.
    Returns (query_or_None, reason_str)."""
    convo = ""
    if recent_history:
        lines = []
        for m in recent_history[-6:]:
            role = "User" if m["role"] == "user" else "Assistant"
            lines.append(f"{role}: {m['content'][:300]}")
        convo = "\n".join(lines)
    user_block = (f"Recent conversation:\n{convo}\n\n" if convo else "") + \
                 f"Latest user message: {user_message}"
    msgs = [{"role": "system", "content": SEARCH_PLANNER_SYSTEM},
            {"role": "user", "content": user_block}]
    try:
        chunks = list(stream_chat(msgs, model, api_base=api_base, api_key=api_key,
                                  temperature=0.0, num_ctx=num_ctx, max_tokens=200))
        raw = "".join(chunks).strip()
    except Exception as e:
        return None, f"(planner error: {e})"
    return _parse_reason_query(raw)


SEARCH_FOLLOWUP_SYSTEM = (
    "You are a research assistant deciding whether another web search is needed. You are given "
    "the user's question and the web search results gathered so far this turn.\n\n"
    "Respond in EXACTLY this two-line format, nothing else, QUERY on the FIRST line always:\n"
    "QUERY: <a new search query (2-6 words) targeting the gap — a different angle than what was "
    "already tried, not a repeat — or NONE if the results already let you fully and accurately "
    "answer the question. Don't chase minor details forever; once the core question is "
    "answerable, use NONE>\n"
    "REASON: <one short sentence — what's still missing, or why the results are already enough>"
)


def plan_followup_query(user_message, tried_queries, accumulated_results_text, model,
                         api_base=None, api_key=None, num_ctx=None):
    """After one or more search rounds, ask the LLM if it needs to search again (and for what).
    Returns (query_or_None, reason_str)."""
    tried = ", ".join(f'"{q}"' for q in tried_queries) if tried_queries else "(none)"
    user_block = (
        f"User's question: {user_message}\n\n"
        f"Queries already tried: {tried}\n\n"
        f"Search results gathered so far:\n{accumulated_results_text or '(none)'}"
    )
    msgs = [{"role": "system", "content": SEARCH_FOLLOWUP_SYSTEM},
            {"role": "user", "content": user_block}]
    try:
        chunks = list(stream_chat(msgs, model, api_base=api_base, api_key=api_key,
                                  temperature=0.0, num_ctx=num_ctx, max_tokens=200))
        raw = "".join(chunks).strip()
    except Exception as e:
        return None, f"(planner error: {e})"
    return _parse_reason_query(raw)


def run_iterative_search(user_message, recent_history, model, api_base=None, api_key=None,
                         num_ctx=None, max_results=5, read_pages=False, max_rounds=3,
                         status_cb=None):
    """Search, let the model judge if that's enough, and if not search again with a new
    query — up to max_rounds times — before handing back combined context for the final answer.
    Surfaces the model's reasoning at each decision point via status_cb(text), so the person
    can see WHY it chose to search, what for, and why it stopped.

    Returns (combined_context, rounds) where rounds is a list of entries:
      {"type": "skip",   "reason": str}                                    — no search needed at all
      {"type": "search", "query": str, "reason": str, "results": [...], "error": str|None}
      {"type": "stop",   "reason": str}                                    — decided results are enough
    in the order they happened."""
    rounds = []
    tried_queries = []
    combined_blocks = []

    query, reason = plan_search_query(user_message, recent_history, model,
                                      api_base=api_base, api_key=api_key, num_ctx=num_ctx)
    if status_cb and reason:
        status_cb(f"💭 {reason}")

    if not query:
        rounds.append({"type": "skip", "reason": reason or "No search needed."})
        return "", rounds

    round_n = 0
    while query and round_n < max_rounds:
        round_n += 1
        if status_cb:
            status_cb(f"🔎 Round {round_n}: searching \u201c{query}\u201d…")
        results, err = web_search(query, max_results=max_results, read_pages=read_pages)
        tried_queries.append(query)
        if results:
            block = f'--- Results for "{query}" ---\n' + results_to_context(results, use_pages=read_pages)
            combined_blocks.append(block)
            rounds.append({"type": "search", "query": query, "reason": reason,
                           "results": results, "error": None})
            if status_cb:
                status_cb(f"✅ Found {len(results)} result(s) for \u201c{query}\u201d")
        else:
            rounds.append({"type": "search", "query": query, "reason": reason,
                           "results": [], "error": err})
            if status_cb:
                status_cb(f"🔇 Search failed for \u201c{query}\u201d: {err}")

        if round_n >= max_rounds:
            if status_cb:
                status_cb(f"⏹️ Reached the {max_rounds}-round limit — answering with what's been gathered.")
            break

        if status_cb:
            status_cb("💭 Checking whether that's enough to answer…")
        query, reason = plan_followup_query(user_message, tried_queries, "\n\n".join(combined_blocks),
                                            model, api_base=api_base, api_key=api_key, num_ctx=num_ctx)
        if not query:
            rounds.append({"type": "stop", "reason": reason or "Enough information gathered."})
            if status_cb and reason:
                status_cb(f"💭 {reason}")
        elif status_cb and reason:
            status_cb(f"💭 {reason}")

    return "\n\n".join(combined_blocks), rounds


# =============================================================================
# PERPLEXITY-STYLE ANSWERING — flatten search rounds into one globally-numbered,
# deduped source list (per-round numbering restarts at [1] each round, which
# collides once you have 2+ rounds), so the final answer's [n] citations match
# one consistent source list, and so we can render clickable source cards.
# =============================================================================
def flatten_search_results(rounds):
    """Flatten all 'search' rounds into one deduped list, in the order first seen."""
    flat = []
    seen = set()
    for rd in rounds:
        if rd.get("type") != "search":
            continue
        for r in rd.get("results", []):
            key = r.get("url") or r.get("title", "")
            if key and key not in seen:
                seen.add(key)
                flat.append(r)
    return flat


def build_cited_context(flat_results, use_pages=False, per_chars=2000):
    """Numbered source text, numbered to match render_source_cards — for injecting into the
    system prompt alongside a citation instruction."""
    blocks = []
    for i, r in enumerate(flat_results, 1):
        pt = r.get("page_text", "")
        usable = use_pages and pt and not str(pt).startswith("[")
        body = (pt if usable else r.get("snippet", "") or "")[:per_chars]
        blocks.append(f"[{i}] {r.get('title', '')}\nURL: {r.get('url', '')}\n{body}")
    return "\n\n".join(blocks)


def _domain_of(url):
    try:
        from urllib.parse import urlparse
        d = urlparse(url).netloc
        return d[4:] if d.startswith("www.") else d
    except Exception:
        return url or ""


def render_source_cards(flat_results):
    """Render a row of small clickable source cards (favicon + [n] + domain + title),
    matching the [n] citations the model was told to use — the Perplexity-style source strip
    shown under an answer."""
    if not flat_results:
        return
    n_cols = min(len(flat_results), 4)
    cols = st.columns(n_cols)
    for i, r in enumerate(flat_results):
        col = cols[i % n_cols]
        url = r.get("url", "") or "#"
        domain = _domain_of(url)
        favicon = f"https://www.google.com/s2/favicons?sz=32&domain={domain}" if domain else ""
        title = (r.get("title") or domain or "source").strip()
        title = title[:65] + "…" if len(title) > 65 else title
        with col:
            st.markdown(
                f'<a href="{url}" target="_blank" style="text-decoration:none;color:inherit;">'
                f'<div style="border:1px solid rgba(128,128,128,0.35);border-radius:10px;'
                f'padding:8px 10px;margin-bottom:8px;font-size:0.8em;min-height:70px;">'
                f'<div style="display:flex;align-items:center;gap:6px;opacity:0.7;">'
                f'<img src="{favicon}" width="14" height="14" style="border-radius:3px;" '
                f'onerror="this.style.display=\'none\'"/><span>[{i + 1}] {domain}</span></div>'
                f'<div style="margin-top:4px;font-weight:500;line-height:1.3;">{title}</div>'
                f'</div></a>',
                unsafe_allow_html=True)


FOLLOWUP_SUGGEST_SYSTEM = (
    "Given the user's question and the assistant's answer, suggest 3 short, natural follow-up "
    "questions the user might want to ask next — the kind a curious reader would click to dig "
    "deeper or explore a related angle. Respond with ONLY the 3 questions, one per line, no "
    "numbering, no bullets, no quotes, no extra text."
)


def suggest_followups(user_message, assistant_answer, model, api_base=None, api_key=None, num_ctx=None):
    """Ask the model for a few related follow-up questions, Perplexity-style. Returns a list
    of up to 3 question strings (empty list on failure)."""
    msgs = [{"role": "system", "content": FOLLOWUP_SUGGEST_SYSTEM},
            {"role": "user", "content": f"User asked: {user_message}\n\nAssistant answered: {assistant_answer[:1500]}"}]
    try:
        chunks = list(stream_chat(msgs, model, api_base=api_base, api_key=api_key,
                                  temperature=0.6, num_ctx=num_ctx, max_tokens=120))
        raw = "".join(chunks).strip()
    except Exception:
        return []
    qs = [ln.strip(" -•*\t") for ln in raw.splitlines() if ln.strip()]
    return qs[:3]


# =============================================================================
# ITERATIVE PAPER SEARCH — same idea as the chat search loop above, but for
# academic literature: turn the user's topic into good search terms, run it
# across whichever sources are enabled, then judge whether the papers found
# are relevant/sufficient — searching again with different terms if not.
# Used by both the Research Crew tab and the Make a Survey tab.
# =============================================================================
PAPER_QUERY_PLANNER_SYSTEM = (
    "You are a research-librarian assistant planning an academic literature search. You are "
    "given the user's topic or description, which may be casual, verbose, or loosely phrased. "
    "Turn it into the best possible search query for academic databases (arXiv, Semantic "
    "Scholar, OpenAlex) — precise technical/academic terminology, 2-6 words, no filler words "
    "like 'papers about' or 'research on'.\n\n"
    "Respond in EXACTLY this two-line format, nothing else, QUERY on the FIRST line always:\n"
    "QUERY: <the search query>\n"
    "REASON: <one short sentence on why you chose these terms>"
)


def plan_paper_query(topic, model, api_base=None, api_key=None, num_ctx=None):
    """Ask the LLM to turn a user's topic/description into good academic search terms.
    Returns (query_or_None, reason_str)."""
    msgs = [{"role": "system", "content": PAPER_QUERY_PLANNER_SYSTEM},
            {"role": "user", "content": f"Topic: {topic}"}]
    try:
        chunks = list(stream_chat(msgs, model, api_base=api_base, api_key=api_key,
                                  temperature=0.0, num_ctx=num_ctx, max_tokens=200))
        raw = "".join(chunks).strip()
    except Exception as e:
        return None, f"(planner error: {e})"
    return _parse_reason_query(raw)


PAPER_EVAL_SYSTEM = (
    "You are a research librarian reviewing literature-search results gathered so far. You are "
    "given the user's topic and the titles of papers found across all queries tried this "
    "session.\n\n"
    "Respond in EXACTLY this two-line format, nothing else, QUERY on the FIRST line always:\n"
    "QUERY: <a new, different search query targeting the gap (2-6 words), or NONE if the papers "
    "gathered are already relevant and sufficient for this topic. Don't chase exhaustive "
    "coverage forever — a solid, on-topic set is enough; use NONE once that's true.>\n"
    "REASON: <one short sentence — are these enough and relevant, or what's missing/off-topic?>"
)


def evaluate_papers_and_maybe_requery(topic, tried_queries, papers, model,
                                       api_base=None, api_key=None, num_ctx=None):
    """After one or more search rounds, ask the LLM whether the papers gathered so far are
    relevant/sufficient for the topic. Returns (new_query_or_None, reason_str)."""
    tried = ", ".join(f'"{q}"' for q in tried_queries) if tried_queries else "(none)"
    listing = "\n".join(f"- [{p.get('source', '?')}] {p.get('title', '')}" for p in papers[:40]) \
        or "(none found yet)"
    user_block = (
        f"Topic: {topic}\n\n"
        f"Queries already tried: {tried}\n\n"
        f"Papers found so far ({len(papers)} total):\n{listing}"
    )
    msgs = [{"role": "system", "content": PAPER_EVAL_SYSTEM},
            {"role": "user", "content": user_block}]
    try:
        chunks = list(stream_chat(msgs, model, api_base=api_base, api_key=api_key,
                                  temperature=0.0, num_ctx=num_ctx, max_tokens=200))
        raw = "".join(chunks).strip()
    except Exception as e:
        return None, f"(planner error: {e})"
    return _parse_reason_query(raw)


def run_iterative_paper_search(topic, model, search_once, api_base=None, api_key=None,
                               num_ctx=None, max_rounds=3, status_cb=None):
    """Plan a search query from the topic, run it via `search_once(query) -> (papers, errors)`,
    then judge whether the papers gathered are relevant/sufficient — searching again with new
    terms (up to max_rounds times) if not. Deduplicates by normalized title across rounds.

    Returns (all_papers, rounds) where rounds is a list of entries:
      {"type": "search", "query": str, "reason": str, "found": int, "added": int, "errors": [...]}
      {"type": "stop", "reason": str}"""
    rounds = []
    tried_queries = []
    all_papers = []
    seen_keys = set()

    query, reason = plan_paper_query(topic, model, api_base=api_base, api_key=api_key, num_ctx=num_ctx)
    if status_cb and reason:
        status_cb(f"💭 {reason}")
    if not query:
        query = topic.strip()
        if status_cb:
            status_cb("💭 Using your topic as-is for the first search.")

    round_n = 0
    while query and round_n < max_rounds:
        round_n += 1
        if status_cb:
            status_cb(f"🔎 Round {round_n}: searching \u201c{query}\u201d…")
        new_papers, errors = search_once(query)
        tried_queries.append(query)
        added = 0
        for p in new_papers:
            k = re.sub(r"\W+", "", p.get("title", "").lower())[:60]
            if k and k not in seen_keys:
                seen_keys.add(k)
                all_papers.append(p)
                added += 1
        rounds.append({"type": "search", "query": query, "reason": reason,
                       "found": len(new_papers), "added": added, "errors": errors})
        if status_cb:
            err_note = f" — ⚠️ {'; '.join(errors)}" if errors else ""
            status_cb(f"✅ Found {len(new_papers)} result(s), {added} new, for \u201c{query}\u201d{err_note}")

        if round_n >= max_rounds:
            if status_cb:
                status_cb(f"⏹️ Reached the {max_rounds}-round limit — using the {len(all_papers)} paper(s) gathered.")
            break

        if status_cb:
            status_cb("💭 Checking whether these papers are relevant/sufficient…")
        query, reason = evaluate_papers_and_maybe_requery(
            topic, tried_queries, all_papers, model,
            api_base=api_base, api_key=api_key, num_ctx=num_ctx)
        if not query:
            rounds.append({"type": "stop", "reason": reason or "Enough relevant papers gathered."})
            if status_cb and reason:
                status_cb(f"💭 {reason}")
        elif status_cb and reason:
            status_cb(f"💭 {reason}")

    return all_papers, rounds



# =============================================================================
def extract_pdf_text(file_obj, max_chars=24000):
    """Return (text, n_pages, truncated). file_obj is a Streamlit UploadedFile."""
    reader = PdfReader(file_obj)
    n_pages = len(reader.pages)
    raw = "\n".join((p.extract_text() or "") for p in reader.pages)
    text = " ".join(raw.split())
    truncated = len(text) > max_chars
    if truncated:
        text = text[:max_chars] + " ...[TRUNCATED — paper continues beyond this point]"
    return text, n_pages, truncated


REVIEWER_SYSTEM = (
    "You are an expert, rigorous but fair academic peer reviewer. You are given the "
    "extracted text of a paper. Write a structured review report using EXACTLY these "
    "numbered sections and headings, in this order:\n\n"
    "1. Summary — 2 to 3 sentences, in plain language, stating what the paper does and claims.\n"
    "2. Questions for the Authors — a numbered list of specific questions about anything "
    "unclear, unjustified, or missing (e.g. undefined terms, unstated assumptions, missing "
    "details needed to reproduce the work).\n"
    "3. Logical Gaps & Inconsistencies — point out claims not supported by the evidence "
    "presented, internal contradictions, overgeneralizations, missing baselines/controls, "
    "confounds, or statistical concerns. Quote or reference the specific claim/section.\n"
    "4. Strengths — what the paper does well.\n"
    "5. Weaknesses / Major Concerns — the most important problems that affect the conclusions.\n"
    "6. Minor Issues — wording, figures/tables, typos, presentation.\n"
    "7. Overall Assessment — one of: Accept / Minor revision / Major revision / Reject, "
    "followed by one or two sentences justifying it.\n\n"
    "IMPORTANT RULES: Base every point ONLY on the paper text provided. If something the "
    "reader would expect (dataset size, ablations, limitations, ethics, etc.) is not present, "
    "say it is 'not reported' rather than inventing it. Be specific and constructive — refer "
    "to the actual claims. Do not invent references, numbers, or results that are not in the text."
)


SUMMARIZE_SYSTEM = (
    "You are an expert research assistant. You are given the extracted text of ONE paper. "
    "Produce a clear, well-structured summary with these sections:\n\n"
    "**TL;DR** — 2 to 3 sentences capturing the whole paper in plain language.\n"
    "**Problem & Motivation** — what problem the paper addresses and why it matters.\n"
    "**Approach / Method** — what the authors actually did.\n"
    "**Key Results** — the main findings, with any specific numbers stated in the paper.\n"
    "**Contributions** — what is novel.\n"
    "**Limitations** — weaknesses or constraints the paper notes (or that are evident).\n"
    "**Takeaways** — 2 to 4 bullet points a reader should remember.\n\n"
    "Base everything ONLY on the provided text. If something expected is missing, say it is "
    "'not reported' rather than inventing it. Do not invent numbers, references, or results."
)

# System prompt for the follow-up chat about the paper (paper text is injected once).
SUMMARIZE_QA_SYSTEM = (
    "You are a helpful research assistant answering questions about ONE specific paper. The "
    "full text of the paper is provided below. Answer using ONLY the paper's content; if the "
    "answer is not in the paper, say so clearly rather than guessing. Be concise and precise, "
    "and quote or point to the relevant part when useful.\n\n"
    "----- PAPER TEXT -----\n{paper}\n----- END PAPER -----"
)


# =============================================================================
# CODING assistant
# =============================================================================
CODING_SYSTEM = (
    "You are an expert software engineer. You are given the current code (which may be empty "
    "if starting from scratch) and a request. Do the following:\n"
    "1. First, a short bulleted list of the changes you are making (or the plan, if new code).\n"
    "2. Then output the COMPLETE, updated, ready-to-run code in a SINGLE fenced code block "
    "(```language ... ```). The code block MUST contain the entire file — never a snippet, "
    "diff, or 'rest unchanged' placeholder.\n"
    "Keep existing working behavior unless asked to change it, fix bugs you notice, and write "
    "clean, commented code. Do not put explanations inside the code block."
)

CODE_LANGS = {
    "python": "py", "javascript": "js", "typescript": "ts", "java": "java",
    "c": "c", "cpp": "cpp", "csharp": "cs", "go": "go", "rust": "rs", "php": "php",
    "ruby": "rb", "swift": "swift", "kotlin": "kt", "html": "html", "css": "css",
    "sql": "sql", "bash": "sh", "json": "json", "yaml": "yaml", "text": "txt",
    "matlab": "m",
}
# map file extension -> language for auto-detect on upload
_EXT_TO_LANG = {v: k for k, v in CODE_LANGS.items()}
_EXT_TO_LANG.update({"jsx": "javascript", "tsx": "typescript", "h": "cpp", "hpp": "cpp",
                     "yml": "yaml", "htm": "html"})


def extract_code_block(text):
    """Pull the largest fenced code block from the model's reply (the full file).

    Tolerates a missing closing fence — this happens when the response gets cut
    off by an output-token limit before the model reaches ```; in that case we
    still grab everything after the opening fence instead of falling back to the
    raw (fenced) reply, which previously caused long files to be silently dropped."""
    closed = re.findall(r"```[a-zA-Z0-9_+\-]*\n(.*?)```", text, re.S)
    if closed:
        return max(closed, key=len).rstrip("\n")
    # No closing fence found — likely truncated mid-file. Take everything after
    # the LAST opening fence (the code block, not any earlier prose fences).
    opens = list(re.finditer(r"```[a-zA-Z0-9_+\-]*\n", text))
    if opens:
        return text[opens[-1].end():].rstrip("\n")
    return text.strip()


# =============================================================================
# AGENTIC CODING MODE — a Claude-Code-style loop that works with ANY configured
# LLM backend: read/write real files and run real shell commands in a chosen
# project folder. Since this grants genuine filesystem/shell access, two safety
# measures are non-negotiable: (1) every path is sandboxed to the chosen project
# folder — nothing can read/write/traverse outside it — and (2) every write_file
# and run_command action is shown to the user and requires explicit approval
# before it executes, unless they deliberately turn that off. No native
# function-calling API is used (support varies wildly across local/Ollama/custom
# OpenAI-compatible models) — instead the model emits one JSON action per turn,
# which works identically regardless of which provider is behind it.
# =============================================================================
_AGENT_IGNORE_DIRS = {".git", "node_modules", "__pycache__", ".venv", "venv",
                      ".mypy_cache", ".pytest_cache", "dist", "build", ".idea", ".vscode"}


def _agent_safe_path(project_dir, rel_path):
    """Resolve rel_path against project_dir and verify it doesn't escape it (blocks '../'
    traversal and absolute-path overrides). Raises ValueError if it would escape."""
    base = os.path.realpath(project_dir)
    target = os.path.realpath(os.path.join(base, rel_path or "."))
    if target != base and not target.startswith(base + os.sep):
        raise ValueError(f"path '{rel_path}' escapes the project folder — refused")
    return target


def agent_list_dir(project_dir, rel_path=".", max_entries=300):
    try:
        root = _agent_safe_path(project_dir, rel_path)
    except ValueError as e:
        return f"ERROR: {e}"
    if not os.path.isdir(root):
        return f"ERROR: '{rel_path}' is not a directory"
    lines = []
    for dirpath, dirnames, filenames in os.walk(root):
        dirnames[:] = [d for d in dirnames if d not in _AGENT_IGNORE_DIRS and not d.startswith(".")]
        rel_dir = os.path.relpath(dirpath, root)
        for fn in sorted(filenames):
            p = fn if rel_dir == "." else os.path.join(rel_dir, fn)
            lines.append(p)
            if len(lines) >= max_entries:
                lines.append("... (truncated)")
                return "\n".join(lines)
    return "\n".join(lines) if lines else "(empty directory)"


def agent_read_file(project_dir, rel_path, max_chars=20000):
    try:
        path = _agent_safe_path(project_dir, rel_path)
    except ValueError as e:
        return f"ERROR: {e}"
    if not os.path.isfile(path):
        return f"ERROR: '{rel_path}' does not exist"
    try:
        with open(path, "r", encoding="utf-8", errors="replace") as f:
            text = f.read()
        if len(text) > max_chars:
            text = text[:max_chars] + f"\n...[truncated, {len(text)} chars total]"
        return text
    except Exception as e:
        return f"ERROR reading file: {e}"


def agent_write_file(project_dir, rel_path, content):
    try:
        path = _agent_safe_path(project_dir, rel_path)
    except ValueError as e:
        return False, str(e)
    try:
        os.makedirs(os.path.dirname(path), exist_ok=True)
        with open(path, "w", encoding="utf-8") as f:
            f.write(content)
        return True, None
    except Exception as e:
        return False, str(e)


def agent_run_command(project_dir, command, timeout=60):
    try:
        result = subprocess.run(
            command, shell=True, cwd=project_dir, capture_output=True,
            text=True, timeout=timeout)
        return result.stdout, result.stderr, result.returncode
    except subprocess.TimeoutExpired:
        return "", f"Command timed out after {timeout}s", -1
    except Exception as e:
        return "", str(e), -1


CODING_AGENT_SYSTEM = (
    "You are an autonomous coding agent working in a real project folder on the user's "
    "computer. You can list directories, read files, write/create files, and run shell "
    "commands (tests, linters, git, package managers, compilers/interpreters for any "
    "language — Python, MATLAB, etc. — and so on) to accomplish the user's task.\n\n"
    "Respond with ONLY a single JSON object describing your NEXT action — nothing else, no "
    "markdown fences, no explanation outside the JSON:\n"
    '{"action": "list_dir", "path": "."}\n'
    '{"action": "read_file", "path": "relative/path.py"}\n'
    '{"action": "write_file", "path": "relative/path.py", "content": "...full file content..."}\n'
    '{"action": "run_command", "command": "pytest -q"}\n'
    '{"action": "done", "summary": "what you accomplished"}\n\n'
    "Rules:\n"
    "- Paths are relative to the project folder — never an absolute path or one starting "
    "with '..'.\n"
    "- write_file always sends the COMPLETE file content, never a diff, snippet, or "
    "'rest unchanged' placeholder.\n"
    "- Explore before you edit: list_dir / read_file the relevant files first so changes are "
    "grounded in what's actually there, not guessed.\n"
    "- Exactly ONE action per response. After each action you'll be shown its result (or told "
    "the user rejected it) and asked for the next one.\n"
    "- Call 'done' as soon as the task is genuinely complete — don't keep exploring or making "
    "unnecessary changes once it is.\n\n"
    "ENVIRONMENT / COMMAND-FAILURE RULES — read carefully, this avoids a common dead end:\n"
    "- A command failing does NOT mean a tool or library is missing — it often means this "
    "particular command resolved to the wrong interpreter/environment on the user's PATH. "
    "Never conclude something is unavailable, and never write a workaround to avoid using it, "
    "after only ONE failed attempt.\n"
    "- On Windows specifically: 'python' and 'pip' on PATH can silently point at an unrelated "
    "or broken virtual environment (e.g. one missing pip entirely), even while a perfectly "
    "working Python installation exists elsewhere on the machine. If 'python ...' or "
    "'python -m pip ...' fails with something like 'No module named pip', immediately retry "
    "the SAME goal with the Windows Python Launcher instead: 'py -m pip install X', "
    "'py script.py', 'py -3 ...'. Try that before concluding anything is missing or "
    "unreachable.\n"
    "- More generally, if one interpreter/tool invocation fails in a way that looks like an "
    "environment/PATH problem (not a real error in your code or a genuinely missing package), "
    "try the other common spelling for it (python/python3/py, pip/pip3/py -m pip, etc.) before "
    "giving up on that approach."
)


def _agent_environment_note(python_path=None):
    """A short, accurate description of the host OS (and, if available, a specific Python
    interpreter and/or MATLAB known to work), generated fresh each run rather than assumed, so
    the agent's environment guidance is grounded in the real machine it's on."""
    note = f"[Host environment: {platform.system()} {platform.release()} — {platform.machine()}]"
    if python_path:
        note += (f"\n[A working Python interpreter with this app's dependencies already "
                 f"installed is at: {python_path} — prefer it over a bare 'python'/'pip' "
                 f"command if you need Python or pip for this task, e.g. "
                 f'"{python_path}" -m pip install X or "{python_path}" script.py]')
    matlab_path = shutil.which("matlab")
    if matlab_path:
        note += (
            f"\n[MATLAB is installed and on PATH at: {matlab_path}. To run a .m script or "
            f"function non-interactively (no GUI, exits automatically, output goes to "
            f'stdout), use: matlab -batch "script_name" (no .m extension, run from the '
            f"folder containing it, or use addpath first) — this is the correct way to run "
            f"MATLAB code from a shell command, not just 'matlab script.m'.]"
        )
    return note


def agent_next_action(history, model, api_base=None, api_key=None, num_ctx=None, max_tokens=2000):
    """Ask the LLM for its next action given the agent's conversation history so far.
    Returns (action_dict_or_None, raw_reply_text)."""
    msgs = [{"role": "system", "content": CODING_AGENT_SYSTEM}] + history
    try:
        chunks = list(stream_chat(msgs, model, api_base=api_base, api_key=api_key,
                                  temperature=0.2, num_ctx=num_ctx, max_tokens=max_tokens))
        raw = "".join(chunks).strip()
    except Exception as e:
        return None, f"(LLM call failed: {e})"
    cleaned = re.sub(r"^```(json)?\s*|\s*```$", "", raw, flags=re.I).strip()
    try:
        return json.loads(cleaned), raw
    except Exception:
        m = re.search(r"\{.*\}", cleaned, re.S)
        if m:
            try:
                return json.loads(m.group(0)), raw
            except Exception:
                pass
    return None, raw


# =============================================================================
# EDITABLE CREW TASKS  (defaults; user can modify/add/delete in the Tasks tab)
# Placeholders: {topic} is replaced with the research query. The selected papers'
# full text is auto-injected into the FIRST task. Tasks run in order, each one
# receiving the previous task's output as context.
# =============================================================================
AGENT_ROLES = ["Paper Analyst", "Scientific Writer & Reviewer",
               "Research & Innovation Strategist"]

DEFAULT_TASKS = [
    {"agent": "Paper Analyst",
     "description": ("For EACH paper provided below, write a detailed summary covering: "
                     "(1) methodology, (2) dataset or participants, (3) key results with any "
                     "numbers stated, (4) limitations. Use ONLY the provided text; never invent "
                     "papers, authors, or numbers."),
     "expected_output": "One detailed, grounded summary per paper."},
    {"agent": "Scientific Writer & Reviewer",
     "description": ("Using ONLY the summaries in your context, write a structured literature "
                     "review strictly about '{topic}', with concrete methods and numbers. End "
                     "with at least one clearly stated research gap. Do not introduce any topic "
                     "outside '{topic}'."),
     "expected_output": "A structured, on-topic literature review ending with a research gap."},
    {"agent": "Research & Innovation Strategist",
     "description": ("Based on the review (strictly about '{topic}'), propose 2 original research "
                     "ideas that address the identified gap. For EACH idea give: a) title, "
                     "b) abstract, c) methodology, d) intended impact."),
     "expected_output": "Two on-topic paper proposals, each with title, abstract, methodology, impact."},
]


# =============================================================================
# Chat streaming (Ollama native, or pure requests for OpenAI/Google)
# bypassing litellm entirely for compatibility
# =============================================================================
def _inject_user_memory(messages):
    """Appends the signed-in user's accumulated memory context (set as the global
    USER_MEMORY_CONTEXT right after login) onto the outgoing system message, or adds one
    if there isn't one. stream_chat is the one function nearly every tool in this app
    calls to reach an LLM, so hooking it here personalizes every tool — Survey, Novelty,
    Write Paper, Summarize, Review, Coding, Presentation, Chat — without each tool
    needing its own wiring. Uses globals().get() so it's a no-op (not a crash) on any
    call made before login finishes setting USER_MEMORY_CONTEXT."""
    ctx = globals().get("USER_MEMORY_CONTEXT", "")
    if not ctx or not messages:
        return messages
    messages = [dict(m) for m in messages]  # don't mutate the caller's list/dicts
    if messages[0].get("role") == "system":
        messages[0]["content"] = str(messages[0].get("content", "")).rstrip() + "\n\n" + ctx
    else:
        messages.insert(0, {"role": "system", "content": ctx})
    return messages


def stream_chat(messages, model, api_base=None, api_key=None, temperature=0.7, num_ctx=None, max_tokens=None):
    messages = _inject_user_memory(messages)
    # 1. Ollama Native API
    if model.startswith("ollama_chat/") or model.startswith("ollama/"):
        name = model.split("/", 1)[1]
        url = (api_base or "http://localhost:11434").rstrip("/") + "/api/chat"
        options = {"temperature": temperature}
        if num_ctx:
            options["num_ctx"] = int(num_ctx)
        if max_tokens:
            options["num_predict"] = int(max_tokens)
        payload = {"model": name, "messages": messages, "stream": True, "options": options}
        headers = {}
        if api_key:
            headers["Authorization"] = f"Bearer {api_key}"
        with requests.post(url, json=payload, stream=True, timeout=600, verify=False, headers=headers) as r:
            r.raise_for_status()
            for line in r.iter_lines():
                if not line:
                    continue
                data = json.loads(line.decode("utf-8"))
                chunk = data.get("message", {}).get("content", "")
                if chunk:
                    yield chunk
                if data.get("done"):
                    break
        return

    # 2. Anthropic (Claude) Messages API — a genuinely different wire format from OpenAI's:
    # the system prompt is a top-level `system` field (not a role:"system" message — the
    # Claude API rejects that), auth uses an `x-api-key` header plus a required
    # `anthropic-version` header (not `Authorization: Bearer`), and the streaming SSE events
    # are typed (`content_block_delta` etc.) rather than OpenAI's `choices[0].delta.content`.
    if model.startswith("anthropic/") or model.startswith("claude/"):
        actual_model = model.split("/", 1)[1]
        url = (api_base.rstrip("/") if api_base else "https://api.anthropic.com/v1") + "/messages"
        headers = {"Content-Type": "application/json", "anthropic-version": "2023-06-01"}
        if api_key:
            headers["x-api-key"] = api_key

        system_prompt = ""
        claude_messages = []
        for m in messages:
            if m.get("role") == "system":
                system_prompt = (system_prompt + "\n\n" + m["content"]).strip() if system_prompt else m["content"]
            else:
                claude_messages.append({"role": m["role"], "content": m["content"]})

        payload = {
            "model": actual_model,
            "messages": claude_messages,
            "stream": True,
            "temperature": temperature,
            "max_tokens": int(max_tokens) if max_tokens else 4096,  # Claude requires max_tokens
        }
        if system_prompt:
            payload["system"] = system_prompt

        try:
            with requests.post(url, json=payload, headers=headers, stream=True, timeout=120) as r:
                if r.status_code != 200:
                    yield f"⚠️ API Error ({r.status_code}): {r.text}"
                    return
                for line in r.iter_lines():
                    if not line:
                        continue
                    decoded = line.decode("utf-8").strip()
                    if not decoded.startswith("data:"):
                        continue
                    data_str = decoded[5:].strip()
                    if not data_str:
                        continue
                    try:
                        data_json = json.loads(data_str)
                    except json.JSONDecodeError:
                        continue
                    if data_json.get("type") == "content_block_delta":
                        delta = data_json.get("delta", {})
                        if delta.get("type") == "text_delta":
                            chunk = delta.get("text", "")
                            if chunk:
                                yield chunk
                    elif data_json.get("type") == "message_stop":
                        break
        except Exception as e:
            yield f"⚠️ Connection Error: {e}"
        return

    # 3. Pure REST requests for Google AI Studio / any OpenAI-compatible endpoint (OpenAI
    # itself, Qwen/DashScope, Together, Groq, Fireworks, DeepSeek, a local vLLM server, etc.)
    headers = {"Content-Type": "application/json"}
    if api_key:
        headers["Authorization"] = f"Bearer {api_key}"

    if model.startswith("gemini/"):
        actual_model = model.split("/", 1)[1]
        url = "https://generativelanguage.googleapis.com/v1beta/openai/chat/completions"
    elif model.startswith("openai/"):
        actual_model = model.split("/", 1)[1]
        url = (api_base.rstrip("/") + "/chat/completions") if api_base else "https://api.openai.com/v1/chat/completions"
    else:
        actual_model = model
        url = (api_base.rstrip("/") + "/chat/completions") if api_base else "https://api.openai.com/v1/chat/completions"

    payload = {
        "model": actual_model,
        "messages": messages,
        "stream": True,
        "temperature": temperature
    }
    if max_tokens:
        payload["max_tokens"] = int(max_tokens)

    try:
        with requests.post(url, json=payload, headers=headers, stream=True, timeout=60) as r:
            if r.status_code != 200:
                yield f"⚠️ API Error ({r.status_code}): {r.text}"
                return
            for line in r.iter_lines():
                if line:
                    decoded = line.decode("utf-8").strip()
                    if decoded.startswith("data:"):
                        data_str = decoded[5:].strip()
                        if data_str == "[DONE]":
                            break
                        try:
                            data_json = json.loads(data_str)
                            choices = data_json.get("choices", [])
                            if choices:
                                chunk = choices[0].get("delta", {}).get("content", "")
                                if chunk:
                                    yield chunk
                        except json.JSONDecodeError:
                            pass
    except Exception as e:
        yield f"⚠️ Connection Error: {e}"


# =============================================================================
# TEXT-TO-SPEECH  (edge-tts: free Microsoft neural voices, no API key)
#   pip install edge-tts
# =============================================================================
def clean_for_speech(text):
    """Strip markdown/links/code so the voice reads naturally."""
    t = re.sub(r"```.*?```", " (code block) ", text, flags=re.S)
    t = re.sub(r"`([^`]*)`", r"\1", t)
    t = re.sub(r"\[([^\]]+)\]\([^)]+\)", r"\1", t)   # [text](url) -> text
    t = re.sub(r"https?://\S+", " ", t)
    t = re.sub(r"[#*_>|~]", " ", t)
    t = re.sub(r"\s+", " ", t)
    return t.strip()


def tts_to_mp3(text, voice="en-US-AriaNeural", rate="+0%"):
    """Return MP3 bytes for the given text, or (None, error_message)."""
    try:
        import edge_tts
    except ImportError:
        return None, "Install the voice engine first:  pip install edge-tts"

    import asyncio, tempfile
    spoken = clean_for_speech(text)
    if not spoken:
        return None, "Nothing to speak."

    fd, path = tempfile.mkstemp(suffix=".mp3")
    os.close(fd)

    async def _gen():
        com = edge_tts.Communicate(spoken, voice, rate=rate)
        await com.save(path)

    try:
        try:
            loop = asyncio.new_event_loop()
            loop.run_until_complete(_gen())
            loop.close()
        except RuntimeError:
            asyncio.run(_gen())
        with open(path, "rb") as f:
            data = f.read()
        return data, None
    except Exception as e:
        return None, f"Voice generation failed: {e}"
    finally:
        try:
            os.remove(path)
        except Exception:
            pass


# A few good voices (incl. Arabic, since it's handy in Egypt/Saudi)
TTS_VOICES = {
    "Aria (US, female)": "en-US-AriaNeural",
    "Guy (US, male)": "en-US-GuyNeural",
    "Jenny (US, female)": "en-US-JennyNeural",
    "Sonia (UK, female)": "en-GB-SoniaNeural",
    "Ryan (UK, male)": "en-GB-RyanNeural",
    "Salma (Egypt, Arabic)": "ar-EG-SalmaNeural",
    "Hamed (Egypt, Arabic)": "ar-EG-ShakirNeural",
    "Zariyah (Saudi, Arabic)": "ar-SA-ZariyahNeural",
}


# =============================================================================
# UI
# =============================================================================
st.set_page_config(page_title="Research Crew + Chat", page_icon="🔬", layout="wide")

# =============================================================================
# LOGIN GATE — nothing below this renders until the visitor is signed in.
# =============================================================================
if "auth_user" not in st.session_state:
    st.title("🔬 Research Crew  +  💬 Chat")
    st.caption("Sign in to use the tools. Each account has its own private data on this workstation "
              "— nobody else can see your papers, sessions, or memory.")

    tab_login, tab_signup = st.tabs(["🔑 Log In", "🆕 Create Account"])

    with tab_login:
        with st.form("login_form"):
            login_ident = st.text_input("Username or email")
            login_pw = st.text_input("Password", type="password")
            login_submitted = st.form_submit_button("Log In", type="primary", use_container_width=True)
        if login_submitted:
            auth_username, auth_err = authenticate(login_ident, login_pw)
            if auth_username:
                st.session_state["auth_user"] = auth_username
                record_login(auth_username)
                st.rerun()
            else:
                st.error(auth_err or "Couldn't log in.")

    with tab_signup:
        with st.form("signup_form"):
            signup_username = st.text_input("Username", help="3-32 characters: letters, numbers, underscore, dot, or hyphen.")
            signup_email = st.text_input("Email address")
            signup_pw = st.text_input("Password", type="password", help="At least 8 characters.")
            signup_pw2 = st.text_input("Confirm password", type="password")
            signup_submitted = st.form_submit_button("Create Account", type="primary", use_container_width=True)
        if signup_submitted:
            if signup_pw != signup_pw2:
                st.error("Passwords don't match.")
            else:
                signup_ok, signup_err = create_account(signup_username, signup_email, signup_pw)
                if signup_ok:
                    st.success("Account created! Switch to the 'Log In' tab to sign in.")
                else:
                    st.error(signup_err)

    st.stop()

# ---- Signed in: give every user their own private data folder on disk, ------
# isolated from every other account. Reassigning these here (rather than at
# their original definitions above) means every function that reads/writes
# them — save_coding_session, load_cc_session, load_memory, etc. — picks up
# the correct per-user path, since Python looks up a global by its current
# value at call time, not at the time the function was defined.
_CURRENT_USER = st.session_state["auth_user"]
_CURRENT_IS_ADMIN = is_admin_user(_CURRENT_USER)
_USER_DIR = os.path.join(BASE_DIR, "user_data", _safe_username_for_path(_CURRENT_USER))
os.makedirs(_USER_DIR, exist_ok=True)
MEMORY_DIR = os.path.join(_USER_DIR, "chat_memory")
WP_PROJECTS_DIR = os.path.join(_USER_DIR, "write_paper_projects")
WP_PDFS_DIR = os.path.join(_USER_DIR, "write_paper_project_pdfs")
SURVEY_PROJECTS_DIR = os.path.join(_USER_DIR, "survey_projects")
SURVEY_PDFS_DIR = os.path.join(_USER_DIR, "survey_project_pdfs")
CREW_PROJECTS_DIR = os.path.join(_USER_DIR, "research_crew_projects")
CREW_PDFS_DIR = os.path.join(_USER_DIR, "research_crew_project_pdfs")
CODING_SESSIONS_DIR = os.path.join(_USER_DIR, "coding_agent_sessions")
CC_SESSIONS_DIR = os.path.join(_USER_DIR, "claude_code_sessions")
CHAT_CONVOS_DIR = os.path.join(_USER_DIR, "chat_conversations")
os.makedirs(CHAT_CONVOS_DIR, exist_ok=True)

# The Agent and Claude Code tabs let this user point at a real folder on this
# machine for the model to read/write/run commands in — on a shared workstation,
# that's confined to this sandbox so one account can never point at another's
# (or the wider filesystem's) files. See path_in_sandbox() and its call sites
# in the Agent/Claude Code tabs below.
USER_SANDBOX_DIR = os.path.join(_USER_DIR, "projects")
os.makedirs(USER_SANDBOX_DIR, exist_ok=True)

# Global per-user memory: builds up automatically as this account uses the tools
# (topics researched, papers written, languages coded in, etc.) and is injected into
# every tool's LLM calls via _inject_user_memory() inside stream_chat, plus the
# Research Crew task description directly (the one tool that bypasses stream_chat).
# Keyed by username itself now that accounts exist — see the memory-fact logging
# calls alongside each tool's usage-event logging below.
USER_MEMORY_CONTEXT = get_user_memory_context(_CURRENT_USER)

st.title("🔬 Research Crew  +  💬 Chat")

# --- chat_input positioning fix ----------------------------------------------
# Streamlit auto-pins st.chat_input to the bottom of the browser viewport, but
# ONLY when it's called at the root of the script. Once it's nested inside
# st.tabs() (as it is in the Chat tab), that fixed/sticky positioning breaks and
# the box can render stuck near the top of the tab instead of after the last
# message. Overriding it to `position: static` makes it render inline instead —
# i.e. exactly where the code calls it, which is right after the message list.
st.markdown("""
<style>
div[data-testid="stBottomBlockContainer"],
div[data-testid="stBottom"],
div[data-testid="stChatInput"] {
    position: static !important;
    bottom: auto !important;
}
</style>
""", unsafe_allow_html=True)

with st.sidebar:
    su1, su2 = st.columns([2.2, 1])
    with su1:
        st.caption(f"👤 **{_CURRENT_USER}**" + (" 🛡️ admin" if _CURRENT_IS_ADMIN else ""))
    with su2:
        if st.button("Log out", use_container_width=True, key="logout_btn"):
            st.session_state.pop("auth_user", None)
            st.rerun()
    st.divider()

    st.header("⚙️ Configuration")
    
    provider_choice = st.selectbox("AI Provider", [
        "Google AI Studio", "Ollama (Local/Cloud)", "Anthropic (Claude)",
        "DeepSeek (OpenAI)", "DeepSeek (Anthropic)",
        "Qwen (DashScope)", "Qwen (Anthropic)", "Custom (Anthropic-compatible)",
        "Custom (OpenAI-compatible)",
    ], help="💡 Don't have an API key? Google AI Studio is free — get a key at aistudio.google.com")

    st.caption("💡 **No API key?** Select **Google AI Studio** — it's free. "
               "Get a key at [aistudio.google.com](https://aistudio.google.com) → Get API Key → copy it here.")
    
    selected_api_key = None
    selected_api_base = None
    
    if provider_choice == "Google AI Studio":
        st.subheader("Google AI Studio connection")
        google_key = st.text_input("Google AI API Key", type="password",
                                   value=os.environ.get("GEMINI_API_KEY", ""),
                                   help="Paste your API key from Google AI Studio.")
        selected_api_key = google_key or None

        # Dynamically fetch available models directly from Google based on API Key
        google_models = list_google_models(selected_api_key)

        st.divider()
        st.subheader("Research-crew model")
        crew_model = st.selectbox("Gemini Model", google_models, index=default_model_index(google_models, provider="google"))
        global_model_string = f"gemini/{crew_model}"
        crew_model_string = f"gemini/{crew_model}"

    elif provider_choice == "Ollama (Local/Cloud)":
        st.subheader("Ollama connection")
        base_url = st.text_input("Ollama base URL", value="http://localhost:11434",
                                 help="Local: http://localhost:11434  ·  "
                                      "Ollama Cloud (direct): https://ollama.com")
        ollama_key = st.text_input("Ollama API key (only for Ollama Cloud)", type="password",
                                   value=os.environ.get("OLLAMA_API_KEY", ""),
                                   help="Leave blank for local use. For cloud, paste the key from "
                                        "ollama.com → Settings → Keys.")
        selected_api_base = base_url
        selected_api_key = ollama_key or None

        installed = list_ollama_models(base_url, api_key=ollama_key or None)
        st.caption(f"✅ {len(installed)} model(s) found" if installed else "⚠️ No Ollama models found")

        st.divider()
        st.subheader("Research-crew model")
        if installed:
            crew_model = st.selectbox("Crew model", installed, index=default_model_index(installed, provider="ollama"))
            global_model_string = f"ollama_chat/{crew_model}"
            crew_model_string = f"ollama/{crew_model}"
        else:
            crew_model = st.text_input("Crew model", value="llama3.2")
            global_model_string = f"ollama_chat/{crew_model}"
            crew_model_string = f"ollama/{crew_model}"

    elif provider_choice == "Anthropic (Claude)":
        st.subheader("Anthropic connection")
        anthropic_key = st.text_input("Anthropic API Key", type="password",
                                      value=os.environ.get("ANTHROPIC_API_KEY", ""),
                                      help="From console.anthropic.com → Settings → API Keys.")
        selected_api_key = anthropic_key or None
        selected_api_base = None

        st.divider()
        st.subheader("Research-crew model")
        crew_model = st.text_input(
            "Claude model name", value="claude-sonnet-4-5",
            help="e.g. claude-opus-4-1, claude-sonnet-4-5, claude-haiku-4-5 — check "
                 "docs.anthropic.com/en/docs/about-claude/models for the current list, "
                 "since it changes over time.")
        global_model_string = f"anthropic/{crew_model}"
        crew_model_string = f"anthropic/{crew_model}"

    elif provider_choice == "DeepSeek (OpenAI)":
        st.subheader("DeepSeek (OpenAI-compatible) connection")
        ds_openai_base = st.text_input(
            "API base URL",
            value="https://api.deepseek.com/v1",
            help="DeepSeek's OpenAI-compatible endpoint.")
        ds_openai_key = st.text_input("API Key", type="password",
                                      value=os.environ.get("DEEPSEEK_API_KEY", ""),
                                      help="From platform.deepseek.com → API keys.")
        selected_api_key = ds_openai_key or None
        selected_api_base = ds_openai_base

        ds_openai_models = list_openai_compatible_models(ds_openai_base, ds_openai_key or None) if ds_openai_base else []
        if ds_openai_base:
            st.caption(f"✅ {len(ds_openai_models)} model(s) found" if ds_openai_models
                      else "⚠️ No models found — enter one manually below")

        st.divider()
        st.subheader("Research-crew model")
        if ds_openai_models:
            crew_model = st.selectbox("DeepSeek model", ds_openai_models,
                                      index=ds_openai_models.index("deepseek-chat") if "deepseek-chat" in ds_openai_models else 0)
        else:
            crew_model = st.text_input("DeepSeek model name", value="deepseek-chat",
                                       help="e.g. deepseek-chat, deepseek-reasoner.")
        global_model_string = f"openai/{crew_model}"
        crew_model_string = f"openai/{crew_model}"

    elif provider_choice == "DeepSeek (Anthropic)":
        st.subheader("DeepSeek (Anthropic-compatible) connection")
        st.caption("DeepSeek via an Anthropic Messages-format endpoint — useful with proxies "
                  "or gateways that translate Anthropic API calls to DeepSeek.")
        ds_anthropic_base = st.text_input(
            "API base URL", value="",
            placeholder="e.g. https://your-proxy.example.com",
            help="An Anthropic-compatible endpoint that proxies to DeepSeek models.")
        ds_anthropic_key = st.text_input("API Key", type="password",
                                         value=os.environ.get("DEEPSEEK_API_KEY", ""),
                                         help="DeepSeek API key (or proxy key).")
        selected_api_key = ds_anthropic_key or None
        selected_api_base = ds_anthropic_base or None

        ds_anthropic_models = list_anthropic_compatible_models(ds_anthropic_base, ds_anthropic_key or None) if ds_anthropic_base else []
        if ds_anthropic_base:
            st.caption(f"✅ {len(ds_anthropic_models)} model(s) found" if ds_anthropic_models
                      else "⚠️ No models found — enter one manually below")

        st.divider()
        st.subheader("Research-crew model")
        if ds_anthropic_models:
            crew_model = st.selectbox("Model", ds_anthropic_models,
                                      index=ds_anthropic_models.index("deepseek-chat") if "deepseek-chat" in ds_anthropic_models else 0)
        else:
            crew_model = st.text_input("Model name", value="deepseek-chat",
                                       placeholder="e.g. deepseek-chat, deepseek-reasoner")
        global_model_string = f"anthropic/{crew_model}"
        crew_model_string = f"anthropic/{crew_model}"

    elif provider_choice == "Qwen (DashScope)":
        st.subheader("Qwen connection")
        qwen_base = st.text_input(
            "API base URL",
            value="https://token-plan.ap-southeast-1.maas.aliyuncs.com/compatible-mode/v1",
            help="Qwen's OpenAI-compatible endpoint.")
        qwen_key = st.text_input("API Key", type="password",
                                 value=os.environ.get("DASHSCOPE_API_KEY", ""),
                                 help="From your Alibaba Cloud / Model Studio console.")
        selected_api_key = qwen_key or None
        selected_api_base = qwen_base

        qwen_models = list_openai_compatible_models(qwen_base, qwen_key or None)
        st.caption(f"✅ {len(qwen_models)} model(s) found" if qwen_models else "⚠️ No models found — enter one manually below")

        st.divider()
        st.subheader("Research-crew model")
        if qwen_models:
            crew_model = st.selectbox("Qwen model", qwen_models,
                                      index=qwen_models.index("qwen-plus") if "qwen-plus" in qwen_models else 0)
        else:
            crew_model = st.text_input("Qwen model name", value="qwen-plus",
                                       help="e.g. qwen-plus, qwen-max, qwen-turbo.")
        global_model_string = f"openai/{crew_model}"
        crew_model_string = f"openai/{crew_model}"

    elif provider_choice == "Qwen (Anthropic)":
        st.subheader("Qwen (Anthropic-compatible) connection")
        st.caption("Alibaba's Anthropic Messages-format endpoint — speaks the same API shape "
                  "Claude does, so this also works natively with the Claude Code tab (no "
                  "translating proxy needed there).")
        qwen_anthropic_base = st.text_input(
            "API base URL", value="https://token-plan.ap-southeast-1.maas.aliyuncs.com/apps/anthropic",
            help="Qwen's Anthropic-compatible endpoint.")
        qwen_anthropic_key = st.text_input("API Key", type="password",
                                           value=os.environ.get("DASHSCOPE_API_KEY", ""),
                                           help="From your Alibaba Cloud / Model Studio console.")
        selected_api_key = qwen_anthropic_key or None
        selected_api_base = qwen_anthropic_base

        # The Anthropic-format endpoint's own /v1/models doesn't return a model list here —
        # DashScope's OpenAI-compatible endpoint (same Alibaba account/key) does, and the two
        # expose the same underlying Qwen models, so borrow its listing for the picker while
        # still sending actual requests to the Anthropic-format base URL above.
        qwen_anthropic_models = list_openai_compatible_models(
            "https://token-plan.ap-southeast-1.maas.aliyuncs.com/compatible-mode/v1",
            qwen_anthropic_key or None)
        st.caption(f"✅ {len(qwen_anthropic_models)} model(s) found" if qwen_anthropic_models
                  else "⚠️ No models found — enter one manually below")

        st.divider()
        st.subheader("Research-crew model")
        if qwen_anthropic_models:
            crew_model = st.selectbox(
                "Qwen model", qwen_anthropic_models,
                index=qwen_anthropic_models.index("qwen3-max") if "qwen3-max" in qwen_anthropic_models else 0)
        else:
            crew_model = st.text_input(
                "Qwen model name", value="qwen3-max",
                help="e.g. qwen3-max, qwen3-coder-plus — check Alibaba Model Studio's docs for the "
                     "current list of models exposed on this endpoint.")
        global_model_string = f"anthropic/{crew_model}"
        crew_model_string = f"anthropic/{crew_model}"

    elif provider_choice == "Custom (Anthropic-compatible)":
        st.subheader("Custom Anthropic-compatible connection")
        st.caption("Works with any provider that speaks the Anthropic Messages API format — "
                  "LiteLLM proxy, OpenRouter, a local Claude Code proxy, etc.")
        custom_anthropic_base = st.text_input(
            "API base URL", value="",
            placeholder="e.g. https://your-proxy.example.com",
            help="The base URL of the Anthropic-compatible endpoint.")
        custom_anthropic_key = st.text_input("API Key", type="password",
                                            value=os.environ.get("CUSTOM_ANTHROPIC_API_KEY", ""))
        selected_api_key = custom_anthropic_key or None
        selected_api_base = custom_anthropic_base or None

        custom_anthropic_models = list_anthropic_compatible_models(custom_anthropic_base, custom_anthropic_key or None) if custom_anthropic_base else []
        if custom_anthropic_base:
            st.caption(f"✅ {len(custom_anthropic_models)} model(s) found" if custom_anthropic_models
                      else "⚠️ No models found — enter one manually below")

        st.divider()
        st.subheader("Research-crew model")
        if custom_anthropic_models:
            crew_model = st.selectbox("Model", custom_anthropic_models,
                                      index=default_model_index(custom_anthropic_models, provider="ollama"))
        else:
            crew_model = st.text_input("Model name", value="",
                                       placeholder="e.g. claude-sonnet-4-5")
        global_model_string = f"anthropic/{crew_model}"
        crew_model_string = f"anthropic/{crew_model}"

    else:  # "Custom (OpenAI-compatible)"
        st.subheader("Custom OpenAI-compatible connection")
        st.caption("Works with any provider that speaks the OpenAI chat-completions format — "
                  "Together AI, Groq, Fireworks, DeepSeek, OpenRouter, a local vLLM/llama.cpp "
                  "server, etc.")
        custom_base = st.text_input(
            "API base URL", value="",
            placeholder="e.g. https://api.groq.com/openai/v1",
            help="The base URL up to (not including) /chat/completions.")
        custom_key = st.text_input("API Key", type="password",
                                   value=os.environ.get("CUSTOM_API_KEY", ""))
        selected_api_key = custom_key or None
        selected_api_base = custom_base or None

        custom_models = list_openai_compatible_models(custom_base, custom_key or None) if custom_base else []
        if custom_base:
            st.caption(f"✅ {len(custom_models)} model(s) found" if custom_models else "⚠️ No models found — enter one manually below")

        st.divider()
        st.subheader("Research-crew model")
        if custom_models:
            crew_model = st.selectbox("Model", custom_models,
                                      index=default_model_index(custom_models, provider="ollama"))
        else:
            crew_model = st.text_input("Model name", value="",
                                       placeholder="e.g. meta-llama/Llama-3.3-70B-Instruct-Turbo")
        global_model_string = f"openai/{crew_model}"
        crew_model_string = f"openai/{crew_model}"

    crew_temperature = st.slider("Crew temperature", 0.0, 1.5, 0.4, 0.05)
    num_ctx = st.select_slider("Context window (num_ctx)",
                               options=[2048, 4096, 8192, 16384, 32768, 65536, 131072],
                               value=32768,
                               help="Local models: keep modest (8k–16k). Cloud models "
                                    "support large windows — 65k–131k lets you feed whole papers.")
    crew_max_tokens = st.slider("Max output tokens/step", 512, 8192, 2048, 256)

    st.divider()
    with st.expander("🛑 Stop the app"):
        st.caption("This shuts down the whole app (kills the Python process). You'll need to "
                   "re-run `streamlit run app.py` to start it again.")
        confirm_stop = st.checkbox("Yes, I want to stop the app", key="confirm_stop")
        if st.button("🛑 Stop & Exit now", disabled=not confirm_stop):
            st.warning("Shutting down… you can close this browser tab.")
            import time as _t
            _t.sleep(1)
            os._exit(0)

_TAB_LABELS = ["🔬 Research Crew", "📚 Make a Survey", "🧭 Check Novelty", "✍️ Write Paper", "📖 Summarize Paper",
              "📝 Review Paper", "💻 Coding", "🖼️ Make Presentation", "💬 Chat with any model"]
if _CURRENT_IS_ADMIN:
    _TAB_LABELS.append("🛡️ Admin")
_TABS = st.tabs(_TAB_LABELS)
(tab_research, tab_survey, tab_novelty, tab_writepaper, tab_summarize, tab_review, tab_coding,
 tab_present, tab_chat) = _TABS[:9]
if _CURRENT_IS_ADMIN:
    tab_admin = _TABS[9]

# =========================== RESEARCH CREW TAB (Sources + Tasks + Crew) =====
with tab_research:
    st.subheader("🔬 Autonomous Academic Writing & Innovation Crew")
    st.caption("Stage 1: configure sources & tasks. Stage 2: find papers. Stage 3: choose items. "
               "Stage 4: execute crew instructions.")

    CREW_STATE_KEYS = ["research_query", "research_query_used", "candidates", "research_result"]

    st.markdown("#### 📁 Projects")
    st.caption("Save your progress — topic, gathered papers, and the crew's output — and pick it "
              "back up later. The PDFs of papers you analyze are kept in their own folder "
              "alongside the project too.")
    crew_projects = list_crew_projects()
    crpc1, crpc2, crpc3 = st.columns([3, 1, 1])
    with crpc1:
        crew_selected_project = st.selectbox(
            "Saved projects", ["(none)"] + crew_projects, key="crew_project_picker",
            label_visibility="collapsed")
    with crpc2:
        if st.button("📂 Load", disabled=(crew_selected_project == "(none)"), use_container_width=True, key="crew_proj_load"):
            payload, err = load_crew_project(crew_selected_project)
            if err:
                st.error(f"Could not load: {err}")
            else:
                for k in CREW_STATE_KEYS:
                    if k in payload:
                        st.session_state[k] = payload[k]
                st.session_state["crew_active_project"] = crew_selected_project
                st.session_state["_pending_crew_project_name"] = crew_selected_project
                st.rerun()
    with crpc3:
        if st.button("🗑️ Delete", disabled=(crew_selected_project == "(none)"), use_container_width=True, key="crew_proj_delete"):
            delete_crew_project(crew_selected_project)
            if st.session_state.get("crew_active_project") == crew_selected_project:
                st.session_state.pop("crew_active_project", None)
            st.rerun()

    if "_pending_crew_project_name" in st.session_state:
        st.session_state["crew_project_name"] = st.session_state.pop("_pending_crew_project_name")

    crsc1, crsc2 = st.columns([3, 1])
    with crsc1:
        crew_save_name = st.text_input(
            "Project name", value=st.session_state.get("crew_active_project", ""),
            key="crew_project_name", label_visibility="collapsed",
            placeholder="Name this project to save it…")
    with crsc2:
        if st.button("💾 Save Project", disabled=not crew_save_name.strip(), use_container_width=True, key="crew_proj_save"):
            state = {k: st.session_state.get(k) for k in CREW_STATE_KEYS}
            ok, err = save_crew_project(crew_save_name.strip(), state)
            if ok:
                st.session_state["crew_active_project"] = crew_save_name.strip()
                st.success(f"Saved '{crew_save_name.strip()}'.")
            else:
                st.error(f"Could not save: {err}")

    if st.session_state.get("crew_active_project"):
        st.caption(f"📌 Working on: **{st.session_state['crew_active_project']}** "
                  f"— remember to 💾 Save Project after changes.")
        _crew_pdfs = list_project_pdfs(CREW_PDFS_DIR, st.session_state["crew_active_project"])
        if _crew_pdfs:
            with st.expander(f"📄 {len(_crew_pdfs)} saved PDF(s) for this project"):
                for _fn in _crew_pdfs:
                    _fpath = os.path.join(CREW_PDFS_DIR, re.sub(r"[^A-Za-z0-9_-]+", "_", st.session_state["crew_active_project"]), _fn)
                    try:
                        with open(_fpath, "rb") as _fh:
                            st.download_button(f"⬇️ {_fn}", data=_fh.read(), file_name=_fn, key=f"crewpdf_{_fn}")
                    except Exception:
                        st.caption(f"⚠️ {_fn} (couldn't read)")

    if st.button("🆕 New Project (clears current work)", key="crew_reset"):
        for key in list(st.session_state.keys()):
            if (key.startswith("crew_") or key.startswith("pick_") or key == "research_result"
                    or key in ("research_query", "research_query_used", "candidates")):
                del st.session_state[key]
        st.rerun()

    st.divider()

    with st.expander("📁 Sources & Retrieval", expanded=False):
        use_arxiv = st.checkbox("arXiv (physics, CS, math)", value=True)
        use_s2 = st.checkbox("Semantic Scholar (all fields, incl. medical/bio)", value=False,
                             help="Free tier is heavily rate-limited (frequent 429 errors). "
                                  "Enable only if you need it.")
        use_openalex = st.checkbox("OpenAlex (IEEE, Elsevier, Springer, ACM, …)", value=True,
                                   help="Indexes the major publishers. Returns title + abstract for "
                                        "all, and an open-access PDF link when one is available.")
        per_source = st.slider("Results per source", 2, 10, 5)
        sort_choice = st.radio("Sort arXiv by", ["relevance", "newest"], horizontal=True)
        chars_per_paper = st.select_slider("Chars of full text / paper",
                                           options=[3000, 6000, 9000, 12000, 20000, 40000, 80000],
                                           value=6000)

    with st.expander("🧩 Crew Tasks", expanded=False):
        st.caption("These are the steps the research crew runs, in order. Edit the instructions, "
                   "change which agent does each step, or add/remove steps. Use {topic} to insert "
                   "your research query. The selected papers are automatically given to the first task.")

        if "tasks" not in st.session_state:
            st.session_state.tasks = [dict(t) for t in DEFAULT_TASKS]

        top = st.columns([1, 1, 5])
        with top[0]:
            if st.button("➕ Add task"):
                st.session_state.tasks.append(
                    {"agent": AGENT_ROLES[0], "description": "", "expected_output": ""})
                st.rerun()
        with top[1]:
            if st.button("↩️ Reset to defaults"):
                st.session_state.tasks = [dict(t) for t in DEFAULT_TASKS]
                st.rerun()

        if not st.session_state.tasks:
            st.warning("No tasks defined. Add at least one task or reset to defaults.")

        st.caption("Steps (expand each to edit — not nested inside this panel, per Streamlit's "
                  "expander rules, so they list directly below):")

    # Per-task expanders must NOT be nested inside the "Crew Tasks" expander above — Streamlit
    # disallows an expander inside another expander — so they render as their own top-level list.
    delete_idx = None
    for i, t in enumerate(st.session_state.get("tasks", [])):
        chain = " · receives previous task's output" if i > 0 else " · receives the selected papers"
        with st.expander(f"Step {i+1}: {t['agent']}{chain}", expanded=(i == 0)):
            t["agent"] = st.selectbox(
                "Agent for this step", AGENT_ROLES,
                index=AGENT_ROLES.index(t["agent"]) if t["agent"] in AGENT_ROLES else 0,
                key=f"task_agent_{i}")
            t["description"] = st.text_area(
                "Instructions (use {topic} for the research query)",
                value=t["description"], height=130, key=f"task_desc_{i}")
            t["expected_output"] = st.text_input(
                "Expected output (a short description of the deliverable)",
                value=t["expected_output"], key=f"task_out_{i}")
            cols = st.columns([1, 1, 4])
            with cols[0]:
                if i > 0 and st.button("⬆️ Move up", key=f"up_{i}"):
                    st.session_state.tasks[i-1], st.session_state.tasks[i] = \
                        st.session_state.tasks[i], st.session_state.tasks[i-1]
                    st.rerun()
            with cols[1]:
                if st.button("🗑️ Delete", key=f"del_{i}"):
                    delete_idx = i
    if delete_idx is not None:
        st.session_state.tasks.pop(delete_idx)
        st.rerun()

    if st.session_state.get("tasks"):
        st.info(f"**{len(st.session_state.tasks)}** task(s) defined. "
                "They'll run in this order when you analyze selected papers below.")

    st.divider()

    query = st.text_input("Research topic:", "wearable for autism", key="research_query")
    wcol1, wcol2, wcol3 = st.columns([1.3, 1, 1.3])
    with wcol1:
        crew_web = st.checkbox("🌐 Also blend in web results (DuckDuckGo)", value=False, key="crew_web")
    with wcol2:
        crew_web_read = st.checkbox("Read full pages", value=True, key="crew_web_read", disabled=not crew_web)
    with wcol3:
        rc_max_rounds = st.slider(
            "Max search rounds", 1, 5, 3, key="rc_max_rounds",
            help="The model turns your topic into search terms, checks whether what it finds is "
                 "relevant/sufficient, and searches again with different terms if not — up to "
                 "this many rounds.")

    with st.expander("📁 Add your own papers (uploads or a local folder)", expanded=False):
        up_files = st.file_uploader("Upload one or more PDFs", type=["pdf"], accept_multiple_files=True, key="crew_uploads")
        folder_path = st.text_input("…or a local folder path", value="", key="crew_folder")

    if st.button("🔍 Find papers"):
        any_search = use_arxiv or use_s2 or use_openalex
        cands, errors = [], []

        def _rc_search_once(q):
            papers, errs = [], []
            if use_arxiv:
                try: papers += search_arxiv_meta(q, n=per_source, sort=sort_choice)
                except Exception as e: errs.append(f"arXiv: {e}")
            if use_s2:
                try: papers += search_semantic_scholar(q, n=per_source)
                except Exception as e: errs.append(f"Semantic Scholar: {e}")
            if use_openalex:
                try: papers += search_openalex(q, n=per_source)
                except Exception as e: errs.append(f"OpenAlex: {e}")
            if crew_web:
                wres, werr = web_search(q, max_results=per_source, read_pages=crew_web_read)
                if werr:
                    errs.append(werr)
                elif wres:
                    for r in wres:
                        papers.append({"source": "Web", "title": r["title"] or "web result",
                                       "authors": "", "published": "web", "pdf_url": "",
                                       "landing_url": r["url"], "abstract": r["snippet"],
                                       "doi": "", "full_text": "", "venue": "web"})
            return papers, errs

        if any_search or crew_web:
            if query.strip():
                with st.status("🧑‍🔬 Planning & running the search…", expanded=True) as sbox:
                    def _status(msg):
                        sbox.write(msg)
                    cands, search_rounds = run_iterative_paper_search(
                        query, global_model_string, _rc_search_once,
                        api_base=selected_api_base, api_key=selected_api_key, num_ctx=num_ctx,
                        max_rounds=int(rc_max_rounds), status_cb=_status)
                    n_searches = sum(1 for r in search_rounds if r["type"] == "search")
                    sbox.update(label=f"🧑‍🔬 Searched {n_searches} round(s) — {len(cands)} unique paper(s)",
                               state="complete")
            else:
                st.warning("Enter a research topic first, or add your own papers below.")
        else:
            st.info("No sources enabled — add your own papers below, or enable a source in "
                    "📁 Sources & Retrieval above.")

        if up_files:
            try: cands += uploaded_pdfs_to_candidates(up_files, chars=chars_per_paper)
            except Exception as e: errors.append(f"Uploads: {e}")
        if folder_path.strip():
            fol, ferr = folder_pdfs_to_candidates(folder_path.strip(), chars=chars_per_paper)
            if ferr: errors.append(ferr)
            cands += fol

        seen, unique = set(), []
        for p in cands:
            k = re.sub(r"\W+", "", p["title"].lower())[:60]
            if k and k not in seen: seen.add(k); unique.append(p)
        st.session_state.candidates = unique
        st.session_state.research_query_used = query
        st.session_state.pop("research_result", None)
        for e in errors: st.warning(e)
        st.rerun()

    candidates = st.session_state.get("candidates", [])
    if candidates:
        st.markdown(f"#### 🧑‍⚖️ Paper Gatekeeper — {len(candidates)} found.")
        cc1, cc2 = st.columns(2)
        with cc1:
            if st.button("✅ Select all", key="rc_select_all"):
                for i in range(len(candidates)): st.session_state[f"pick_{i}"] = True
        with cc2:
            if st.button("⬜ Clear all", key="rc_clear_all"):
                for i in range(len(candidates)): st.session_state[f"pick_{i}"] = False

        for i, p in enumerate(candidates):
            col_a, col_b = st.columns([0.06, 0.94])
            with col_a:
                st.checkbox(f"Include paper {i+1}", value=st.session_state.get(f"pick_{i}", True), key=f"pick_{i}", label_visibility="collapsed")
            with col_b:
                with st.expander(f"{p['source']} — {p['title']} ({p['published']})"):
                    st.markdown(p["abstract"])

        selected = [p for i, p in enumerate(candidates) if st.session_state.get(f"pick_{i}", True)]
        st.info(f"**{len(selected)}** paper(s) selected.")

        if st.button("🚀 Analyze selected papers", type="primary", disabled=not selected):
            log_event(_CURRENT_USER, "research_crew_run", detail=f"{len(selected)} paper(s)")
            _crew_topic = st.session_state.get('research_query_used', query)
            add_memory_fact(_CURRENT_USER, f"Researched the topic: {_crew_topic}")
            save_project_source_pdfs(CREW_PDFS_DIR, _crew_topic, selected)
            # IMPORTANT: The 'Crew' tab uses the 'crewai' library, which fundamentally relies on 'litellm'.
            # If litellm is not installed, the agents here will fail to run.
            st.markdown("#### 🕵️ Crew Analysis Live Stream")
            log_ph = st.empty()
            log_q = queue.Queue()
            holder = {}
            sel_copy = [dict(p) for p in selected]

            def agent_cb(name):
                def cb(step):
                    ts = datetime.now().strftime("%H:%M:%S")
                    msg = f"\n---\n##### 🤖 {name} · `{ts}`\n"
                    th = getattr(step, "thought", None)
                    out = getattr(step, "output", None)
                    if th: msg += f"💡 **Thinking:**\n{th}\n\n"
                    if out: msg += f"✅ **Answer:**\n{out}\n\n"
                    log_q.put(msg)
                return cb

            def pipeline(papers, cfg, out_dict, ctx):
                add_script_run_ctx(threading.current_thread(), ctx)
                try:
                    for p in papers: ensure_fulltext(p, chars=cfg["chars"])
                    context_block = papers_to_context_block(papers)

                    if cfg.get("web"):
                        wres, werr = web_search(cfg["query"], max_results=5, read_pages=cfg.get("web_read", True), chars_per_page=4000)
                        if wres: context_block += "\n\n" + results_to_context(wres, use_pages=cfg.get("web_read", True))

                    os.environ["OLLAMA_CONTEXT_LENGTH"] = str(cfg["num_ctx"])
                    if cfg.get("api_key"):
                        os.environ["OLLAMA_API_KEY"] = cfg["api_key"]
                        os.environ["GEMINI_API_KEY"] = cfg["api_key"]
                        os.environ["ANTHROPIC_API_KEY"] = cfg["api_key"]
                        os.environ["OPENAI_API_KEY"] = cfg["api_key"]  # covers Qwen/Custom via litellm's openai/ prefix

                    llm = LLM(model=cfg["crew_model_string"], base_url=cfg["base_url"], temperature=cfg["temperature"], max_tokens=cfg["max_tokens"])

                    analyst = Agent(role="Paper Analyst", goal="Extract methods and data.", backstory="Summarize context details accurately.", llm=llm, step_callback=agent_cb("Paper Analyst"))
                    writer = Agent(role="Scientific Writer & Reviewer", goal="Write reviews.", backstory="Structured expert formatting.", llm=llm, step_callback=agent_cb("Writer & Reviewer"))
                    strategist = Agent(role="Research & Innovation Strategist", goal="Propose ideas.", backstory="Ground concepts cleanly.", llm=llm, step_callback=agent_cb("Strategist"))

                    task_specs = cfg["tasks"]
                    task_objs = []
                    for idx, spec in enumerate(task_specs):
                        desc = (spec.get("description") or "").replace("{topic}", cfg["query"])
                        if idx == 0: desc += "\n\n----- PAPERS -----\n" + context_block
                        if idx == 0 and globals().get("USER_MEMORY_CONTEXT"):
                            desc += "\n\n----- " + globals()["USER_MEMORY_CONTEXT"]
                        
                        ag = {"Paper Analyst": analyst, "Scientific Writer & Reviewer": writer, "Research & Innovation Strategist": strategist}.get(spec.get("agent"), analyst)
                        kw = dict(description=desc, expected_output=spec.get("expected_output") or "Output", agent=ag)
                        if task_objs: kw["context"] = [task_objs[-1]]
                        task_objs.append(Task(**kw))

                    crew = Crew(agents=[analyst, writer, strategist], tasks=task_objs, verbose=True)
                    result = str(crew.kickoff())
                    out_dict.update(result=result, bib=to_bibtex(papers), ts=datetime.now().strftime("%Y%m%d_%H%M%S"))
                except Exception as e: out_dict["error"] = str(e)
                finally: log_q.put(None)

            cfg = dict(base_url=selected_api_base, crew_model_string=crew_model_string, temperature=crew_temperature,
                       num_ctx=num_ctx, max_tokens=crew_max_tokens, chars=chars_per_paper,
                       query=st.session_state.get("research_query_used", query), api_key=selected_api_key,
                       web=st.session_state.get("crew_web", False), web_read=st.session_state.get("crew_web_read", True),
                       tasks=[dict(t) for t in st.session_state.get("tasks", DEFAULT_TASKS)])
                       
            worker = threading.Thread(target=pipeline, args=(sel_copy, cfg, holder, get_script_run_ctx()))
            worker.start()
            buf = ""
            while True:
                try:
                    ch = log_q.get(timeout=0.1)
                    if ch is None: break
                    buf += ch
                    log_ph.markdown(buf)
                except queue.Empty:
                    if not worker.is_alive(): break
            worker.join()
            if "result" in holder: st.session_state.research_result = holder
            elif "error" in holder: st.error(f"Error: {holder['error']}")

    res = st.session_state.get("research_result")
    if res:
        st.success("🎉 Complete!")
        rc1, rc2 = st.columns(2)
        with rc1:
            st.download_button("📄 Report (.md)", data=res["result"], file_name=f"report_{res['ts']}.md")
        with rc2:
            rc_pdf = md_to_pdf_bytes(res["result"], title="Research Crew Report")
            if rc_pdf:
                st.download_button("📕 Report (.pdf)", data=rc_pdf, file_name=f"report_{res['ts']}.pdf",
                                   mime="application/pdf")
            else:
                st.caption("PDF needs: pip install reportlab")
        st.markdown(res["result"])

# =============================== SURVEY TAB =================================
with tab_survey:
    st.subheader("📚 Make a Survey")
    st.caption("Searches every enabled source (arXiv, Semantic Scholar, OpenAlex, web) plus your "
               "own uploaded/local papers, lets you pick which to include, then writes a long, "
               "structured survey — section by section — with citations, figures, and LaTeX export.")

    st.markdown("#### 📁 Projects")
    st.caption("Save your progress — topic, gathered sources, and the written survey — and pick "
              "it back up later, even after closing the app. The PDFs of papers you use are kept "
              "in their own folder alongside the project too.")
    sv_projects = list_survey_projects()
    svpc1, svpc2 = st.columns([4, 1])
    with svpc1:
        sv_selected_project = st.selectbox(
            "Saved projects", ["🆕 (start new)"] + sv_projects, key="sv_project_picker")
    with svpc2:
        st.markdown("<div style='height:28px'></div>", unsafe_allow_html=True)
        if st.button("🗑️ Delete", disabled=(sv_selected_project == "🆕 (start new)"), use_container_width=True, key="sv_proj_delete"):
            delete_survey_project(sv_selected_project)
            if st.session_state.get("sv_active_project") == sv_selected_project:
                st.session_state.pop("sv_active_project", None)
            st.session_state.pop("_sv_loaded_project", None)
            st.rerun()

    # Auto-load: as soon as the dropdown's selection differs from what's currently loaded,
    # load it — no separate Load button/click needed.
    if sv_selected_project != "🆕 (start new)" and sv_selected_project != st.session_state.get("_sv_loaded_project"):
        payload, err = load_survey_project(sv_selected_project)
        if err:
            st.error(f"Could not load: {err}")
        else:
            for k in SURVEY_STATE_KEYS:
                if k in payload:
                    st.session_state[k] = payload[k]
            st.session_state["sv_active_project"] = sv_selected_project
            st.session_state["_sv_loaded_project"] = sv_selected_project
            st.session_state["_pending_sv_project_name"] = sv_selected_project
            st.rerun()
    elif sv_selected_project == "🆕 (start new)" and st.session_state.get("_sv_loaded_project"):
        st.session_state["_sv_loaded_project"] = None

    if "_pending_sv_project_name" in st.session_state:
        st.session_state["sv_project_name"] = st.session_state.pop("_pending_sv_project_name")

    svsc1, svsc2 = st.columns([3, 1])
    with svsc1:
        sv_save_name = st.text_input(
            "Project name", value=st.session_state.get("sv_active_project", ""),
            key="sv_project_name", label_visibility="collapsed",
            placeholder="Auto-named from your topic once you start — edit here to rename")
    with svsc2:
        if st.button("✏️ Rename", disabled=not sv_save_name.strip(), use_container_width=True, key="sv_proj_rename"):
            old_name = st.session_state.get("sv_active_project")
            new_name = sv_save_name.strip()
            state = {k: st.session_state.get(k) for k in SURVEY_STATE_KEYS}
            ok, err = save_survey_project(new_name, state)
            if ok:
                if old_name and old_name != new_name:
                    delete_survey_project(old_name)
                st.session_state["sv_active_project"] = new_name
                st.success(f"Renamed to '{new_name}'.")
            else:
                st.error(f"Could not rename: {err}")

    if st.session_state.get("sv_active_project"):
        st.caption(f"📌 Working on: **{st.session_state['sv_active_project']}** "
                  f"— saves automatically as you gather sources and write the survey.")
        _sv_pdfs = list_project_pdfs(SURVEY_PDFS_DIR, st.session_state["sv_active_project"])
        if _sv_pdfs:
            with st.expander(f"📄 {len(_sv_pdfs)} saved PDF(s) for this project"):
                for _fn in _sv_pdfs:
                    _fpath = os.path.join(SURVEY_PDFS_DIR, re.sub(r"[^A-Za-z0-9_-]+", "_", st.session_state["sv_active_project"]), _fn)
                    try:
                        with open(_fpath, "rb") as _fh:
                            st.download_button(f"⬇️ {_fn}", data=_fh.read(), file_name=_fn, key=f"svpdf_{_fn}")
                    except Exception:
                        st.caption(f"⚠️ {_fn} (couldn't read)")

    if st.button("🆕 New Project (clears current work)", key="sv_reset"):
        for key in list(st.session_state.keys()):
            if key.startswith("sv_") or key.startswith("svpick_") or key in ("survey_result", "_sv_loaded_project"):
                del st.session_state[key]
        st.rerun()

    st.divider()

    sv_topic = st.text_input("Survey topic:", "human body communication for wearables",
                             key="sv_topic")

    sc1, sc2, sc3 = st.columns(3)
    with sc1:
        sv_per_source = st.slider("Papers per source", 3, 20, 8, key="sv_per")
    with sc2:
        sv_detail = st.select_slider("Length / detail",
                                     options=["Concise", "Standard", "Comprehensive"],
                                     value="Standard", key="sv_detail")
    with sc3:
        sv_digest_chars = st.select_slider("Evidence per paper (chars)",
                                           options=[800, 1500, 2500, 4000], value=1500,
                                           key="sv_digest")

    sc4, sc5, sc6 = st.columns(3)
    with sc4:
        sv_web = st.checkbox("🌐 Also include web results", value=False, key="sv_web")
    with sc5:
        sv_fulltext = st.checkbox("📄 Download full text for open-access PDFs (slower, deeper)",
                                  value=False, key="sv_full")
    with sc6:
        sv_make_figures = st.checkbox("📊 Auto-generate figures with AI-written code", value=True, key="sv_figs",
                                      help="For each section (except Conclusion), the model writes its "
                                           "own matplotlib Python code — a block diagram, flowchart, "
                                           "taxonomy tree, or data chart, whichever fits that section "
                                           "best — which is then executed and embedded. Not a fixed "
                                           "template; the model decides what to draw each time.")

    sv_max_rounds = st.slider(
        "Max search rounds", 1, 5, 3, key="sv_max_rounds",
        help="The model turns your topic into search terms, checks whether what it finds is "
             "relevant/sufficient, and searches again with different terms if not — up to this "
             "many rounds.")

    sv_selection_mode = st.radio(
        "Source selection", ["🤖 Automatic — use everything found", "🧑‍⚖️ Manual — I'll pick"],
        horizontal=True, key="sv_selection_mode")

    with st.expander("📁 Add your own papers to the survey", expanded=False):
        sv_uploads = st.file_uploader("Upload PDFs", type=["pdf"], accept_multiple_files=True,
                                      key="sv_uploads")
        sv_folder = st.text_input("…or a local folder (scanned recursively)", value="",
                                  key="sv_folder")

    detail_mult = {"Concise": 0.6, "Standard": 1.0, "Comprehensive": 1.7}[sv_detail]

    # ---- Stage 1: Gather sources, then screen for relevance ------------------
    if st.button("🔍 Gather Sources", type="primary"):
        papers, errors = [], []

        def _sv_search_once(q):
            found, errs = [], []
            if use_arxiv:
                try: found += search_arxiv_meta(q, n=sv_per_source, sort=sort_choice)
                except Exception as e: errs.append(f"arXiv: {e}")
            if use_s2:
                try: found += search_semantic_scholar(q, n=sv_per_source)
                except Exception as e: errs.append(f"Semantic Scholar: {e}")
            if use_openalex:
                try: found += search_openalex(q, n=sv_per_source)
                except Exception as e: errs.append(f"OpenAlex: {e}")
            if sv_web:
                wres, werr = web_search(q, max_results=6, read_pages=False)
                if werr:
                    errs.append(werr)
                elif wres:
                    for r in wres:
                        found.append({"source": "Web", "title": r["title"] or "web result",
                                      "authors": "", "published": "web", "pdf_url": "",
                                      "landing_url": r["url"], "abstract": r["snippet"],
                                      "doi": "", "full_text": "", "venue": "web"})
            return found, errs

        if sv_topic.strip() and (use_arxiv or use_s2 or use_openalex or sv_web):
            with st.status("🧑‍🔬 Planning & running the search…", expanded=True) as sbox:
                def _status(msg):
                    sbox.write(msg)
                papers, search_rounds = run_iterative_paper_search(
                    sv_topic, global_model_string, _sv_search_once,
                    api_base=selected_api_base, api_key=selected_api_key, num_ctx=num_ctx,
                    max_rounds=int(sv_max_rounds), status_cb=_status)
                n_searches = sum(1 for r in search_rounds if r["type"] == "search")
                sbox.update(label=f"🧑‍🔬 Searched {n_searches} round(s) — {len(papers)} unique paper(s)",
                           state="complete")
        elif not sv_topic.strip():
            st.warning("Enter a survey topic first, or add your own papers below.")

        if sv_uploads:
            with st.spinner("Reading uploaded PDFs…"):
                try: papers += uploaded_pdfs_to_candidates(sv_uploads, chars=8000)
                except Exception as e: errors.append(f"Uploads: {e}")
        if sv_folder.strip():
            with st.spinner("Scanning local folder…"):
                fol, ferr = folder_pdfs_to_candidates(sv_folder.strip(), chars=8000)
                if ferr: errors.append(ferr)
                papers += fol

        seen, unique = set(), []
        for p in papers:
            k = re.sub(r"\W+", "", p["title"].lower())[:60]
            if k and k not in seen:
                seen.add(k); unique.append(p)
        for e in errors:
            st.warning(e)

        removed = []
        if unique and sv_topic.strip():
            with st.status("🧹 Checking relevance to the topic…", expanded=True) as rbox:
                rbox.write(f"Reviewing {len(unique)} gathered source(s) against “{sv_topic}”…")
                unique, removed = filter_relevant_papers(
                    sv_topic, unique, global_model_string, api_base=selected_api_base,
                    api_key=selected_api_key, num_ctx=num_ctx)
                if removed:
                    for r in removed:
                        rbox.write(f"🗑️ Dropped: **{r['title']}** — _{r.get('_removal_reason', 'off-topic')}_")
                rbox.update(label=f"🧹 Kept {len(unique)} relevant source(s), dropped {len(removed)}",
                           state="complete")

        # Reset any stale gatekeeper picks / prior result now that we have a fresh candidate set
        for key in list(st.session_state.keys()):
            if key.startswith("svpick_"):
                del st.session_state[key]
        st.session_state.sv_candidates = unique
        st.session_state.sv_removed = removed
        st.session_state.pop("survey_result", None)
        if not st.session_state.get("sv_active_project"):
            st.session_state["sv_active_project"] = autoname_session(sv_topic, list_survey_projects())
        save_survey_project(st.session_state["sv_active_project"],
                            {k: st.session_state.get(k) for k in SURVEY_STATE_KEYS})
        if not unique:
            st.error("No sources gathered. Enable sources in Research Crew's 📁 Sources & "
                    "Retrieval, add uploads/folder, and enter a topic.")
        st.rerun()

    sv_removed = st.session_state.get("sv_removed", [])
    if sv_removed:
        with st.expander(f"🗑️ {len(sv_removed)} source(s) removed as off-topic"):
            for r in sv_removed:
                st.markdown(f"- **[{r['source']}]** {r['title']} — _{r.get('_removal_reason', 'off-topic')}_")

    # ---- Stage 2: Select (or auto-use) + Write ------------------------------
    sv_candidates = st.session_state.get("sv_candidates", [])
    if sv_candidates:
        is_manual = sv_selection_mode.startswith("🧑‍⚖️")

        if is_manual:
            st.markdown(f"#### 🧑‍⚖️ Source Gatekeeper — {len(sv_candidates)} found")
            gc1, gc2 = st.columns(2)
            with gc1:
                if st.button("✅ Select all", key="sv_select_all"):
                    for i in range(len(sv_candidates)): st.session_state[f"svpick_{i}"] = True
            with gc2:
                if st.button("⬜ Clear all", key="sv_clear_all"):
                    for i in range(len(sv_candidates)): st.session_state[f"svpick_{i}"] = False

            for i, p in enumerate(sv_candidates):
                col_a, col_b = st.columns([0.06, 0.94])
                with col_a:
                    st.checkbox(f"Include source {i+1}", value=st.session_state.get(f"svpick_{i}", True),
                               key=f"svpick_{i}", label_visibility="collapsed")
                with col_b:
                    with st.expander(f"{p['source']} — {p['title']} ({p['published']})"):
                        st.markdown(p["abstract"])
                        render_paper_card_extra(p, key_prefix=f"sv_card_{i}")

            selected_sources = [p for i, p in enumerate(sv_candidates)
                                if st.session_state.get(f"svpick_{i}", True)]
            st.info(f"**{len(selected_sources)}** of {len(sv_candidates)} source(s) selected.")
        else:
            selected_sources = sv_candidates
            st.markdown(f"📋 **{len(sv_candidates)} source(s) gathered** (automatic mode — all included)")
            for i, p in enumerate(sv_candidates):
                with st.expander(f"{p['source']} — {p['title']} ({p['published']})"):
                    st.markdown(p["abstract"])
                    render_paper_card_extra(p, key_prefix=f"sv_autocard_{i}")

        with st.expander("➕ Pull more sources", expanded=False):
            sv_more_query = st.text_input(
                "Additional search query (optional — leave blank to reuse the survey topic):",
                key="sv_more_query")
            if st.button("🔍 Pull More Sources", key="sv_pull_more"):
                topic_to_use = sv_more_query.strip() or sv_topic

                def _sv_search_once_more(q):
                    found, errs = [], []
                    if use_arxiv:
                        try: found += search_arxiv_meta(q, n=sv_per_source, sort=sort_choice)
                        except Exception as e: errs.append(f"arXiv: {e}")
                    if use_s2:
                        try: found += search_semantic_scholar(q, n=sv_per_source)
                        except Exception as e: errs.append(f"Semantic Scholar: {e}")
                    if use_openalex:
                        try: found += search_openalex(q, n=sv_per_source)
                        except Exception as e: errs.append(f"OpenAlex: {e}")
                    if sv_web:
                        wres, werr = web_search(q, max_results=6, read_pages=False)
                        if werr:
                            errs.append(werr)
                        elif wres:
                            for r in wres:
                                found.append({"source": "Web", "title": r["title"] or "web result",
                                              "authors": "", "published": "web", "pdf_url": "",
                                              "landing_url": r["url"], "abstract": r["snippet"],
                                              "doi": "", "full_text": "", "venue": "web"})
                    return found, errs

                with st.status("🧑‍🔬 Searching for more sources…", expanded=True) as mbox:
                    def _mstatus(msg):
                        mbox.write(msg)
                    more_papers, _more_rounds = run_iterative_paper_search(
                        topic_to_use, global_model_string, _sv_search_once_more,
                        api_base=selected_api_base, api_key=selected_api_key, num_ctx=num_ctx,
                        max_rounds=int(sv_max_rounds), status_cb=_mstatus)

                    existing_keys = {re.sub(r"\W+", "", p["title"].lower())[:60] for p in sv_candidates}
                    added = []
                    for p in more_papers:
                        k = re.sub(r"\W+", "", p["title"].lower())[:60]
                        if k and k not in existing_keys:
                            existing_keys.add(k)
                            added.append(p)

                    new_removed = []
                    if added:
                        added, new_removed = filter_relevant_papers(
                            topic_to_use, added, global_model_string, api_base=selected_api_base,
                            api_key=selected_api_key, num_ctx=num_ctx)

                    mbox.update(label=f"🧑‍🔬 Added {len(added)} new relevant source(s)"
                                     + (f", dropped {len(new_removed)} off-topic" if new_removed else ""),
                               state="complete")

                if added:
                    st.session_state.sv_candidates = sv_candidates + added
                    st.session_state.sv_removed = st.session_state.get("sv_removed", []) + new_removed
                    st.rerun()
                else:
                    st.info("No new unique, relevant sources found — try a more specific additional query.")

        if st.button("📝 Write Survey", type="primary", disabled=not selected_sources):
            log_event(_CURRENT_USER, "survey_write", detail=f"{len(selected_sources)} source(s)")
            add_memory_fact(_CURRENT_USER, f"Wrote a literature survey about: {sv_topic}")
            save_project_source_pdfs(SURVEY_PDFS_DIR, sv_topic, selected_sources)
            papers = selected_sources

            if sv_fulltext:
                prog = st.progress(0.0, text="Reading full text…")
                for j, p in enumerate(papers):
                    if p.get("pdf_url") and not p.get("full_text"):
                        ensure_fulltext(p, chars=sv_digest_chars * 2)
                    prog.progress((j + 1) / len(papers))
                prog.empty()

            digest = build_source_digest(papers, per_chars=sv_digest_chars)
            outline = " | ".join(t for t, _, _ in SURVEY_SECTIONS)
            refs = numbered_references(papers)

            full_parts = [f"# Survey: {sv_topic}\n",
                          f"*Generated from {len(papers)} sources on "
                          f"{datetime.now():%Y-%m-%d}.*\n"]
            figs = []

            for title, guidance, base_words in SURVEY_SECTIONS:
                target = int(base_words * detail_mult)
                st.markdown(f"### {title}")
                sys_msg = SURVEY_SYSTEM
                user_msg = (
                    f"SURVEY TOPIC: {sv_topic}\n"
                    f"FULL OUTLINE (for context, do not rewrite other sections): {outline}\n\n"
                    f"NUMBERED SOURCES:\n{digest}\n\n"
                    f"Now write the section '{title}'. {guidance} "
                    f"Aim for roughly {target} words. Cite sources as [n] where relevant.")
                try:
                    txt = st.write_stream(stream_chat(
                        [{"role": "system", "content": sys_msg},
                         {"role": "user", "content": user_msg}],
                        global_model_string, api_base=selected_api_base, temperature=0.5,
                        num_ctx=num_ctx, api_key=selected_api_key))
                except Exception as e:
                    txt = ""
                    st.error(f"Section '{title}' failed: {e}")

                if title == "Abstract":
                    # Safety net: strip any citation markers even if the model slipped one in
                    # despite the instruction, since the abstract must stand alone.
                    txt = re.sub(r"\s?\[\d+(?:\s*,\s*\d+)*\]", "", txt)

                full_parts.append(f"## {title}\n\n{txt}\n")

                # One figure per section, except Conclusion — the model decides what kind of
                # figure fits (block diagram, flowchart, taxonomy, data chart, etc.), writes
                # the matplotlib code itself, and it's executed and embedded right here.
                if sv_make_figures and txt.strip() and not title.startswith("8."):
                    is_abstract = (title == "Abstract")
                    with st.status(f"🎨 Drawing a figure for “{title}”…", expanded=False) as figbox:
                        def _figstatus(msg):
                            figbox.write(msg)
                        fig_result = generate_section_figure(
                            sv_topic, title, txt, global_model_string,
                            is_graphical_abstract=is_abstract, api_base=selected_api_base,
                            api_key=selected_api_key, num_ctx=num_ctx, status_cb=_figstatus)
                        figbox.update(label=f"🎨 Figure for “{title}”: {'done' if fig_result else 'skipped'}",
                                     state="complete")
                    if fig_result:
                        idx = len(figs)
                        figs.append(fig_result)
                        st.image(fig_result["png"], caption=f"Figure {idx + 1}. {fig_result['caption']}")
                        full_parts.append(f"[[FIGURE:{idx}]]\n\n*Figure {idx + 1}. {fig_result['caption']}*\n")

            full_parts.append("## References\n\n" + refs)
            survey_md_master = "\n".join(full_parts)  # keeps [[FIGURE:n]] markers — feeds PDF/LaTeX
            display_md = materialize_markdown_images(survey_md_master, figs)  # real images — for .md/preview

            fig_filenames = {i: f"fig{i}.png" for i in range(len(figs))}
            bib_str = to_bibtex(papers)
            try:
                tex_str = survey_to_latex(survey_md_master, f"Survey: {sv_topic}", figs, fig_filenames)
                tex_zip = build_latex_zip(tex_str, figs, fig_filenames, bib_str=bib_str, tex_filename="survey.tex")
            except Exception as e:
                tex_zip = None
                st.warning(f"LaTeX export failed: {e}")

            st.session_state["survey_result"] = {
                "md": display_md, "md_master": survey_md_master, "figs": figs,
                "bib": bib_str, "tex_zip": tex_zip,
                "ts": datetime.now().strftime("%Y%m%d_%H%M%S")}
            st.success("🎉 Survey complete.")

            # Autosave — no manual "Save Project" click needed. Auto-names from the topic
            # the first time; if a project is already active (e.g. loaded from disk), it
            # just resaves under that same name so it stays up to date.
            if not st.session_state.get("sv_active_project"):
                st.session_state["sv_active_project"] = autoname_session(sv_topic, list_survey_projects())
            _sv_autosave_state = {k: st.session_state.get(k) for k in SURVEY_STATE_KEYS}
            save_survey_project(st.session_state["sv_active_project"], _sv_autosave_state)

    sv = st.session_state.get("survey_result")
    if sv:
        d1, d2, d3, d4, d5 = st.columns(5)
        with d1:
            st.download_button("📄 Survey (.md)", data=sv["md"], file_name=f"survey_{sv['ts']}.md")
        with d2:
            st.download_button("📄 Survey (.txt)", data=sv["md"], file_name=f"survey_{sv['ts']}.txt")
        with d3:
            pdf_bytes = md_to_pdf_bytes(sv["md_master"], title="Survey", figs=sv.get("figs"))
            if pdf_bytes:
                st.download_button("📕 Survey (.pdf)", data=pdf_bytes, file_name=f"survey_{sv['ts']}.pdf", mime="application/pdf")
            else:
                st.caption("PDF needs: pip install reportlab")
        with d4:
            st.download_button("📚 References (.bib)", data=sv["bib"], file_name=f"survey_refs_{sv['ts']}.bib")
        with d5:
            if sv.get("tex_zip"):
                st.download_button("📦 LaTeX (.zip)", data=sv["tex_zip"], file_name=f"survey_{sv['ts']}_latex.zip",
                                   mime="application/zip",
                                   help="Contains survey.tex, figure PNGs, and references.bib — a complete compilable project.")
            else:
                st.caption("LaTeX export unavailable")

# =============================== CHECK NOVELTY TAB ===========================
with tab_novelty:
    st.subheader("🧭 Check Novelty")
    st.caption("Describe your idea. The model searches academic papers and the general web for "
               "prior art, tells you whether it's novel and why, lists related work, and "
               "suggests ways to strengthen it.")

    idea_text = st.text_area(
        "Describe your idea:", height=150, key="nov_idea",
        placeholder="e.g. A wearable that uses human body communication instead of RF radio to "
                    "sync ECG timestamps between two devices worn by the same person, for "
                    "low-power outdoor group fitness monitoring...")

    nc1, nc2, nc3 = st.columns(3)
    with nc1:
        nov_include_web = st.checkbox(
            "🌐 Also search the general web (products, patents, blogs)", value=True, key="nov_web")
    with nc2:
        nov_paper_rounds = st.slider("Max paper search rounds", 1, 5, 3, key="nov_paper_rounds",
                                     help="Uses the sources enabled in Research Crew's "
                                          "📁 Sources & Retrieval (arXiv / Semantic Scholar / OpenAlex).")
    with nc3:
        nov_web_rounds = st.slider("Max web search rounds", 1, 5, 2, key="nov_web_rounds",
                                   disabled=not nov_include_web)

    if st.button("🔎 Check Novelty", type="primary", disabled=not idea_text.strip()):
        log_event(_CURRENT_USER, "novelty_check")
        add_memory_fact(_CURRENT_USER, f"Checked the novelty of this research idea: {idea_text.strip()[:200]}")
        def _nov_search_once(q):
            found, errs = [], []
            if use_arxiv:
                try: found += search_arxiv_meta(q, n=per_source, sort=sort_choice)
                except Exception as e: errs.append(f"arXiv: {e}")
            if use_s2:
                try: found += search_semantic_scholar(q, n=per_source)
                except Exception as e: errs.append(f"Semantic Scholar: {e}")
            if use_openalex:
                try: found += search_openalex(q, n=per_source)
                except Exception as e: errs.append(f"OpenAlex: {e}")
            return found, errs

        papers, web_hits = [], []
        with st.status("🧭 Searching for prior art…", expanded=True) as sbox:
            def _status(msg):
                sbox.write(msg)

            if use_arxiv or use_s2 or use_openalex:
                papers, _paper_rounds = run_iterative_paper_search(
                    idea_text, global_model_string, _nov_search_once,
                    api_base=selected_api_base, api_key=selected_api_key, num_ctx=num_ctx,
                    max_rounds=int(nov_paper_rounds), status_cb=_status)
            else:
                sbox.write("💭 No paper databases enabled in Research Crew's 📁 Sources & "
                          "Retrieval — skipping paper search.")

            if nov_include_web:
                _, web_rounds = run_iterative_search(
                    idea_text, [], global_model_string,
                    api_base=selected_api_base, api_key=selected_api_key, num_ctx=num_ctx,
                    max_results=6, read_pages=False, max_rounds=int(nov_web_rounds),
                    status_cb=_status)
                for rd in web_rounds:
                    if rd["type"] == "search" and rd.get("results"):
                        web_hits.extend(rd["results"])

            sbox.update(label=f"🧭 Gathered {len(papers)} paper(s) and {len(web_hits)} web source(s)",
                       state="complete")

        if not papers and not web_hits:
            st.warning("No prior-art sources found. Enable at least one source in Research "
                      "Crew's 📁 Sources & Retrieval, or turn on general web search above, then "
                      "try again.")
        else:
            digest, total_n = build_novelty_digest(papers, web_hits, per_chars=chars_per_paper)
            references = numbered_novelty_references(papers, web_hits)

            st.markdown("#### 📝 Novelty Report")
            user_msg = (
                f"IDEA:\n{idea_text}\n\n"
                f"NUMBERED PRIOR-ART SOURCES ({total_n} total):\n{digest}\n\n"
                f"Write the novelty report now."
            )
            try:
                report_txt = st.write_stream(stream_chat(
                    [{"role": "system", "content": NOVELTY_SYSTEM},
                     {"role": "user", "content": user_msg}],
                    global_model_string, api_base=selected_api_base, temperature=0.3,
                    num_ctx=num_ctx, api_key=selected_api_key, max_tokens=crew_max_tokens))
            except Exception as e:
                report_txt = ""
                st.error(f"Novelty check failed: {e}")

            if report_txt:
                st.markdown("##### 📚 Sources checked")
                st.markdown(references or "(none)")

                ts = datetime.now().strftime("%Y%m%d_%H%M%S")
                full_md = (
                    f"# Novelty Report\n\n"
                    f"*Checked on {datetime.now():%Y-%m-%d}, against {total_n} source(s).*\n\n"
                    f"## Original Idea\n\n{idea_text}\n\n"
                    f"{report_txt}\n\n"
                    f"## Sources\n\n{references}"
                )
                st.session_state["novelty_result"] = {"md": full_md, "ts": ts}

    nov = st.session_state.get("novelty_result")
    if nov:
        nv1, nv2 = st.columns(2)
        with nv1:
            st.download_button("📄 Novelty Report (.md)", data=nov["md"], file_name=f"novelty_{nov['ts']}.md")
        with nv2:
            nov_pdf = md_to_pdf_bytes(nov["md"], title="Novelty Report")
            if nov_pdf:
                st.download_button("📕 Novelty Report (.pdf)", data=nov_pdf, file_name=f"novelty_{nov['ts']}.pdf",
                                   mime="application/pdf")
            else:
                st.caption("PDF needs: pip install reportlab")

# =============================== WRITE PAPER TAB =============================
with tab_writepaper:
    st.subheader("✍️ Write Paper")
    st.caption("A guided pipeline: describe your idea → the model searches the literature → "
               "you curate the sources → check novelty → approve/adjust the idea → review an "
               "editable outline → the paper gets written section by section with figures → "
               "fine-tune with AI → export.")

    st.markdown("#### 📁 Projects")
    st.caption("Save your progress at any stage — idea, gathered sources, novelty report, "
              "outline, and the written draft with its figures — and pick it back up later, "
              "even after closing the app.")
    wp_projects = list_wp_projects()
    pc1, pc2, pc3 = st.columns([3, 1, 1])
    with pc1:
        wp_selected_project = st.selectbox(
            "Saved projects", ["(none)"] + wp_projects, key="wp_project_picker",
            label_visibility="collapsed")
    with pc2:
        if st.button("📂 Load", disabled=(wp_selected_project == "(none)"), use_container_width=True, key="wp_proj_load"):
            payload, err = load_wp_project(wp_selected_project)
            if err:
                st.error(f"Could not load: {err}")
            else:
                # Safe to set widget-backed keys (wp_idea, wp_web, ...) here: this section runs
                # before any of those widgets are instantiated later in this same script pass.
                for k in WP_STATE_KEYS:
                    if k in payload:
                        st.session_state[k] = payload[k]
                st.session_state["wp_active_project"] = wp_selected_project
                # "wp_project_name" is a widget-owned key once rendered once — queue the new
                # value for next run rather than relying on value=, which gets silently ignored
                # for an already-instantiated widget (same class of issue as the per-section
                # editor fix elsewhere in this tab).
                st.session_state["_pending_wp_project_name"] = wp_selected_project
                st.rerun()
    with pc3:
        if st.button("🗑️ Delete", disabled=(wp_selected_project == "(none)"), use_container_width=True, key="wp_proj_delete"):
            delete_wp_project(wp_selected_project)
            if st.session_state.get("wp_active_project") == wp_selected_project:
                st.session_state.pop("wp_active_project", None)
            st.rerun()

    if "_pending_wp_project_name" in st.session_state:
        st.session_state["wp_project_name"] = st.session_state.pop("_pending_wp_project_name")

    sc1, sc2 = st.columns([3, 1])
    with sc1:
        wp_save_name = st.text_input(
            "Project name", value=st.session_state.get("wp_active_project", ""),
            key="wp_project_name", label_visibility="collapsed",
            placeholder="Name this project to save it…")
    with sc2:
        if st.button("💾 Save Project", disabled=not wp_save_name.strip(), use_container_width=True, key="wp_proj_save"):
            state = {k: st.session_state.get(k) for k in WP_STATE_KEYS}
            ok, err = save_wp_project(wp_save_name.strip(), state)
            if ok:
                st.session_state["wp_active_project"] = wp_save_name.strip()
                st.success(f"Saved '{wp_save_name.strip()}'.")
            else:
                st.error(f"Could not save: {err}")

    if st.session_state.get("wp_active_project"):
        st.caption(f"📌 Working on: **{st.session_state['wp_active_project']}** "
                  f"— remember to 💾 Save Project after changes.")

    if st.button("🆕 New Project (clears current work)", key="wp_reset"):
        for key in list(st.session_state.keys()):
            if key.startswith("wp_") or key.startswith("wppick_"):
                del st.session_state[key]
        st.rerun()

    st.divider()

    # ---- Step 1: Idea + search configuration --------------------------------
    st.markdown("#### 1️⃣ Your idea")
    wp_idea = st.text_area(
        "Describe your idea:", height=140, key="wp_idea",
        placeholder="e.g. A wearable that uses human body communication instead of RF radio to "
                    "sync ECG timestamps between two devices worn by the same person, for "
                    "low-power outdoor group fitness monitoring...")

    wc1, wc2, wc3 = st.columns(3)
    with wc1:
        wp_web = st.checkbox("🌐 Also include general web results", value=False, key="wp_web")
    with wc2:
        wp_web_full = st.checkbox("📄 Read full web page content", value=False, key="wp_web_full",
                                  disabled=not wp_web,
                                  help="Fetch and read each web result's actual page content, not "
                                       "just the search snippet — slower, but gives the writer "
                                       "much more to work with.")
    with wc3:
        wp_max_rounds = st.slider(
            "Max search rounds", 1, 5, 3, key="wp_max_rounds",
            help="Uses the sources enabled in Research Crew's 📁 Sources & Retrieval "
                 "(arXiv / Semantic Scholar / OpenAlex) and 'Results per source' for paper count.")

    if st.button("🔍 Search Literature", type="primary", disabled=not wp_idea.strip()):
        def _wp_search_once(q):
            found, errs = [], []
            if use_arxiv:
                try: found += search_arxiv_meta(q, n=per_source, sort=sort_choice)
                except Exception as e: errs.append(f"arXiv: {e}")
            if use_s2:
                try: found += search_semantic_scholar(q, n=per_source)
                except Exception as e: errs.append(f"Semantic Scholar: {e}")
            if use_openalex:
                try: found += search_openalex(q, n=per_source)
                except Exception as e: errs.append(f"OpenAlex: {e}")
            if wp_web:
                wres, werr = web_search(q, max_results=per_source, read_pages=wp_web_full)
                if werr:
                    errs.append(werr)
                elif wres:
                    for r in wres:
                        page = r.get("page_text", "")
                        full_text = page if (wp_web_full and page and not page.startswith("[")) else ""
                        found.append({"source": "Web", "title": r["title"] or "web result",
                                      "authors": "", "published": "web", "pdf_url": "",
                                      "landing_url": r["url"], "abstract": r["snippet"],
                                      "doi": "", "full_text": full_text, "venue": "web"})
            return found, errs

        with st.status("🧑‍🔬 Planning & running the search…", expanded=True) as sbox:
            def _status(msg):
                sbox.write(msg)
            papers, search_rounds = run_iterative_paper_search(
                wp_idea, global_model_string, _wp_search_once,
                api_base=selected_api_base, api_key=selected_api_key, num_ctx=num_ctx,
                max_rounds=int(wp_max_rounds), status_cb=_status)
            n_searches = sum(1 for r in search_rounds if r["type"] == "search")
            sbox.update(label=f"🧑‍🔬 Searched {n_searches} round(s) — {len(papers)} unique paper(s)",
                       state="complete")

        removed = []
        if papers:
            with st.status("🧹 Checking relevance to the idea…", expanded=True) as rbox:
                rbox.write(f"Reviewing {len(papers)} gathered source(s)…")
                papers, removed = filter_relevant_papers(
                    wp_idea, papers, global_model_string, api_base=selected_api_base,
                    api_key=selected_api_key, num_ctx=num_ctx)
                for r in removed:
                    rbox.write(f"🗑️ Dropped: **{r['title']}** — _{r.get('_removal_reason', 'off-topic')}_")
                rbox.update(label=f"🧹 Kept {len(papers)} relevant source(s), dropped {len(removed)}",
                           state="complete")

        # Fresh search invalidates every downstream stage
        for key in list(st.session_state.keys()):
            if key.startswith("wppick_"):
                del st.session_state[key]
        for key in ("wp_sources", "wp_novelty_report", "wp_approved_idea", "wp_title",
                   "wp_sections", "wp_paper"):
            st.session_state.pop(key, None)
        st.session_state.wp_candidates = papers
        st.session_state.wp_removed = removed
        if not papers:
            st.error("No sources gathered. Enable a source in Research Crew's 📁 Sources & "
                    "Retrieval, or turn on general web results above.")
        st.rerun()

    wp_removed = st.session_state.get("wp_removed", [])
    if wp_removed:
        with st.expander(f"🗑️ {len(wp_removed)} source(s) auto-removed as off-topic"):
            for r in wp_removed:
                st.markdown(f"- **[{r['source']}]** {r['title']} — _{r.get('_removal_reason', 'off-topic')}_")

    # ---- Step 2: Curate sources ----------------------------------------------
    wp_candidates = st.session_state.get("wp_candidates", [])
    if wp_candidates:
        st.divider()
        st.markdown(f"#### 2️⃣ Curate sources — {len(wp_candidates)} found")
        st.caption("Already screened for relevance automatically — remove any more you don't want included.")
        gc1, gc2 = st.columns(2)
        with gc1:
            if st.button("✅ Select all", key="wp_select_all"):
                for i in range(len(wp_candidates)): st.session_state[f"wppick_{i}"] = True
        with gc2:
            if st.button("⬜ Clear all", key="wp_clear_all"):
                for i in range(len(wp_candidates)): st.session_state[f"wppick_{i}"] = False

        for i, p in enumerate(wp_candidates):
            col_a, col_b = st.columns([0.06, 0.94])
            with col_a:
                st.checkbox(f"Include source {i+1}", value=st.session_state.get(f"wppick_{i}", True),
                           key=f"wppick_{i}", label_visibility="collapsed")
            with col_b:
                with st.expander(f"{p['source']} — {p['title']} ({p['published']})"):
                    st.markdown(p["abstract"])
                    render_paper_card_extra(p, key_prefix=f"wp_card_{i}")

        wp_selected = [p for i, p in enumerate(wp_candidates) if st.session_state.get(f"wppick_{i}", True)]
        st.info(f"**{len(wp_selected)}** of {len(wp_candidates)} source(s) selected.")

        with st.expander("➕ Pull more sources", expanded=False):
            wp_more_query = st.text_input(
                "Additional search query (optional — leave blank to reuse your idea):",
                key="wp_more_query")
            if st.button("🔍 Pull More Sources", key="wp_pull_more"):
                topic_to_use = wp_more_query.strip() or wp_idea

                def _wp_search_once_more(q):
                    found, errs = [], []
                    if use_arxiv:
                        try: found += search_arxiv_meta(q, n=per_source, sort=sort_choice)
                        except Exception as e: errs.append(f"arXiv: {e}")
                    if use_s2:
                        try: found += search_semantic_scholar(q, n=per_source)
                        except Exception as e: errs.append(f"Semantic Scholar: {e}")
                    if use_openalex:
                        try: found += search_openalex(q, n=per_source)
                        except Exception as e: errs.append(f"OpenAlex: {e}")
                    if wp_web:
                        wres, werr = web_search(q, max_results=per_source, read_pages=wp_web_full)
                        if werr:
                            errs.append(werr)
                        elif wres:
                            for r in wres:
                                page = r.get("page_text", "")
                                full_text = page if (wp_web_full and page and not page.startswith("[")) else ""
                                found.append({"source": "Web", "title": r["title"] or "web result",
                                              "authors": "", "published": "web", "pdf_url": "",
                                              "landing_url": r["url"], "abstract": r["snippet"],
                                              "doi": "", "full_text": full_text, "venue": "web"})
                    return found, errs

                with st.status("🧑‍🔬 Searching for more sources…", expanded=True) as mbox:
                    def _mstatus(msg):
                        mbox.write(msg)
                    more_papers, _more_rounds = run_iterative_paper_search(
                        topic_to_use, global_model_string, _wp_search_once_more,
                        api_base=selected_api_base, api_key=selected_api_key, num_ctx=num_ctx,
                        max_rounds=int(wp_max_rounds), status_cb=_mstatus)

                    existing_keys = {re.sub(r"\W+", "", p["title"].lower())[:60] for p in wp_candidates}
                    added = []
                    for p in more_papers:
                        k = re.sub(r"\W+", "", p["title"].lower())[:60]
                        if k and k not in existing_keys:
                            existing_keys.add(k)
                            added.append(p)

                    new_removed = []
                    if added:
                        added, new_removed = filter_relevant_papers(
                            topic_to_use, added, global_model_string, api_base=selected_api_base,
                            api_key=selected_api_key, num_ctx=num_ctx)

                    mbox.update(label=f"🧑‍🔬 Added {len(added)} new relevant source(s)"
                                     + (f", dropped {len(new_removed)} off-topic" if new_removed else ""),
                               state="complete")

                if added:
                    st.session_state.wp_candidates = wp_candidates + added
                    st.session_state.wp_removed = st.session_state.get("wp_removed", []) + new_removed
                    st.rerun()
                else:
                    st.info("No new unique, relevant sources found — try a more specific additional query.")

        with st.expander("📖 Read full content for selected sources", expanded=False):
            st.caption("One setting, applied to everything: academic PDFs get their full text "
                      "extracted, web results get their actual page content fetched (instead of "
                      "just the search snippet) — both capped at the same character limit.")
            wp_max_chars_per_paper = st.slider(
                "Max characters per source", 1000, 20000, 4000, step=500, key="wp_max_chars_per_paper")
            if st.button("📖 Read Full Content for All Selected Sources", key="wp_read_all"):
                prog = st.progress(0.0, text="Reading…")
                n = len(wp_selected) or 1
                for j, p in enumerate(wp_selected):
                    if not p.get("full_text"):
                        ensure_fulltext(p, chars=wp_max_chars_per_paper)
                    prog.progress((j + 1) / n, text=f"Read {j + 1}/{n}: {p['title'][:60]}")
                prog.empty()
                st.success(f"Read full content for {len(wp_selected)} selected source(s) "
                          f"(already-read sources were skipped).")
                st.rerun()

        if st.button("✅ Confirm Sources", type="primary", disabled=not wp_selected):
            st.session_state.wp_sources = wp_selected
            for key in ("wp_novelty_report", "wp_approved_idea", "wp_title", "wp_sections", "wp_paper"):
                st.session_state.pop(key, None)
            st.rerun()

    # ---- Step 3: Novelty check + idea approval -------------------------------
    wp_sources = st.session_state.get("wp_sources", [])
    if wp_sources:
        st.divider()
        st.markdown(f"#### 3️⃣ Novelty check — using {len(wp_sources)} confirmed source(s)")

        if st.button("🔎 Check Novelty & Idea Strength"):
            digest, total_n = build_novelty_digest(wp_sources, [], per_chars=1500)
            references = numbered_novelty_references(wp_sources, [])
            user_msg = (f"IDEA:\n{wp_idea}\n\nNUMBERED PRIOR-ART SOURCES ({total_n} total):\n{digest}\n\n"
                       f"Write the novelty report now.")
            try:
                report_txt = st.write_stream(stream_chat(
                    [{"role": "system", "content": NOVELTY_SYSTEM},
                     {"role": "user", "content": user_msg}],
                    global_model_string, api_base=selected_api_base, temperature=0.3,
                    num_ctx=num_ctx, api_key=selected_api_key, max_tokens=crew_max_tokens))
                st.session_state.wp_novelty_report = report_txt
            except Exception as e:
                st.error(f"Novelty check failed: {e}")

        if st.session_state.get("wp_novelty_report"):
            st.markdown(st.session_state["wp_novelty_report"])
            st.markdown("##### ✏️ Approve or adjust your idea")
            st.caption("Edit the idea below to incorporate the suggestions above, or leave it as-is.")
            wp_idea_edit = st.text_area("Idea (editable):", value=wp_idea, height=140, key="wp_idea_edit")
            if st.button("✅ Approve Idea & Continue", type="primary"):
                st.session_state.wp_approved_idea = wp_idea_edit.strip() or wp_idea
                for key in ("wp_title", "wp_sections", "wp_paper"):
                    st.session_state.pop(key, None)
                st.rerun()

    # ---- Step 4: Outline suggestion + editing --------------------------------
    wp_approved_idea = st.session_state.get("wp_approved_idea", "")
    if wp_approved_idea:
        st.divider()
        st.markdown("#### 4️⃣ Title & outline")

        if not st.session_state.get("wp_sections"):
            if st.button("📝 Suggest Title & Outline", type="primary"):
                digest = build_source_digest(wp_sources, per_chars=wp_max_chars_per_paper)
                with st.spinner("Drafting a title and outline…"):
                    title, sections = suggest_paper_outline(
                        wp_approved_idea, digest, global_model_string, api_base=selected_api_base,
                        api_key=selected_api_key, num_ctx=num_ctx)
                for s in sections:
                    s.setdefault("figure", "abstract" not in s["title"].lower()
                                and "conclusion" not in s["title"].lower()
                                and "reference" not in s["title"].lower())
                st.session_state.wp_title = title
                st.session_state.wp_sections = sections
                st.session_state.pop("wp_paper", None)
                st.rerun()

        if st.session_state.get("wp_sections"):
            st.markdown("#### 5️⃣ Review & edit the outline")
            wp_title = st.text_input("Paper title:", value=st.session_state.get("wp_title", ""), key="wp_title_input")

            top = st.columns([1, 1, 5])
            with top[0]:
                if st.button("➕ Add section"):
                    st.session_state.wp_sections.append({"title": "New Section", "description": "", "figure": True})
                    st.rerun()
            with top[1]:
                if st.button("↩️ Re-suggest outline"):
                    st.session_state.pop("wp_sections", None)
                    st.session_state.pop("wp_title", None)
                    st.rerun()

            delete_idx = None
            for i, sec in enumerate(st.session_state.wp_sections):
                with st.expander(f"Section {i+1}: {sec['title']}", expanded=False):
                    sec["title"] = st.text_input("Section title", value=sec["title"], key=f"wpsec_title_{i}")
                    sec["description"] = st.text_area("What should this section cover?",
                                                       value=sec["description"], height=80, key=f"wpsec_desc_{i}")
                    sec["figure"] = st.checkbox("Include a figure for this section",
                                               value=sec.get("figure", True), key=f"wpsec_fig_{i}")
                    cols = st.columns([1, 1, 4])
                    with cols[0]:
                        if i > 0 and st.button("⬆️ Move up", key=f"wpsec_up_{i}"):
                            st.session_state.wp_sections[i-1], st.session_state.wp_sections[i] = \
                                st.session_state.wp_sections[i], st.session_state.wp_sections[i-1]
                            st.rerun()
                    with cols[1]:
                        if st.button("🗑️ Delete", key=f"wpsec_del_{i}"):
                            delete_idx = i
            if delete_idx is not None:
                st.session_state.wp_sections.pop(delete_idx)
                st.rerun()

            if st.button("🖋️ Write Paper", type="primary", disabled=not st.session_state.wp_sections):
                log_event(_CURRENT_USER, "write_paper_generate")
                wp_title = st.session_state.get("wp_title_input", wp_title)
                add_memory_fact(_CURRENT_USER, f"Wrote a paper titled: {wp_title}")
                save_project_source_pdfs(WP_PDFS_DIR, wp_title, wp_sources)
                sections = st.session_state.wp_sections
                digest = build_source_digest(wp_sources, per_chars=wp_max_chars_per_paper)
                outline_str = " | ".join(s["title"] for s in sections)
                refs = numbered_references(wp_sources)
                generated_date = f"{datetime.now():%Y-%m-%d}"

                sections_data = []
                figs = []

                for sec in sections:
                    title, desc, want_fig = sec["title"], sec["description"], sec.get("figure", True)
                    is_abstract = "abstract" in title.lower()
                    st.markdown(f"### {title}")
                    user_msg = (
                        f"PAPER TITLE: {wp_title}\n"
                        f"AUTHOR'S IDEA: {wp_approved_idea}\n"
                        f"FULL OUTLINE (for context, do not rewrite other sections): {outline_str}\n\n"
                        f"NUMBERED SOURCES:\n{digest}\n\n"
                        f"Now write the section '{title}'. {desc} Aim for roughly 400-600 words. "
                        f"Cite sources as [n] where relevant."
                    )
                    try:
                        txt = st.write_stream(stream_chat(
                            [{"role": "system", "content": PAPER_SECTION_SYSTEM},
                             {"role": "user", "content": user_msg}],
                            global_model_string, api_base=selected_api_base, temperature=0.5,
                            num_ctx=num_ctx, api_key=selected_api_key))
                    except Exception as e:
                        txt = ""
                        st.error(f"Section '{title}' failed: {e}")

                    if is_abstract:
                        txt = re.sub(r"\s?\[\d+(?:\s*,\s*\d+)*\]", "", txt)

                    entry = {"title": title, "text": txt, "figure_idx": None, "figure_caption": ""}

                    if want_fig and txt.strip():
                        with st.status(f"🎨 Drawing a figure for “{title}”…", expanded=False) as figbox:
                            def _figstatus(msg):
                                figbox.write(msg)
                            fig_result = generate_section_figure(
                                wp_title, title, txt, global_model_string,
                                is_graphical_abstract=is_abstract, api_base=selected_api_base,
                                api_key=selected_api_key, num_ctx=num_ctx, status_cb=_figstatus)
                            figbox.update(label=f"🎨 Figure for “{title}”: {'done' if fig_result else 'skipped'}",
                                         state="complete")
                        if fig_result:
                            idx = len(figs)
                            figs.append(fig_result)
                            entry["figure_idx"] = idx
                            entry["figure_caption"] = fig_result["caption"]
                            st.image(fig_result["png"], caption=f"Figure {idx + 1}. {fig_result['caption']}")

                    sections_data.append(entry)

                md_master = rebuild_paper_markdown(wp_title, sections_data, refs, generated_date, len(wp_sources))

                st.session_state.wp_paper = {
                    "title": wp_title, "sections_data": sections_data, "refs": refs,
                    "generated_date": generated_date, "source_count": len(wp_sources), "figs": figs,
                    "md_master": md_master, "md": materialize_markdown_images(md_master, figs),
                    "bib": to_bibtex(wp_sources), "ts": datetime.now().strftime("%Y%m%d_%H%M%S")}
                st.success("🎉 Paper draft complete.")

    # ---- Step 6-8: Edit, AI fine-tune, export ---------------------------------
    wp_paper = st.session_state.get("wp_paper")
    if wp_paper:
        st.divider()
        st.markdown("#### 6️⃣ Edit, fine-tune, and export")

        def _apply_pending(pending_key, real_key):
            """Streamlit forbids writing to st.session_state[real_key] once that widget has
            already been instantiated THIS run — doing so raises StreamlitAPIException (the
            bug this fixes). To update a widget's displayed value programmatically instead:
            stash the new value under `pending_key`, call st.rerun(), and this function — called
            at the TOP of the next run, before the widget is created — applies it safely."""
            if pending_key in st.session_state:
                st.session_state[real_key] = st.session_state.pop(pending_key)

        def _wp_resync(paper):
            """Rebuild md_master/md from sections_data after any per-section change, and queue
            the result for the whole-document editor too — otherwise that box would keep
            showing its old content for the same reason the per-section box did."""
            paper["md_master"] = rebuild_paper_markdown(
                paper["title"], paper["sections_data"], paper["refs"],
                paper["generated_date"], paper["source_count"])
            paper["md"] = materialize_markdown_images(paper["md_master"], paper["figs"])
            st.session_state.wp_paper = paper
            st.session_state["_pending_wp_edit_area"] = paper["md_master"]

        if wp_paper.get("sections_data"):
            st.markdown("##### ✏️ Edit sections individually")
            st.caption("Edit text directly, or give the model an instruction to revise just this "
                      "section — much faster than rewriting the whole paper for a small change.")
            for i, sec in enumerate(wp_paper["sections_data"]):
                with st.expander(f"{sec['title']}", expanded=False):
                    _apply_pending(f"_pending_wpsecedit_text_{i}", f"wpsecedit_text_{i}")
                    new_text = st.text_area("Section text:", value=sec["text"], height=180,
                                           key=f"wpsecedit_text_{i}")
                    if st.button("💾 Save this section", key=f"wpsecedit_save_{i}"):
                        sec["text"] = new_text
                        _wp_resync(wp_paper)
                        st.success(f"Saved '{sec['title']}'.")
                        st.rerun()

                    fi1, fi2 = st.columns([3, 1])
                    with fi1:
                        instr = st.text_input(
                            "Fine-tune instruction:", key=f"wpsecedit_instr_{i}",
                            placeholder="e.g. 'make this more concise', 'add more technical detail'")
                    with fi2:
                        st.markdown("<div style='height:28px'></div>", unsafe_allow_html=True)
                        finetune_clicked = st.button("🪄 Fine-tune", key=f"wpsecedit_finetune_{i}",
                                                    disabled=not instr.strip())
                    if finetune_clicked:
                        with st.spinner(f"Revising '{sec['title']}'…"):
                            revised = finetune_section(
                                new_text, instr, global_model_string, api_base=selected_api_base,
                                api_key=selected_api_key, num_ctx=num_ctx)
                        if revised:
                            if "abstract" in sec["title"].lower():
                                revised = re.sub(r"\s?\[\d+(?:\s*,\s*\d+)*\]", "", revised)
                            sec["text"] = revised
                            # Queue the box's new value for next run rather than writing to its
                            # key directly here — it was already instantiated a few lines above
                            # in this same run, so a direct write would raise the exact
                            # exception this pending-value pattern exists to avoid.
                            st.session_state[f"_pending_wpsecedit_text_{i}"] = revised
                            _wp_resync(wp_paper)
                            st.success(f"Revised '{sec['title']}'.")
                            st.rerun()
                        else:
                            st.error("Fine-tuning failed — try again or rephrase the instruction.")

                    if sec.get("figure_idx") is not None and sec["figure_idx"] < len(wp_paper["figs"]):
                        fidx = sec["figure_idx"]
                        st.image(wp_paper["figs"][fidx]["png"], caption=f"Figure {fidx + 1}. {sec['figure_caption']}")
                        if st.button("🔄 Regenerate this figure", key=f"wpsecedit_regenfig_{i}"):
                            with st.spinner("Redrawing figure…"):
                                new_fig = generate_section_figure(
                                    wp_paper["title"], sec["title"], sec["text"], global_model_string,
                                    is_graphical_abstract=("abstract" in sec["title"].lower()),
                                    api_base=selected_api_base, api_key=selected_api_key, num_ctx=num_ctx)
                            if new_fig:
                                wp_paper["figs"][fidx] = new_fig
                                sec["figure_caption"] = new_fig["caption"]
                                _wp_resync(wp_paper)
                                st.success("Figure regenerated.")
                                st.rerun()
                            else:
                                st.warning("Figure regeneration failed — keeping the existing one.")

        st.markdown("##### 📄 Edit whole document")
        st.caption("Editing here is independent from the per-section editor above — using a "
                  "per-section action afterward will rebuild the document from its section "
                  "texts, which can overwrite manual edits made only here.")
        _apply_pending("_pending_wp_edit_area", "wp_edit_area")
        edited = st.text_area("Edit the paper directly (Markdown — [[FIGURE:n]] markers become "
                              "images in the exports):", value=wp_paper["md_master"], height=400,
                              key="wp_edit_area")
        if st.button("💾 Save manual edits", key="wp_save_manual_edits"):
            wp_paper["md_master"] = edited
            wp_paper["md"] = materialize_markdown_images(edited, wp_paper["figs"])
            st.session_state.wp_paper = wp_paper
            st.success("Saved.")
            st.rerun()

        st.markdown("##### 🪄 Fine-tune entire paper with AI")
        wp_instruction = st.text_area("Describe what to change:", height=80, key="wp_finetune_instruction",
                                      placeholder="e.g. 'Make the Related Work section more concise', "
                                                  "'expand the Results section with more detail', "
                                                  "'tighten the Introduction and fix any repetition'")
        if st.button("🪄 Fine-tune with AI", type="primary", disabled=not wp_instruction.strip()):
            with st.spinner("Revising the paper…"):
                revised = finetune_paper(wp_paper["md_master"], wp_instruction, global_model_string,
                                        api_base=selected_api_base, api_key=selected_api_key,
                                        num_ctx=num_ctx, max_tokens=crew_max_tokens)
            if revised:
                wp_paper["md_master"] = revised
                wp_paper["md"] = materialize_markdown_images(revised, wp_paper["figs"])
                st.session_state.wp_paper = wp_paper
                # Queued for next run — "wp_edit_area" was already instantiated above, in this
                # same run, so writing to it directly here would raise the same exception.
                st.session_state["_pending_wp_edit_area"] = revised
                st.success("Revised — see the updated draft below and in the text box above.")
                st.rerun()
            else:
                st.error("Fine-tuning failed — the model didn't return usable output. Try again "
                        "or rephrase the instruction.")

        st.markdown("##### 📄 Current draft")
        st.markdown(wp_paper["md"])

        d1, d2, d3, d4, d5 = st.columns(5)
        with d1:
            st.download_button("📄 Paper (.md)", data=wp_paper["md"], file_name=f"paper_{wp_paper['ts']}.md")
        with d2:
            st.download_button("📄 Paper (.txt)", data=wp_paper["md"], file_name=f"paper_{wp_paper['ts']}.txt")
        with d3:
            pdf_bytes = md_to_pdf_bytes(wp_paper["md_master"], title=wp_paper["title"], figs=wp_paper["figs"])
            if pdf_bytes:
                st.download_button("📕 Paper (.pdf)", data=pdf_bytes, file_name=f"paper_{wp_paper['ts']}.pdf",
                                   mime="application/pdf")
            else:
                st.caption("PDF needs: pip install reportlab")
        with d4:
            st.download_button("📚 References (.bib)", data=wp_paper["bib"], file_name=f"paper_refs_{wp_paper['ts']}.bib")
        with d5:
            try:
                fig_filenames = {i: f"fig{i}.png" for i in range(len(wp_paper["figs"]))}
                tex_str = survey_to_latex(wp_paper["md_master"], wp_paper["title"], wp_paper["figs"], fig_filenames)
                tex_zip = build_latex_zip(tex_str, wp_paper["figs"], fig_filenames, bib_str=wp_paper["bib"],
                                          tex_filename="paper.tex")
                st.download_button("📦 LaTeX (.zip)", data=tex_zip, file_name=f"paper_{wp_paper['ts']}_latex.zip",
                                   mime="application/zip",
                                   help="Contains paper.tex, figure PNGs, and references.bib.")
            except Exception as e:
                st.caption(f"LaTeX export unavailable: {e}")

# =============================== SUMMARIZE TAB =============================
with tab_summarize:
    st.subheader("📖 Summarize Paper")
    st.caption("Upload one paper. The model reads it and writes a clear summary — then you can ask follow-up questions and it answers using the paper.")

    sm_up = st.file_uploader("Upload a paper (PDF)", type=["pdf"], key="sm_pdf")

    scol1, scol2 = st.columns([1.4, 1])
    with scol1:
        st.info(f"Active Model Environment: **{global_model_string}**")
    with scol2:
        sm_maxchars = st.select_slider("Max chars read", key="sm_chars",
                                       options=[8000, 16000, 24000, 40000, 60000, 100000, 200000], value=40000)
    sm_speak = st.toggle("🔊 Read the summary aloud", value=False, key="sm_speak")

    if sm_up is not None:
        fname = getattr(sm_up, "name", "paper.pdf")
        if st.session_state.get("sm_file_name") != fname or st.session_state.get("sm_maxchars_used") != sm_maxchars:
            try:
                text, n_pages, truncated = extract_pdf_text(sm_up, max_chars=sm_maxchars)
            except Exception as e:
                text, n_pages, truncated = "", 0, False
                st.error(f"Could not read PDF: {e}")
            st.session_state["sm_paper_text"] = text
            st.session_state["sm_file_name"] = fname
            st.session_state["sm_maxchars_used"] = sm_maxchars
            st.session_state["sm_pages"] = n_pages
            st.session_state["sm_truncated"] = truncated
            st.session_state.pop("sm_summary", None)
            st.session_state["sm_chat"] = []

    paper_text = st.session_state.get("sm_paper_text", "")
    if paper_text:
        n_pages = st.session_state.get("sm_pages", 0)
        truncated = st.session_state.get("sm_truncated", False)
        tok = int(len(paper_text) / 4) + 400
        note = "  ⚠️ may exceed context window" if tok > num_ctx * 0.95 else ""
        msg = f"Read **{n_pages}** page(s), **{len(paper_text):,}** characters (≈ {tok:,} tokens){note}."
        if truncated:
            msg += "  Truncated — raise *Max chars read* / context if needed."
        st.info(msg)

        if st.button("📖 Summarize", type="primary", disabled=len(paper_text.strip()) == 0):
            log_event(_CURRENT_USER, "summarize_paper")
            add_memory_fact(_CURRENT_USER, f"Read/summarized a paper: {st.session_state.get('sm_file_name', 'a paper')}")
            with st.spinner("Reading and summarizing…"):
                try:
                    summary = st.write_stream(stream_chat(
                        [{"role": "system", "content": SUMMARIZE_SYSTEM},
                         {"role": "user", "content": f"Paper text:\n\n{paper_text}"}],
                        global_model_string, api_base=selected_api_base, temperature=0.3,
                        num_ctx=num_ctx, api_key=selected_api_key))
                except Exception as e:
                    summary = ""
                    st.error(f"Summary failed: {e}")
            if summary:
                st.session_state["sm_summary"] = summary
                if sm_speak:
                    with st.spinner("🔊 Generating speech…"):
                        audio, terr = tts_to_mp3(summary)
                    if audio:
                        st.audio(audio, format="audio/mpeg", autoplay=True)

        if st.session_state.get("sm_summary"):
            st.markdown("#### 📋 Summary")
            st.markdown(st.session_state["sm_summary"])
            base = re.sub(r"\.pdf$", "", st.session_state.get("sm_file_name", "paper"), flags=re.I)
            sm1, sm2 = st.columns(2)
            with sm1:
                st.download_button("💾 Download summary (.md)", data=st.session_state["sm_summary"], file_name=f"summary_{base}.md")
            with sm2:
                sm_pdf = md_to_pdf_bytes(st.session_state["sm_summary"], title=f"Summary — {base}")
                if sm_pdf:
                    st.download_button("📕 Download summary (.pdf)", data=sm_pdf, file_name=f"summary_{base}.pdf",
                                       mime="application/pdf")
                else:
                    st.caption("PDF needs: pip install reportlab")

            st.divider()
            st.markdown("#### 💬 Ask about this paper")
            if "sm_chat" not in st.session_state:
                st.session_state["sm_chat"] = []
            for m in st.session_state["sm_chat"]:
                with st.chat_message(m["role"]):
                    st.markdown(m["content"])
            q = st.chat_input("Ask a question about the paper…", key="sm_q")
            if q:
                st.session_state["sm_chat"].append({"role": "user", "content": q})
                with st.chat_message("user"):
                    st.markdown(q)
                sys_msg = SUMMARIZE_QA_SYSTEM.replace("{paper}", paper_text)
                msgs = [{"role": "system", "content": sys_msg}] + st.session_state["sm_chat"]
                with st.chat_message("assistant"):
                    try:
                        ans = st.write_stream(stream_chat(
                            msgs, global_model_string, api_base=selected_api_base,
                            temperature=0.3, num_ctx=num_ctx, api_key=selected_api_key))
                    except Exception as e:
                        ans = f"⚠️ Error: {e}"; st.error(ans)
                    if sm_speak and ans and not str(ans).startswith("⚠️"):
                        audio, terr = tts_to_mp3(ans)
                        if audio:
                            st.audio(audio, format="audio/mpeg", autoplay=True)
                st.session_state["sm_chat"].append({"role": "assistant", "content": ans})

# =============================== REVIEW TAB =================================
with tab_review:
    st.subheader("📝 Peer-Review a Paper")
    st.caption("Upload a PDF. The model reads it and writes a structured review report.")

    up = st.file_uploader("Upload a paper (PDF)", type=["pdf"], key="review_pdf")

    rcol1, rcol2 = st.columns([1.3, 1])
    with rcol1:
        st.info(f"Active Model Environment: **{global_model_string}**")
    with rcol2:
        rev_maxchars = st.select_slider("Max chars read", key="rev_chars",
                                        options=[8000, 16000, 24000, 40000, 60000, 100000, 200000, 400000], value=60000)

    rev_speak = st.toggle("🔊 Also read the review aloud", value=False, key="rev_speak")

    if up is not None:
        try:
            paper_text, n_pages, truncated = extract_pdf_text(up, max_chars=rev_maxchars)
        except Exception as e:
            paper_text, n_pages, truncated = "", 0, False
            st.error(f"Could not read PDF: {e}")

        tok = int(len(paper_text) / 4) + 400
        note = "  ⚠️ may exceed context window" if tok > num_ctx * 0.95 else ""
        info = f"Read **{n_pages}** page(s), extracted **{len(paper_text):,}** characters (≈ {tok:,} tokens){note}."
        if truncated:
            info += "  The paper was truncated to fit — raise *Max chars read* / context if needed."
        st.info(info)

        disabled = len(paper_text.strip()) == 0
        if st.button("🧑‍⚖️ Generate review", type="primary", disabled=disabled):
            log_event(_CURRENT_USER, "review_paper")
            add_memory_fact(_CURRENT_USER, f"Peer-reviewed a paper: {up.name}")
            messages = [
                {"role": "system", "content": REVIEWER_SYSTEM},
                {"role": "user", "content": f"Here is the extracted text of the paper to review:\n\n{paper_text}"},
            ]
            with st.spinner("Reading the paper and writing the review…"):
                try:
                    review = st.write_stream(
                        stream_chat(messages, global_model_string,
                                    api_base=selected_api_base, temperature=0.3,
                                    num_ctx=num_ctx, api_key=selected_api_key))
                except Exception as e:
                    review = ""
                    st.error(f"Review failed: {e}")
            if review:
                st.session_state["last_review"] = review
                st.session_state["last_review_name"] = getattr(up, "name", "paper")
                if rev_speak:
                    with st.spinner("🔊 Generating speech…"):
                        audio, terr = tts_to_mp3(review)
                    if terr: st.caption(f"🔇 {terr}")
                    elif audio: st.audio(audio, format="audio/mpeg", autoplay=True)

    if st.session_state.get("last_review"):
        rv = st.session_state["last_review"]
        base = re.sub(r"\.pdf$", "", st.session_state.get("last_review_name", "paper"), flags=re.I)
        rev_ts = datetime.now().strftime("%Y%m%d_%H%M%S")
        rv1, rv2 = st.columns(2)
        with rv1:
            st.download_button("💾 Download review (.md)", data=rv, file_name=f"review_{base}_{rev_ts}.md")
        with rv2:
            rv_pdf = md_to_pdf_bytes(rv, title=f"Peer Review — {base}")
            if rv_pdf:
                st.download_button("📕 Download review (.pdf)", data=rv_pdf, file_name=f"review_{base}_{rev_ts}.pdf",
                                   mime="application/pdf")
            else:
                st.caption("PDF needs: pip install reportlab")

# =============================== CODING TAB ==================================
with tab_coding:
    st.subheader("💻 Coding Assistant")

    # ---- Project Manager --------------------------------------------------
    # A "project" is just a named subfolder of this user's sandbox. It gives every
    # coding mode (Quick Edit, Agent, Claude Code) a shared place to work, browse,
    # and download from — instead of each mode juggling its own raw folder path.
    def list_sandbox_projects():
        try:
            return sorted(d for d in os.listdir(USER_SANDBOX_DIR)
                          if os.path.isdir(os.path.join(USER_SANDBOX_DIR, d)))
        except Exception:
            return []

    st.markdown("##### 📁 Project")
    _projects = list_sandbox_projects()
    prj1, prj2 = st.columns([3, 1.4])
    with prj1:
        if _projects:
            _cur = st.session_state.get("coding_project")
            _idx = _projects.index(_cur) if _cur in _projects else 0
            active_project = st.selectbox("Active project", _projects, index=_idx, key="coding_project_picker")
        else:
            active_project = None
            st.caption("No projects yet — create one to give the coding tools a place to work.")
    with prj2:
        with st.expander("➕ New project"):
            new_proj_name = st.text_input("Name", key="new_project_name", label_visibility="collapsed",
                                          placeholder="e.g. my-first-project")
            if st.button("Create", key="create_project_btn", use_container_width=True):
                safe = re.sub(r"[^A-Za-z0-9_ -]+", "_", new_proj_name.strip())[:60] \
                    or f"project_{datetime.now().strftime('%Y%m%d_%H%M%S')}"
                os.makedirs(os.path.join(USER_SANDBOX_DIR, safe), exist_ok=True)
                st.session_state["coding_project"] = safe
                st.rerun()

    project_path = None
    if active_project:
        st.session_state["coding_project"] = active_project
        project_path = os.path.join(USER_SANDBOX_DIR, active_project)
        os.makedirs(project_path, exist_ok=True)

        # Sync Agent/Claude Code's folder fields to the active project whenever it
        # changes — but only then, so it doesn't fight the user editing them by hand.
        if st.session_state.get("_coding_project_synced") != active_project:
            st.session_state["agent_project_dir"] = project_path
            st.session_state["cc_project_dir"] = project_path
            st.session_state["_coding_project_synced"] = active_project

        with st.expander(f"📂 Files in '{active_project}'", expanded=False):
            _all_files = []
            for _root, _dirs, _files in os.walk(project_path):
                for _fn in _files:
                    _full = os.path.join(_root, _fn)
                    _rel = os.path.relpath(_full, project_path)
                    _all_files.append((_rel, _full))
            if not _all_files:
                st.caption("Empty — files created by Agent or Claude Code here will show up.")
            else:
                for _rel, _full in sorted(_all_files):
                    fc1, fc2 = st.columns([5, 1])
                    with fc1:
                        try:
                            _size = os.path.getsize(_full)
                        except OSError:
                            _size = 0
                        st.caption(f"📄 {_rel}  ({_size:,} bytes)")
                    with fc2:
                        try:
                            with open(_full, "rb") as _fh:
                                st.download_button("⬇️", data=_fh.read(),
                                                  file_name=os.path.basename(_rel),
                                                  key=f"dl_{_rel}", use_container_width=True)
                        except Exception:
                            st.caption("—")

                _zip_buf = io.BytesIO()
                with zipfile.ZipFile(_zip_buf, "w", zipfile.ZIP_DEFLATED) as _zf:
                    for _rel, _full in _all_files:
                        _zf.write(_full, _rel)
                st.download_button(f"⬇️ Download entire project ({active_project}.zip)",
                                  data=_zip_buf.getvalue(), file_name=f"{active_project}.zip",
                                  key="dl_project_zip", use_container_width=True)

    st.divider()

    coding_mode = st.radio(
        "Mode", ["✏️ Quick Edit", "🤖 Agent (local files + terminal)",
                 "🧭 Claude Code (native CLI, any LLM)", "🖥️ CDesktop"],
        index=2, horizontal=True, key="coding_mode",
        help="Quick Edit: paste/upload one file and iterate on it in the browser. "
             "Agent: this app's own built-in agent — reads/writes files and runs shell "
             "commands, with your approval on every write and command, using whichever "
             "model is configured in the sidebar. "
             "Claude Code: runs the official open-source `claude` CLI as a real subprocess "
             "in a project folder, pointed at Ollama, an OpenAI-compatible proxy, or "
             "Anthropic's API — so it isn't tied to any one provider. "
             "CDesktop: launches the CDesktop desktop agent via npx cdesktop — opens a "
             "code-aware chat in a new browser tab.")

    if coding_mode == "✏️ Quick Edit":
        st.caption("Upload, paste, or start from scratch. Tell the model what to change, and it will give you the complete updated code.")

        if "coding_code" not in st.session_state:
            st.session_state["coding_code"] = ""
        if "coding_lang" not in st.session_state:
            st.session_state["coding_lang"] = "python"
        if "coding_chat" not in st.session_state:
            st.session_state["coding_chat"] = []

        c_col1, c_col2 = st.columns([1, 1])
        with c_col1:
            code_up = st.file_uploader("Upload a code file", key="code_up")
            if code_up is not None:
                fname = getattr(code_up, "name", "")
                ext = fname.split(".")[-1].lower() if "." in fname else ""
                content = code_up.getvalue().decode("utf-8", errors="replace")
                
                if st.session_state.get("last_code_up") != fname:
                    st.session_state["coding_code"] = content
                    if ext in _EXT_TO_LANG:
                        st.session_state["coding_lang"] = _EXT_TO_LANG[ext]
                    st.session_state["last_code_up"] = fname
                    st.session_state["coding_chat"] = [] 
                    st.rerun()

            st.session_state["coding_lang"] = st.selectbox(
                "Language", 
                list(CODE_LANGS.keys()), 
                index=list(CODE_LANGS.keys()).index(st.session_state["coding_lang"]) if st.session_state["coding_lang"] in CODE_LANGS else 0
            )

        with c_col2:
            st.info(f"Active Environment Target: **{global_model_string}**")
            st.markdown("<br/>", unsafe_allow_html=True)
            if st.button("🧹 Clear & Start Over", key="code_clear"):
                st.session_state["coding_code"] = ""
                st.session_state["coding_chat"] = []
                st.session_state.pop("last_code_up", None)
                st.rerun()

        st.markdown("### Current Code")
        current_code = st.text_area("Edit code here manually, or leave blank to start from scratch:", 
                                    value=st.session_state["coding_code"], 
                                    height=300, key="code_editor")
        if current_code != st.session_state["coding_code"]:
             st.session_state["coding_code"] = current_code

        if st.session_state["coding_code"].strip():
            ext = CODE_LANGS.get(st.session_state["coding_lang"], "txt")
            st.download_button("💾 Download Final Code", 
                               data=st.session_state["coding_code"], 
                               file_name=f"script_{datetime.now():%Y%m%d_%H%M%S}.{ext}")

        st.markdown("### Instructions")
        
        for m in st.session_state["coding_chat"]:
            with st.chat_message(m["role"]):
                st.markdown(m["content"])
                
        code_prompt = st.chat_input("E.g., 'Add error handling', 'Refactor into a class', 'Write a function to fetch data'")
        
        if code_prompt:
            log_event(_CURRENT_USER, "coding_quick_edit")
            add_memory_fact(_CURRENT_USER, f"Uses {st.session_state.get('coding_lang', 'code')} for coding tasks")
            st.session_state["coding_chat"].append({"role": "user", "content": code_prompt})
            
            with st.chat_message("user"):
                st.markdown(code_prompt)
            
            user_msg = f"CURRENT CODE:\n```{st.session_state['coding_lang']}\n{st.session_state['coding_code']}\n```\n\nREQUEST: {code_prompt}"
            
            msgs = [{"role": "system", "content": CODING_SYSTEM}]
            for m in st.session_state["coding_chat"][:-1]:
                msgs.append(m)
            msgs.append({"role": "user", "content": user_msg})

            with st.chat_message("assistant"):
                with st.spinner("Processing structural adjustments..."):
                    try:
                        reply = st.write_stream(stream_chat(
                            msgs, 
                            global_model_string, 
                            api_base=selected_api_base, 
                            temperature=0.2, 
                            num_ctx=num_ctx, 
                            api_key=selected_api_key,
                            max_tokens=crew_max_tokens,
                        ))

                        fence_count = reply.count("```")
                        truncated = "```" in reply and fence_count % 2 == 1
                        new_code = extract_code_block(reply)
                        st.session_state["coding_chat"].append({"role": "assistant", "content": reply})

                        if new_code and new_code.strip() and new_code.strip() != reply.strip():
                            st.session_state["coding_code"] = new_code
                            if truncated:
                                st.warning("⚠️ Response looked cut off before the closing ``` — "
                                           "applied what was generated so far, but increase "
                                           "'Max output tokens/step' in the sidebar and try again "
                                           "to get the complete file.")
                            else:
                                st.success("Code engine update applied!")
                    except Exception as e:
                        st.error(f"Failed: {e}")
            st.rerun()

    elif coding_mode == "🤖 Agent (local files + terminal)":
        st.caption("Works on a real folder on this machine using the model configured in the "
                  "sidebar — any provider. Every file write and every shell command is shown "
                  "to you and requires approval before it runs, unless you turn that off below.")

        if "coding_agent" not in st.session_state:
            st.session_state.coding_agent = {
                "project_dir": "", "history": [], "transcript": [],
                "status": "idle", "pending_action": None, "steps": 0, "task_gen": 0,
            }
        agent = st.session_state.coding_agent

        # ---- Sessions: save/load/delete, so a coding session survives closing the app ----
        with st.expander("📁 Sessions", expanded=False):
            coding_sessions = list_coding_sessions()
            se1, se2, se3 = st.columns([3, 1, 1])
            with se1:
                picked_session = st.selectbox(
                    "Saved sessions", ["(none)"] + coding_sessions, key="agent_session_picker",
                    label_visibility="collapsed")
            with se2:
                if st.button("📂 Load", disabled=(picked_session == "(none)"), use_container_width=True, key="agent_sess_load"):
                    payload, err = load_coding_session(picked_session)
                    if err:
                        st.error(f"Could not load: {err}")
                    else:
                        # Safe here: this expander renders before agent_project_dir/agent_task
                        # are instantiated later in this same script pass.
                        st.session_state.coding_agent = {
                            "project_dir": payload.get("project_dir", ""),
                            "history": payload.get("history", []),
                            "transcript": payload.get("transcript", []),
                            "status": "done" if payload.get("history") else "idle",
                            "pending_action": None, "steps": 0,
                            "task_gen": agent.get("task_gen", 0) + 1,
                        }
                        st.session_state["agent_project_dir"] = payload.get("project_dir", "")
                        st.session_state["agent_active_session"] = picked_session
                        st.rerun()
            with se3:
                if st.button("🗑️ Delete", disabled=(picked_session == "(none)"), use_container_width=True, key="agent_sess_delete"):
                    delete_coding_session(picked_session)
                    if st.session_state.get("agent_active_session") == picked_session:
                        st.session_state.pop("agent_active_session", None)
                    st.rerun()

            sv1, sv2 = st.columns([3, 1])
            with sv1:
                session_save_name = st.text_input(
                    "Session name", value=st.session_state.get("agent_active_session", ""),
                    key="agent_session_name", label_visibility="collapsed",
                    placeholder="Auto-named from your first instruction — edit here to rename")
            with sv2:
                if st.button("✏️ Rename", disabled=not session_save_name.strip(), use_container_width=True, key="agent_sess_rename"):
                    old_name = st.session_state.get("agent_active_session")
                    new_name = session_save_name.strip()
                    ok, err = save_coding_session(new_name, agent)
                    if ok:
                        if old_name and old_name != new_name:
                            delete_coding_session(old_name)
                        st.session_state["agent_active_session"] = new_name
                        st.success(f"Renamed to '{new_name}'.")
                    else:
                        st.error(f"Could not rename: {err}")

            if st.session_state.get("agent_active_session"):
                st.caption(f"📌 Working on: **{st.session_state['agent_active_session']}** — saves automatically as you work.")

        def _agent_browse_callback():
            # Runs BEFORE the script reruns and re-instantiates the text_input widget below,
            # so it's safe to write to its session_state key here — doing it inline after the
            # widget has already rendered this run raises a StreamlitAPIException.
            picked = browse_for_folder(initial_dir=st.session_state.get("agent_project_dir", "") or USER_SANDBOX_DIR)
            if not picked:
                st.session_state["agent_browse_result"] = "cancelled"
            elif not (_CURRENT_IS_ADMIN or path_in_sandbox(picked, USER_SANDBOX_DIR)):
                st.session_state["agent_browse_result"] = "outside_sandbox"
            else:
                st.session_state["agent_browse_result"] = "picked"
                st.session_state["agent_project_dir"] = picked

        if _CURRENT_IS_ADMIN:
            st.caption("🛡️ Admin account — no folder restriction, any path on this machine is allowed.")
        else:
            st.caption(f"📁 Folders must be inside your sandbox: `{USER_SANDBOX_DIR}`")
        ac1, ac2, ac3 = st.columns([3, 0.6, 1])
        with ac1:
            project_dir_input = st.text_input(
                "Project folder" if _CURRENT_IS_ADMIN else "Project folder (must be inside your sandbox — see above)",
                value=agent["project_dir"], key="agent_project_dir", placeholder=USER_SANDBOX_DIR,
                disabled=agent["status"] not in ("idle", "done", "stopped", "error"))
        with ac2:
            st.markdown("<div style='height:28px'></div>", unsafe_allow_html=True)  # align with text input
            st.button("📂 Browse", key="agent_browse_dir", use_container_width=True,
                      disabled=agent["status"] not in ("idle", "done", "stopped", "error"),
                      help="Opens a folder picker.",
                      on_click=_agent_browse_callback)
        with ac3:
            agent_auto_approve = st.toggle("Auto-approve", value=False, key="agent_auto_approve",
                                          help="Skip the approval step and execute writes/commands "
                                               "immediately. Off by default — review each one.")

        _browse_result = st.session_state.pop("agent_browse_result", None)
        if _browse_result == "cancelled":
            st.toast("No folder chosen (or the folder picker isn't available here — "
                    "type the path manually instead).")
        elif _browse_result == "outside_sandbox":
            st.toast(f"⛔ That folder is outside your sandbox ({USER_SANDBOX_DIR}) — not allowed.")


        agent_python = st.text_input(
            "Python interpreter for the agent's commands", value=sys.executable,
            key="agent_python",
            help="Machines often have several separate Python installs/venvs, each with its "
                 "own packages — bare 'python'/'pip' can silently resolve to the wrong one. "
                 "Told to the agent as guidance; it decides when to use it. Defaults to the "
                 "interpreter running this app, which is guaranteed to already have its "
                 "dependencies installed.")
        agent_max_steps = st.slider("Max steps per instruction", 5, 40, 20, key="agent_max_steps")

        dir_exists = bool(project_dir_input) and os.path.isdir(project_dir_input)
        dir_sandboxed = _CURRENT_IS_ADMIN or (bool(project_dir_input) and path_in_sandbox(project_dir_input, USER_SANDBOX_DIR))
        dir_ready = dir_exists and dir_sandboxed
        if project_dir_input and not dir_sandboxed:
            st.error(f"⛔ That folder is outside your sandbox. Use a path under: `{USER_SANDBOX_DIR}`")
        elif project_dir_input and not dir_exists:
            st.warning("That folder doesn't exist yet — create it first (e.g. via the Browse "
                      "dialog's 'New Folder' option, or a file explorer)"
                      + ("." if _CURRENT_IS_ADMIN else ", inside your sandbox."))

        # ---- Render the transcript so far ------------------------------------
        for entry in agent["transcript"]:
            t = entry["type"]
            if t == "task":
                st.info(f"🎯 **Task:** {entry['text']}")
            elif t == "action":
                st.markdown(f"**Step {entry.get('step','')} — `{entry['action']}`**")
                st.code(entry["detail"], language="text")
            elif t == "observation":
                with st.expander(f"↳ result", expanded=False):
                    st.code(entry["text"][:4000], language="text")
            elif t == "rejected":
                st.warning(f"↳ You rejected this action.")
            elif t == "error":
                st.error(entry["text"])
            elif t == "done":
                st.success(f"✅ **Done:** {entry['text']}")

        # ---- Drive the loop ---------------------------------------------------
        if agent["status"] == "running":
            if agent["steps"] >= agent_max_steps:
                agent["status"] = "stopped"
                agent["transcript"].append({"type": "error", "text": f"Stopped — hit the {agent_max_steps}-step limit for this instruction. You can give it another instruction to keep going."})
                st.rerun()

            with st.spinner("Thinking about the next step…"):
                action, raw = agent_next_action(
                    agent["history"], global_model_string, api_base=selected_api_base,
                    api_key=selected_api_key, num_ctx=num_ctx, max_tokens=crew_max_tokens)

            if action is None:
                agent["status"] = "error"
                agent["transcript"].append({"type": "error", "text": f"Couldn't parse the model's action: {raw[:300]}"})
                st.rerun()

            agent["history"].append({"role": "assistant", "content": raw})
            act_type = action.get("action")
            agent["steps"] += 1

            if act_type == "done":
                agent["status"] = "done"
                agent["transcript"].append({"type": "done", "text": action.get("summary", "Task complete.")})
                st.rerun()

            elif act_type in ("list_dir", "read_file"):
                # Read-only — safe to execute immediately, no approval needed.
                path = action.get("path", ".")
                agent["transcript"].append({"type": "action", "step": agent["steps"], "action": act_type, "detail": path})
                if act_type == "list_dir":
                    result = agent_list_dir(agent["project_dir"], path)
                else:
                    result = agent_read_file(agent["project_dir"], path)
                agent["transcript"].append({"type": "observation", "text": result})
                agent["history"].append({"role": "user", "content": f"OBSERVATION ({act_type} {path}):\n{result}"})
                st.rerun()

            elif act_type in ("write_file", "run_command"):
                if agent_auto_approve:
                    agent["pending_action"] = action
                    agent["status"] = "executing_approved"
                else:
                    agent["pending_action"] = action
                    agent["status"] = "awaiting_approval"
                st.rerun()

            else:
                agent["status"] = "error"
                agent["transcript"].append({"type": "error", "text": f"Unknown action type: {act_type!r}"})
                st.rerun()

        # ---- Approval gate for write_file / run_command -----------------------
        if agent["status"] in ("awaiting_approval", "executing_approved"):
            pending = agent["pending_action"]
            act_type = pending.get("action")

            if act_type == "write_file":
                detail = f"WRITE FILE: {pending.get('path','')}\n\n{pending.get('content','')}"
                st.markdown(f"##### ✍️ Wants to write `{pending.get('path','')}`")
                st.code(pending.get("content", ""), language="text", line_numbers=True)
            else:
                detail = f"RUN COMMAND: {pending.get('command','')}"
                st.markdown("##### ⚡ Wants to run a command")
                st.code(pending.get("command", ""), language="bash")

            if agent["status"] == "executing_approved":
                approve, reject = True, False
            else:
                bc1, bc2 = st.columns(2)
                with bc1:
                    approve = st.button("✅ Approve & run", type="primary", key=f"agent_approve_{agent['steps']}")
                with bc2:
                    reject = st.button("❌ Reject", key=f"agent_reject_{agent['steps']}")

            if approve:
                agent["transcript"].append({"type": "action", "step": agent["steps"], "action": act_type, "detail": detail})
                if act_type == "write_file":
                    ok, err = agent_write_file(agent["project_dir"], pending.get("path", ""), pending.get("content", ""))
                    result = "OK — file written." if ok else f"ERROR: {err}"
                else:
                    out, errout, code = agent_run_command(agent["project_dir"], pending.get("command", ""))
                    result = f"exit code: {code}\nSTDOUT:\n{out}\nSTDERR:\n{errout}"
                agent["transcript"].append({"type": "observation", "text": result})
                agent["history"].append({"role": "user", "content": f"OBSERVATION ({act_type}):\n{result}"})
                agent["pending_action"] = None
                agent["status"] = "running"
                st.rerun()
            elif reject:
                agent["transcript"].append({"type": "rejected"})
                agent["history"].append({"role": "user", "content": f"The user REJECTED this action ({act_type}). Try a different approach."})
                agent["pending_action"] = None
                agent["status"] = "running"
                st.rerun()

        if agent["status"] not in ("idle",):
            if st.button("🆕 Clear Conversation (start a brand new session)", key="agent_stop"):
                st.session_state.coding_agent = {
                    "project_dir": agent["project_dir"], "history": [], "transcript": [],
                    "status": "idle", "pending_action": None, "steps": 0, "task_gen": agent.get("task_gen", 0) + 1,
                }
                st.session_state.pop("agent_active_session", None)
                st.rerun()

        # ---- What should the agent do next? — kept at the bottom of the page, ------
        # below the transcript, so it always sits right after the latest activity. --
        if agent["status"] in ("idle", "done", "stopped", "error"):
            st.divider()
            continuing = bool(agent["history"])
            label = "What should the agent do next?" if continuing else "What should the agent do?"
            agent_task = st.text_area(
                label, height=90, key=f"agent_task_{agent['task_gen']}",
                placeholder="e.g. 'Add type hints and docstrings to every function in utils.py', "
                            "'Find and fix the bug causing the KeyError in main.py, then run the tests'")
            btn_label = "▶️ Continue" if continuing else "🚀 Start Agent"
            if st.button(btn_label, type="primary", disabled=not (dir_ready and agent_task.strip())):
                log_event(_CURRENT_USER, "coding_agent_run", detail="continue" if continuing else "start")
                add_memory_fact(_CURRENT_USER, f"Used the Coding Agent for: {agent_task.strip()[:200]}")
                if continuing:
                    agent["history"].append({"role": "user", "content": f"NEW INSTRUCTION: {agent_task.strip()}"})
                    agent["transcript"].append({"type": "task", "text": agent_task.strip()})
                else:
                    agent["project_dir"] = project_dir_input
                    agent["history"] = [{"role": "user",
                                        "content": f"{_agent_environment_note(agent_python)}\n\nTASK: {agent_task.strip()}"}]
                    agent["transcript"] = [{"type": "task", "text": agent_task.strip()}]
                agent["status"] = "running"
                agent["pending_action"] = None
                agent["steps"] = 0  # fresh step budget per instruction; history/transcript stay cumulative
                agent["task_gen"] = agent.get("task_gen", 0) + 1  # fresh empty input box next render
                st.rerun()

        # ---- Autosave — runs on every render where the agent has history. Steps in the
        # loop above each end in st.rerun(), so this only actually reaches disk once the
        # agent settles (awaiting approval, done, error, or stopped) — by then agent["history"]
        # /["transcript"] already reflect every step that happened, so nothing is lost.
        if agent["history"]:
            if not st.session_state.get("agent_active_session"):
                first_task = next((t["text"] for t in agent["transcript"] if t.get("type") == "task"), "session")
                st.session_state["agent_active_session"] = autoname_session(first_task, list_coding_sessions())
            save_coding_session(st.session_state["agent_active_session"], agent)

    elif coding_mode == "🧭 Claude Code (native CLI, any LLM)":
        # ── Initialise session state for hideable panels ──────────────────────
        if "cc_show_sessions" not in st.session_state:
            st.session_state.cc_show_sessions = True
        if "cc_show_files" not in st.session_state:
            st.session_state.cc_show_files = True
        if "cc_editing_turn" not in st.session_state:
            st.session_state.cc_editing_turn = None
        if "cc_artifacts" not in st.session_state:
            st.session_state.cc_artifacts = {}  # {rel_path: {content, turn, timestamp}}
        if "cc_hide_sidebar" not in st.session_state:
            st.session_state.cc_hide_sidebar = False

        # CSS to hide the main Streamlit sidebar when toggled
        if st.session_state.cc_hide_sidebar:
            st.markdown("""
            <style>
            section[data-testid="stSidebar"] { display: none !important; }
            div[data-testid="stAppViewContainer"] > section:first-child { display: none !important; }
            </style>
            """, unsafe_allow_html=True)

        claude_path = shutil.which("claude")
        if not claude_path:
            st.error("The `claude` CLI isn't on this machine's PATH.")
            st.markdown(
                "Install it, then reload this page:\n\n"
                "```bash\nnpm install -g @anthropic-ai/claude-code\n```\n\n"
                "Requires Node.js 18+. See the [Claude Code docs](https://docs.claude.com/en/docs/claude-code) "
                "for other install options.")
            claude_path = "claude"  # fallback so the UI still renders; launch button stays disabled

        if "cc_native" not in st.session_state:
            st.session_state.cc_native = {"turns": [], "run_gen": 0}
        ccn = st.session_state.cc_native

        # ── LLM backend — reuse sidebar provider config ─────────────────────
        cc_env_set, cc_env_clear, cc_model_final = {}, [], crew_model
        _cc_backend_label = ""
        if provider_choice == "Ollama (Local/Cloud)":
            cc_env_set = {
                "ANTHROPIC_BASE_URL": selected_api_base,
                "ANTHROPIC_AUTH_TOKEN": selected_api_key or "ollama",
                "ANTHROPIC_DEFAULT_HAIKU_MODEL": crew_model,
                "ANTHROPIC_DEFAULT_SONNET_MODEL": crew_model,
                "ANTHROPIC_DEFAULT_OPUS_MODEL": crew_model,
                "ANTHROPIC_API_KEY": "",
                "DISABLE_TELEMETRY": "1",
            }
            _cc_backend_label = f"Ollama · {crew_model}"
        elif provider_choice == "Anthropic (Claude)":
            if selected_api_key:
                cc_env_set = {"ANTHROPIC_API_KEY": selected_api_key}
            cc_env_clear = ["ANTHROPIC_BASE_URL", "ANTHROPIC_AUTH_TOKEN"]
            _cc_backend_label = f"Anthropic · {crew_model}"
        elif provider_choice in ("Qwen (Anthropic)", "DeepSeek (Anthropic)", "Custom (Anthropic-compatible)"):
            cc_env_set = {"ANTHROPIC_BASE_URL": selected_api_base or "",
                          "ANTHROPIC_AUTH_TOKEN": selected_api_key or "",
                          "ANTHROPIC_API_KEY": ""}
            _cc_backend_label = f"{provider_choice} · {crew_model}"
        else:
            cc_env_set = {"ANTHROPIC_BASE_URL": selected_api_base or "",
                          "ANTHROPIC_AUTH_TOKEN": selected_api_key or "placeholder",
                          "ANTHROPIC_API_KEY": ""}
            _cc_backend_label = f"{provider_choice} (proxy needed)"

        # ── Project folder ──────────────────────────────────────────────────
        def _cc_browse_callback():
            picked = browse_for_folder(initial_dir=st.session_state.get("cc_project_dir", "") or USER_SANDBOX_DIR)
            if not picked:
                st.session_state["cc_browse_result"] = "cancelled"
            elif not (_CURRENT_IS_ADMIN or path_in_sandbox(picked, USER_SANDBOX_DIR)):
                st.session_state["cc_browse_result"] = "outside_sandbox"
            else:
                st.session_state["cc_browse_result"] = "picked"
                st.session_state["cc_project_dir"] = picked

        cc_project_dir = st.session_state.get("cc_project_dir", "")
        cc_dir_exists = bool(cc_project_dir) and os.path.isdir(cc_project_dir)
        cc_dir_sandboxed = _CURRENT_IS_ADMIN or (bool(cc_project_dir) and path_in_sandbox(cc_project_dir, USER_SANDBOX_DIR))
        cc_dir_ready = cc_dir_exists and cc_dir_sandboxed

        # ── Top toolbar ────────────────────────────────────────────────────
        with st.container():
            tb1, tb2, tb3, tb4, tb5, tb6 = st.columns([1.1, 0.5, 0.5, 2.8, 0.5, 0.5])
            with tb1:
                cc_project_dir = st.text_input(
                    "📁 Project folder", value=cc_project_dir,
                    key="cc_project_dir_v2", placeholder=USER_SANDBOX_DIR,
                    label_visibility="collapsed")
            with tb2:
                st.button("📂", key="cc_browse_dir_v2", use_container_width=True,
                         on_click=_cc_browse_callback,
                         help="Browse for a folder on this machine")
            with tb3:
                st.session_state.cc_show_sessions = st.toggle(
                    "💬", value=st.session_state.cc_show_sessions, key="cc_toggle_sessions",
                    help="Show/hide session list panel")
            with tb4:
                if cc_project_dir and not cc_dir_sandboxed:
                    st.error(f"⛔ Outside sandbox: `{USER_SANDBOX_DIR}`")
                elif cc_project_dir and not cc_dir_exists:
                    st.warning("Folder doesn't exist — create it first.")
                else:
                    st.caption(f"`{cc_project_dir or '(no folder)'}`  ·  {_cc_backend_label}")
            with tb5:
                st.session_state.cc_show_files = st.toggle(
                    "📄", value=st.session_state.cc_show_files, key="cc_toggle_files",
                    help="Show/hide file browser panel")
            with tb6:
                if st.button("⬅️" if not st.session_state.cc_hide_sidebar else "➡️",
                           key="cc_toggle_sidebar_btn", use_container_width=True,
                           help="Show/hide the configuration sidebar"):
                    st.session_state.cc_hide_sidebar = not st.session_state.cc_hide_sidebar
                    st.rerun()

            _cc_browse_result = st.session_state.pop("cc_browse_result", None)
            if _cc_browse_result == "cancelled":
                st.toast("No folder chosen — type the path manually.")
            elif _cc_browse_result == "outside_sandbox":
                st.toast(f"⛔ Outside your sandbox ({USER_SANDBOX_DIR}).")

        # ── Options bar ────────────────────────────────────────────────────
        with st.container():
            oc1, oc2, oc3 = st.columns([2, 1, 1])
            with oc1:
                cc_auto_approve = st.toggle(
                    "⚡ Auto-approve", value=True, key="cc_auto_approve_v2",
                    help="Skip permission prompts (--dangerously-skip-permissions). "
                         "Only use on folders you trust.")
            with oc2:
                has_history = bool(ccn["turns"])
                if has_history:
                    continuing = st.toggle("🔄 Continue", value=True, key="cc_continue_toggle_v2",
                                          help="Use --continue to resume the previous session.")
                else:
                    continuing = False
            with oc3:
                if has_history and st.button("🧹 Clear chat", key="cc_clear_v2", use_container_width=True):
                    ccn["turns"] = []
                    ccn["run_gen"] = ccn.get("run_gen", 0) + 1
                    st.session_state.pop("cc_active_session", None)
                    st.rerun()

        # ── CSS: sticky side panels + scrollable chat so sessions/files
        # stay visible even when the conversation is long ──────────────────
        st.markdown("""
        <style>
        .cc-sticky-left, .cc-sticky-right {
            position: sticky;
            top: 0.5rem;
            max-height: 85vh;
            overflow-y: auto;
            padding-right: 0.3rem;
        }
        .cc-chat-scroll {
            max-height: 58vh;
            overflow-y: auto;
            padding: 0.5rem;
            margin-bottom: 0.5rem;
            border: 1px solid rgba(128,128,128,0.12);
            border-radius: 8px;
        }
        .cc-chat-input-area {
            position: sticky;
            bottom: 0;
            padding-top: 0.3rem;
            z-index: 10;
        }
        .cc-stream-output {
            width: 100% !important;
            min-width: 100% !important;
            white-space: pre-wrap;
            word-break: break-word;
            font-family: "SFMono-Regular", Consolas, "Liberation Mono", Menlo, monospace;
            font-size: 0.82rem;
            line-height: 1.5;
            background: rgba(128,128,128,0.06);
            border-radius: 6px;
            padding: 0.8rem;
            min-height: 120px;
            max-height: 60vh;
            overflow-y: auto;
            box-sizing: border-box;
        }
        </style>
        """, unsafe_allow_html=True)

        st.divider()

        # ── 3-panel layout: sessions | chat | files ─────────────────────────
        _show_left = st.session_state.cc_show_sessions
        _show_right = st.session_state.cc_show_files

        if _show_left and _show_right:
            left_col, mid_col, right_col = st.columns([1.8, 4, 1.8])
        elif _show_left:
            left_col, mid_col, right_col = st.columns([2, 5, 0.02])
        elif _show_right:
            left_col, mid_col, right_col = st.columns([0.02, 5, 2])
        else:
            left_col, mid_col = st.columns([0.02, 5])
            right_col = None

        # ── LEFT PANEL: Sessions ─────────────────────────────────────────────
        if _show_left:
            with left_col:
                st.markdown('<div class="cc-sticky-left">', unsafe_allow_html=True)
                st.caption("##### 💬 Sessions")
                cc_sessions = list_cc_sessions()
                active = st.session_state.get("cc_active_session", "")

                # New session button
                if st.button("➕ New session", key="cc_new_session_btn", use_container_width=True):
                    st.session_state.pop("cc_active_session", None)
                    ccn["turns"] = []
                    ccn["run_gen"] = ccn.get("run_gen", 0) + 1
                    st.rerun()

                if not cc_sessions:
                    st.caption("No saved sessions yet — your first message will create one.")
                else:
                    for sname in cc_sessions:
                        is_active = (sname == active)
                        _btn_label = ("🟢 " if is_active else "   ") + sname[:35]
                        srow1, srow2 = st.columns([5, 1])
                        with srow1:
                            if st.button(_btn_label, key=f"cc_sess_{sname}",
                                        use_container_width=True,
                                        help=f"Load '{sname}'",
                                        type="primary" if is_active else "secondary"):
                                if sname != active:
                                    payload, err = load_cc_session(sname)
                                    if err:
                                        st.error(f"Load failed: {err}")
                                    else:
                                        st.session_state.cc_native = {
                                            "turns": payload.get("turns", []),
                                            "run_gen": ccn.get("run_gen", 0) + 1,
                                        }
                                        st.session_state["cc_project_dir"] = payload.get("project_dir", "")
                                        st.session_state["cc_active_session"] = sname
                                        st.session_state.cc_editing_turn = None
                                        st.rerun()
                        with srow2:
                            if st.button("🗑", key=f"cc_del_{sname}",
                                        help=f"Delete '{sname}'"):
                                delete_cc_session(sname)
                                if active == sname:
                                    st.session_state.pop("cc_active_session", None)
                                    ccn["turns"] = []
                                    ccn["run_gen"] = ccn.get("run_gen", 0) + 1
                                st.rerun()

                    # Rename
                    if active:
                        st.divider()
                        new_name = st.text_input(
                            "Rename", value=active, key="cc_rename_input",
                            label_visibility="collapsed")
                        if st.button("✏️ Rename", key="cc_rename_btn", use_container_width=True,
                                    disabled=(new_name.strip() == active or not new_name.strip())):
                            ok, err = save_cc_session(new_name.strip(), ccn, cc_project_dir)
                            if ok:
                                if active != new_name.strip():
                                    delete_cc_session(active)
                                st.session_state["cc_active_session"] = new_name.strip()
                                st.rerun()
                            else:
                                st.error(f"Rename failed: {err}")
        if _show_left:
            with left_col:
                st.markdown('</div>', unsafe_allow_html=True)  # close cc-sticky-left

        # ── CENTER: Chat ─────────────────────────────────────────────────────
        with mid_col:
            # Editing a previous prompt?
            _editing = st.session_state.cc_editing_turn
            if _editing is not None and 0 <= _editing < len(ccn["turns"]):
                with st.container():
                    st.caption(f"✏️ Editing turn #{_editing + 1}")
                    edited_task = st.text_area(
                        "Edit your prompt", value=ccn["turns"][_editing]["task"],
                        height=120, key=f"cc_edit_{_editing}")
                    ec1, ec2, ec3 = st.columns([1, 1, 2])
                    with ec1:
                        if st.button("✅ Save & re-run", key=f"cc_edit_save_{_editing}", use_container_width=True):
                            # Trim turns from this point onward, re-run with edited prompt
                            ccn["turns"] = ccn["turns"][:_editing]
                            st.session_state.cc_editing_turn = None
                            st.session_state["_cc_edited_prompt"] = edited_task.strip()
                            st.rerun()
                    with ec2:
                        if st.button("Cancel", key=f"cc_edit_cancel_{_editing}", use_container_width=True):
                            st.session_state.cc_editing_turn = None
                            st.rerun()

            # Render conversation
            st.markdown('<div class="cc-chat-scroll">', unsafe_allow_html=True)
            for i, turn in enumerate(ccn["turns"]):
                # User message — editable
                with st.chat_message("user", avatar="🧑"):
                    st.markdown(turn["task"])
                    if st.button("✏️", key=f"cc_edit_btn_{i}",
                                help="Edit this prompt and re-run from here"):
                        st.session_state.cc_editing_turn = i
                        st.rerun()

                # Assistant message — formatted nicely
                with st.chat_message("assistant", avatar="✨"):
                    # Split output into sections: command, tool calls, thinking, final answer
                    _out = turn["output"]
                    # Try to extract and render structured parts
                    _cmd_end = _out.find("\n\n")
                    if _cmd_end > 0 and _out.startswith("$ "):
                        _cmd = _out[:_cmd_end]
                        with st.expander("🔧 Command", expanded=False):
                            st.code(_cmd, language="bash")
                        _body = _out[_cmd_end:].strip()
                    else:
                        _body = _out

                    # Render the body: code blocks get syntax highlighting,
                    # regular text gets markdown
                    st.markdown(_body)
                    _rc = turn["returncode"]
                    if _rc == 0:
                        st.caption("✅ Completed")
                    elif _rc is not None:
                        st.caption(f"⚠️ Exit code {_rc}")

            if not ccn["turns"]:
                st.info("👋 Describe your task below — Claude will read, write, and run code in the project folder.")
            st.markdown('</div>', unsafe_allow_html=True)  # close cc-chat-scroll

            # ── Input area ─────────────────────────────────────────────────
            st.markdown('<div class="cc-chat-input-area">', unsafe_allow_html=True)
            _edited_prompt = st.session_state.pop("_cc_edited_prompt", None)
            _initial_prompt = _edited_prompt if _edited_prompt else ""
            label = "What should Claude do next?" if (ccn["turns"] and not _edited_prompt) else "What should Claude do?"
            cc_task = st.text_area(
                label, height=80, key=f"cc_task_v2_{ccn['run_gen']}",
                value=_initial_prompt,
                placeholder="e.g. 'Add unit tests for the parser module and run them'")

            launch_disabled = not (cc_dir_ready and cc_task.strip() and (cc_model_final or provider_choice == "Anthropic (Claude)"))
            ccb1, ccb2 = st.columns([3, 1])
            with ccb1:
                pass
            with ccb2:
                btn_label = "▶️ Send" if (ccn["turns"] and not _edited_prompt) else "🚀 Send"
                if st.button(btn_label, type="primary", disabled=launch_disabled, key="cc_launch_v2", use_container_width=True):
                    st.session_state["_cc_launch_prompt"] = cc_task.strip()
                    st.session_state["_cc_launch_continuing"] = continuing and not _edited_prompt
                    st.rerun()

            # ── Launch execution (OUTSIDE the button column → full-width display) ──
            _launch_prompt = st.session_state.pop("_cc_launch_prompt", None)
            if _launch_prompt:
                _launch_continuing = st.session_state.pop("_cc_launch_continuing", False)
                log_event(_CURRENT_USER, "claude_code_run", detail="continue" if _launch_continuing else "start")
                add_memory_fact(_CURRENT_USER, f"Used Claude Code for: {_launch_prompt[:200]}")

                # Snapshot project files before this turn (for artifact detection)
                _before_snapshot = {}
                if cc_dir_ready:
                    try:
                        for _r, _ds, _fs in os.walk(cc_project_dir):
                            _ds[:] = [d for d in _ds if not d.startswith(".")]
                            for _fn in _fs:
                                if not _fn.startswith("."):
                                    _fp = os.path.join(_r, _fn)
                                    _rel = os.path.relpath(_fp, cc_project_dir)
                                    _before_snapshot[_rel] = os.path.getmtime(_fp)
                    except Exception:
                        pass

                args = [claude_path, "-p", "--output-format", "stream-json", "--verbose"]
                if _launch_continuing:
                    args.append("--continue")
                if cc_model_final:
                    args += ["--model", cc_model_final]
                if cc_auto_approve:
                    args.append("--dangerously-skip-permissions")

                run_env = os.environ.copy()
                for k in cc_env_clear:
                    run_env.pop(k, None)
                run_env.update(cc_env_set)

                prompt_text = _launch_prompt

                with st.chat_message("user", avatar="🧑"):
                    st.markdown(prompt_text)

                q = queue.Queue()
                ctx = get_script_run_ctx()

                def _cc_worker(args, cwd, env, prompt_text, q, ctx):
                    add_script_run_ctx(threading.current_thread(), ctx)
                    try:
                        proc = subprocess.Popen(args, cwd=cwd, env=env, stdin=subprocess.PIPE,
                                               stdout=subprocess.PIPE, stderr=subprocess.STDOUT,
                                               text=True, bufsize=1, encoding="utf-8", errors="replace")
                        proc.stdin.write(prompt_text)
                        proc.stdin.close()
                        for line in proc.stdout:
                            q.put(line)
                        proc.wait()
                        q.put(("__RC__", proc.returncode))
                    except Exception as e:
                        q.put(f"\n[ERROR launching claude: {e}]\n")
                        q.put(("__RC__", -1))
                    finally:
                        q.put(None)

                with st.chat_message("assistant", avatar="✨"):
                    log_ph = st.empty()
                    buf = f"$ {' '.join(shlex.quote(a) for a in args)}   (prompt piped via stdin)\n\n"
                    log_ph.markdown(
                        f'<pre class="cc-stream-output">⏳ Working…</pre>',
                        unsafe_allow_html=True)

                    worker = threading.Thread(target=_cc_worker, args=(args, cc_project_dir, run_env, prompt_text, q, ctx))
                    worker.start()
                    returncode = None
                    while True:
                        try:
                            item = q.get(timeout=0.1)
                        except queue.Empty:
                            if not worker.is_alive():
                                break
                            continue
                        if item is None:
                            break
                        if isinstance(item, tuple):
                            returncode = item[1]
                            continue
                        stripped = item.strip()
                        if stripped:
                            try:
                                evt = json.loads(stripped)
                                buf += format_claude_code_event(evt)
                            except json.JSONDecodeError:
                                buf += item
                            log_ph.markdown(
                                f'<pre class="cc-stream-output">{html.escape(buf)}</pre>',
                                unsafe_allow_html=True)

                ccn["turns"].append({"task": prompt_text, "output": buf, "returncode": returncode})
                ccn["run_gen"] = ccn.get("run_gen", 0) + 1

                # Detect artifacts: files created or modified during this turn
                _current_turn = len(ccn["turns"]) - 1
                if cc_dir_ready and _before_snapshot:
                    try:
                        for _r, _ds, _fs in os.walk(cc_project_dir):
                            _ds[:] = [d for d in _ds if not d.startswith(".")]
                            for _fn in _fs:
                                if not _fn.startswith("."):
                                    _fp = os.path.join(_r, _fn)
                                    _rel = os.path.relpath(_fp, cc_project_dir)
                                    _prev_mtime = _before_snapshot.get(_rel, 0)
                                    _curr_mtime = os.path.getmtime(_fp)
                                    if _curr_mtime > _prev_mtime + 0.5:  # new or modified
                                        try:
                                            with open(_fp, "r", encoding="utf-8", errors="replace") as _af:
                                                _acontent = _af.read()
                                            st.session_state.cc_artifacts[_rel] = {
                                                "content": _acontent[:50000],
                                                "turn": _current_turn,
                                                "timestamp": _curr_mtime,
                                                "is_new": _prev_mtime == 0,
                                            }
                                        except Exception:
                                            pass
                    except Exception:
                        pass

                if not st.session_state.get("cc_active_session"):
                    st.session_state["cc_active_session"] = autoname_session(prompt_text, list_cc_sessions())
                save_cc_session(st.session_state["cc_active_session"], ccn, cc_project_dir)
                st.rerun()

            st.markdown('</div>', unsafe_allow_html=True)  # close cc-chat-input-area

        # ── RIGHT PANEL: Files / Artifacts ───────────────────────────────────
        if _show_right and right_col is not None:
            with right_col:
                st.markdown('<div class="cc-sticky-right">', unsafe_allow_html=True)

                # ── Artifacts: files created/modified during this session ───
                _artifacts = st.session_state.cc_artifacts
                if _artifacts:
                    st.caption("##### 🏗️ Artifacts")
                    # Group by turn, newest first
                    _by_turn = {}
                    for _rel, _info in _artifacts.items():
                        _t = _info.get("turn", 0)
                        _by_turn.setdefault(_t, []).append((_rel, _info))
                    for _t in sorted(_by_turn, reverse=True):
                        if _t < len(ccn["turns"]):
                            _task_preview = ccn["turns"][_t]["task"][:40]
                            st.caption(f"*Turn {_t+1}:* {_task_preview}…")
                        for _rel, _info in _by_turn[_t]:
                            _is_new = _info.get("is_new", False)
                            _badge = "🆕" if _is_new else "✏️"
                            _fname = os.path.basename(_rel)
                            _suffix = _fname.rsplit(".", 1)[-1] if "." in _fname else ""
                            _lang = _suffix if _suffix in ("py","js","ts","json","md","html","css","yaml","yml","toml","sh","sql","xml") else "text"
                            _content = _info.get("content", "")
                            _lines = len(_content.splitlines()) if _content else 0
                            with st.expander(f"{_badge} {_fname} ({_lines} lines)", expanded=False):
                                st.caption(f"`{_rel}`")
                                if _content:
                                    st.code(_content[:8000], language=_lang, line_numbers=True)
                                    if len(_content) > 8000:
                                        st.caption(f"… truncated ({len(_content):,} total chars)")
                    if _artifacts:
                        st.caption(f"📦 {len(_artifacts)} artifact(s) this session")
                        if st.button("🧹 Clear artifacts", key="cc_clear_artifacts", use_container_width=True):
                            st.session_state.cc_artifacts = {}
                            st.rerun()
                    st.divider()

                st.caption("##### 📄 Files")
                if not cc_dir_ready:
                    st.caption("Set a project folder to browse files.")
                else:
                    # Quick file listing
                    _all_files = []
                    try:
                        for _root, _dirs, _files in os.walk(cc_project_dir):
                            # Skip hidden dirs
                            _dirs[:] = [d for d in _dirs if not d.startswith(".")]
                            for _fn in _files:
                                if not _fn.startswith("."):
                                    _full = os.path.join(_root, _fn)
                                    _rel = os.path.relpath(_full, cc_project_dir)
                                    _all_files.append((_rel, _full))
                    except Exception:
                        _all_files = []

                    if not _all_files:
                        st.caption("No files yet — Claude will create them as it works.")
                    else:
                        # File tree
                        _by_dir = {}
                        for _rel, _full in sorted(_all_files):
                            _d = os.path.dirname(_rel) or "."
                            _by_dir.setdefault(_d, []).append((_rel, _full))

                        for _d in sorted(_by_dir):
                            if _d != ".":
                                st.caption(f"📁 {_d}/")
                            for _rel, _full in _by_dir[_d]:
                                _fname = os.path.basename(_rel)
                                _size = os.path.getsize(_full) if os.path.isfile(_full) else 0
                                _suffix = _fname.rsplit(".", 1)[-1] if "." in _fname else ""
                                _emoji = {"py": "🐍", "js": "🟨", "ts": "🔷", "json": "📋",
                                          "md": "📝", "txt": "📄", "html": "🌐", "css": "🎨",
                                          "yaml": "⚙️", "yml": "⚙️", "toml": "⚙️",
                                          "sh": "💻", "bat": "💻", "ps1": "💻",
                                          "zip": "📦", "gz": "📦", "png": "🖼️", "jpg": "🖼️",
                                          "svg": "🖼️", "pdf": "📕", "ipynb": "📓",
                                          }.get(_suffix, "📄")
                                _label = f"{_emoji} {_fname} ({_size:,}b)"

                                if st.button(_label, key=f"cc_file_{_rel}", use_container_width=True,
                                            help=f"Preview {_fname}"):
                                    st.session_state["_cc_preview_file"] = _full
                                    st.session_state["_cc_preview_rel"] = _rel
                                    st.rerun()

                        # Preview selected file
                        _preview_path = st.session_state.get("_cc_preview_file")
                        if _preview_path:
                            st.divider()
                            _preview_rel = st.session_state.get("_cc_preview_rel", "")
                            st.caption(f"📄 Preview: **{_preview_rel}**")
                            if st.button("✕ Close", key="cc_close_preview"):
                                st.session_state.pop("_cc_preview_file", None)
                                st.session_state.pop("_cc_preview_rel", None)
                                st.rerun()
                            try:
                                with open(_preview_path, "r", encoding="utf-8", errors="replace") as _pf:
                                    _pcontent = _pf.read()
                                # Truncate large files
                                if len(_pcontent) > 50000:
                                    _pcontent = _pcontent[:50000] + "\n\n… (truncated after 50 KB)"
                                _lang = _preview_rel.rsplit(".", 1)[-1] if "." in _preview_rel else "text"
                                st.code(_pcontent, language=_lang, line_numbers=True)
                                with open(_preview_path, "rb") as _pfb:
                                    st.download_button(
                                        "⬇️ Download", data=_pfb.read(),
                                        file_name=os.path.basename(_preview_rel),
                                        key=f"cc_dl_preview_{_preview_rel}", use_container_width=True)
                            except Exception as _pe:
                                st.caption(f"Can't preview: {_pe}")

                        # Download whole project
                        if _all_files:
                            st.divider()
                            _zip_buf = io.BytesIO()
                            with zipfile.ZipFile(_zip_buf, "w", zipfile.ZIP_DEFLATED) as _zf:
                                for _rel, _full in _all_files:
                                    _zf.write(_full, _rel)
                            _proj_name = os.path.basename(cc_project_dir) or "project"
                            st.download_button(
                                f"⬇️ Download all ({_proj_name}.zip)",
                                data=_zip_buf.getvalue(), file_name=f"{_proj_name}.zip",
                                key="cc_dl_zip_v2", use_container_width=True)
        if _show_right and right_col is not None:
            with right_col:
                st.markdown('</div>', unsafe_allow_html=True)  # close cc-sticky-right

        # ── Bottom status bar ────────────────────────────────────────────────
        st.divider()
        b1, b2, b3 = st.columns([2, 2, 1])
        with b1:
            if cc_dir_ready:
                _nfiles = 0
                try:
                    for _, _, files in os.walk(cc_project_dir):
                        _nfiles += len([f for f in files if not f.startswith(".")])
                except Exception:
                    _nfiles = "?"
                st.caption(f"📁 `{cc_project_dir}` — {_nfiles} files")
            else:
                st.caption("📁 No project folder selected")
        with b2:
            st.caption(f"🔌 {_cc_backend_label}")
        with b3:
            has_history = bool(ccn["turns"])
            st.caption(f"💬 {len(ccn['turns'])} turn(s)" if has_history else "💬 New session")

    else:  # 🖥️ CDesktop
        st.subheader("🖥️ CDesktop — Desktop Agent")
        st.caption("CDesktop is a desktop agent that opens a code-aware AI chat in a new "
                   "browser tab. It runs locally via `npx cdesktop` and gives the AI access "
                   "to your project files for reading, writing, and running commands.")

        # Pick a project folder (reuse the sandbox logic from Agent/Claude Code)
        cdesk_dir = st.text_input(
            "📁 Project folder", key="cdesk_project_dir",
            value=USER_SANDBOX_DIR,
            placeholder=USER_SANDBOX_DIR,
            help="Folder where CDesktop will start — it can read/write/run commands here.")

        cdesk_dir_ok = bool(cdesk_dir) and os.path.isdir(cdesk_dir)
        cdesk_sandboxed = _CURRENT_IS_ADMIN or (bool(cdesk_dir) and path_in_sandbox(cdesk_dir, USER_SANDBOX_DIR))

        if cdesk_dir and not cdesk_sandboxed:
            st.error(f"⛔ That folder is outside your sandbox. Use a path under: `{USER_SANDBOX_DIR}`")
        elif cdesk_dir and not cdesk_dir_ok:
            st.warning("Folder doesn't exist — create it first or pick another.")

        st.divider()

        col_a, col_b = st.columns([3, 1])
        with col_a:
            st.info(
                "💡 **What happens:** clicking **Launch** runs `npx cdesktop` in the selected "
                "project folder. A new browser tab will open with the CDesktop chat interface — "
                "it has full access to that folder. Close the browser tab to end the session."
            )
        with col_b:
            launch_disabled = not (cdesk_dir_ok and cdesk_sandboxed)
            if st.button("🚀 Launch CDesktop", type="primary", use_container_width=True,
                         disabled=launch_disabled, key="cdesk_launch"):
                log_event(_CURRENT_USER, "cdesktop_launch", detail=cdesk_dir)
                add_memory_fact(_CURRENT_USER, f"Launched CDesktop in: {cdesk_dir}")

                with st.status("🖥️ Launching CDesktop…", expanded=True) as cdesk_status:
                    st.write(f"Running `npx cdesktop` in `{cdesk_dir}`…")
                    try:
                        # Resolve npx path — on Windows, subprocess.Popen with a list
                        # doesn't resolve .cmd/.bat extensions the way the shell does.
                        npx_cmd = "npx"
                        if sys.platform == "win32":
                            _which = shutil.which("npx") or shutil.which("npx.cmd")
                            if _which:
                                npx_cmd = _which
                        # Run detached so it survives this script run.
                        popen_kwargs = dict(
                            cwd=cdesk_dir,
                            stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
                        if sys.platform == "win32":
                            popen_kwargs["creationflags"] = (
                                subprocess.DETACHED_PROCESS | subprocess.CREATE_NEW_PROCESS_GROUP)
                        else:
                            popen_kwargs["start_new_session"] = True
                        proc = subprocess.Popen([npx_cmd, "cdesktop"], **popen_kwargs)
                        cdesk_status.update(
                            label="✅ CDesktop launched — a new browser tab should open shortly.",
                            state="complete")
                    except Exception as e:
                        cdesk_status.update(
                            label=f"❌ Failed to launch: {e}", state="error")

        st.caption("Make sure Node.js 18+ is installed (`node --version`). "
                  "`npx cdesktop` will download and start CDesktop on first run — "
                  "it opens a new browser tab automatically.")

# =============================== MAKE PRESENTATION TAB =======================
with tab_present:
    st.subheader("🖼️ Make Presentation")
    st.caption("Turn an idea — or something already produced elsewhere in this app — into an "
              "actual slide deck (.pptx). Pick a source, generate an outline, tweak it, then "
              "build the file.")

    if "pres_outline" not in st.session_state:
        st.session_state.pres_outline = None

    # ── Method selector ──────────────────────────────────────────────────
    pres_method = st.radio(
        "Method", ["🧠 LLM generate (.pptx)", "🎯 Use Presenton (local CLI)"],
        horizontal=True, key="pres_method")

    if pres_method == "🎯 Use Presenton (local CLI)":
        # ── Presenton mode ───────────────────────────────────────────────
        st.caption("Uses the `presenton` CLI installed on this machine to build a presentation "
                  "directly from your source material.")

        pres_presenton_path = st.text_input(
            "Presenton path or command",
            value=shutil.which("presenton") or "presenton",
            key="pres_presenton_path",
            help="Full path to the presenton executable, or just 'presenton' if it's on PATH.")

        pres_presenton_args = st.text_input(
            "Extra CLI arguments (optional)", key="pres_presenton_args",
            placeholder="e.g. --theme dark --format pptx")

        st.divider()

        st.markdown("##### 1. Source material")
        _wp_projects = list_wp_projects()
        _sv_projects = list_survey_projects()
        _crew_projects = list_crew_projects()
        _agent_sessions = list_coding_sessions()
        _cc_sessions = list_cc_sessions()
        _pres_conv = st.session_state.get("conversations", {}).get(st.session_state.get("current_conv_id"))

        pres_categories = ["✍️ Type/paste an idea"]
        if _wp_projects:
            pres_categories.append("✍️ Write Paper project")
        if _sv_projects:
            pres_categories.append("📚 Survey project")
        if _crew_projects:
            pres_categories.append("🔬 Research Crew project")
        if st.session_state.get("sm_summary"):
            pres_categories.append("📖 Summarize Paper tab — last summary")
        if _agent_sessions:
            pres_categories.append("🤖 Coding Agent session")
        if _cc_sessions:
            pres_categories.append("🧭 Claude Code session")
        if st.session_state.get("coding_code", "").strip():
            pres_categories.append("💻 Coding tab — current Quick Edit code")
        if _pres_conv and _pres_conv.get("messages"):
            pres_categories.append("💬 Chat tab — current conversation")
        pres_categories.append("📤 Upload file(s)")

        _pres_items_by_category = {
            "✍️ Write Paper project": _wp_projects,
            "📚 Survey project": _sv_projects,
            "🔬 Research Crew project": _crew_projects,
            "🤖 Coding Agent session": _agent_sessions,
            "🧭 Claude Code session": _cc_sessions,
        }

        pc1, pc2 = st.columns(2)
        with pc1:
            pres_category = st.selectbox("Source", pres_categories, key="preson_category")
        pres_item = None
        if pres_category in _pres_items_by_category:
            with pc2:
                pres_item = st.selectbox("Which one", _pres_items_by_category[pres_category], key="preson_item")

        pres_source_text, pres_default_title = "", ""
        # (same source-text extraction as LLM mode)
        if pres_category == "✍️ Type/paste an idea":
            pres_source_text = st.text_area(
                "Describe the idea for the presentation:", height=160, key="preson_manual_idea",
                placeholder="e.g. 'A wearable ECG patch for continuous heart monitoring...'")
        elif pres_category == "✍️ Write Paper project" and pres_item:
            payload, _ = load_wp_project(pres_item)
            wp_data = (payload or {}).get("wp_paper") or {}
            pres_source_text = wp_data.get("md") or (payload or {}).get("wp_idea", "")
            pres_default_title = wp_data.get("title") or pres_item
        elif pres_category == "📚 Survey project" and pres_item:
            payload, _ = load_survey_project(pres_item)
            sv_res = (payload or {}).get("survey_result") or {}
            pres_source_text = sv_res.get("md", "")
            pres_default_title = pres_item
        elif pres_category == "🔬 Research Crew project" and pres_item:
            payload, _ = load_crew_project(pres_item)
            res_data = (payload or {}).get("research_result") or {}
            pres_source_text = res_data.get("result", "")
            pres_default_title = pres_item
        elif pres_category == "📖 Summarize Paper tab — last summary":
            pres_source_text = st.session_state["sm_summary"]
            pres_default_title = re.sub(r"\.pdf$", "", st.session_state.get("sm_file_name", "Paper"), flags=re.I)
        elif pres_category == "🤖 Coding Agent session" and pres_item:
            payload, _ = load_coding_session(pres_item)
            transcript = (payload or {}).get("transcript", [])
            pres_source_text = "\n\n".join(f"{t.get('type', '').upper()}: {t.get('text', '')}"
                                          for t in transcript if t.get("text"))
            pres_default_title = pres_item
        elif pres_category == "🧭 Claude Code session" and pres_item:
            payload, _ = load_cc_session(pres_item)
            turns = (payload or {}).get("turns", [])
            pres_source_text = "\n\n".join(f"TASK: {t['task']}\nRESULT:\n{t['output']}" for t in turns)
            pres_default_title = pres_item
        elif pres_category == "💻 Coding tab — current Quick Edit code":
            pres_source_text = f"```{st.session_state.get('coding_lang', '')}\n{st.session_state['coding_code']}\n```"
        elif pres_category == "💬 Chat tab — current conversation":
            pres_source_text = "\n\n".join(f"{m['role'].upper()}: {m['content']}" for m in _pres_conv["messages"])
        elif pres_category == "📤 Upload file(s)":
            pres_ups = st.file_uploader("Upload file(s), or a .zip of a folder", type=["txt", "md", "pdf", "zip"],
                                       accept_multiple_files=True, key="preson_upload")
            if pres_ups:
                _parts = []
                for _u in pres_ups:
                    _lname = _u.name.lower()
                    try:
                        if _lname.endswith(".pdf"):
                            _txt, _, _ = extract_pdf_text(_u, max_chars=24000)
                            _parts.append(f"--- {_u.name} ---\n{_txt}")
                        elif _lname.endswith(".zip"):
                            with zipfile.ZipFile(_u) as _z:
                                for _zn in _z.namelist():
                                    if _zn.lower().endswith((".txt", ".md")):
                                        _parts.append(f"--- {_zn} ---\n" + _z.read(_zn).decode("utf-8", errors="replace"))
                                    elif _zn.lower().endswith(".pdf"):
                                        try:
                                            _reader = PdfReader(io.BytesIO(_z.read(_zn)))
                                            _ptxt = "\n".join((p.extract_text() or "") for p in _reader.pages[:30])
                                            _parts.append(f"--- {_zn} ---\n{_ptxt[:24000]}")
                                        except Exception:
                                            pass
                        else:
                            _parts.append(f"--- {_u.name} ---\n" + _u.read().decode("utf-8", errors="replace"))
                    except Exception:
                        pass
                pres_source_text = "\n\n".join(_parts)
                pres_default_title = (re.sub(r"\.(zip|pdf|txt|md)$", "", pres_ups[0].name, flags=re.I)
                                     if len(pres_ups) == 1 else "Uploaded files")

        st.divider()
        st.markdown("##### 2. Run Presenton")
        ppt1, ppt2 = st.columns([3, 1])
        with ppt1:
            preson_output_name = st.text_input(
                "Output filename", value=pres_default_title or "presentation",
                key="preson_output_name")
        with ppt2:
            preset_fmt = st.selectbox("Format", ["pptx", "pdf", "html"], key="preson_format")

        if st.button("🚀 Run Presenton", type="primary", disabled=not pres_source_text.strip()):
            log_event(_CURRENT_USER, "presenton_run", detail=pres_category)
            # Write source text to a temp file
            _tmp_src = os.path.join(USER_SANDBOX_DIR, f"_preson_src_{datetime.now().strftime('%Y%m%d_%H%M%S')}.md")
            os.makedirs(os.path.dirname(_tmp_src), exist_ok=True)
            with open(_tmp_src, "w", encoding="utf-8") as _sf:
                _sf.write(f"# {pres_default_title or 'Presentation'}\n\n{pres_source_text}")

            _out_name = preson_output_name.strip() or "presentation"
            _out_file = os.path.join(USER_SANDBOX_DIR, f"_preson_output_{_out_name}.{preset_fmt}")
            _cmd_parts = [pres_presenton_path.strip()]
            if pres_presenton_args.strip():
                _cmd_parts.extend(shlex.split(pres_presenton_args.strip()))
            _cmd_parts.extend(["--input", _tmp_src, "--output", _out_file])

            with st.spinner(f"Running presenton…"):
                try:
                    proc = subprocess.run(
                        _cmd_parts, capture_output=True, text=True,
                        timeout=120, encoding="utf-8", errors="replace")
                    st.session_state["_preson_stdout"] = proc.stdout
                    st.session_state["_preson_stderr"] = proc.stderr
                    st.session_state["_preson_rc"] = proc.returncode
                    st.session_state["_preson_out"] = _out_file if os.path.isfile(_out_file) else None
                except subprocess.TimeoutExpired:
                    st.session_state["_preson_rc"] = -1
                    st.session_state["_preson_stderr"] = "Presenton timed out after 120 seconds."
                    st.session_state["_preson_out"] = None
                except Exception as e:
                    st.session_state["_preson_rc"] = -1
                    st.session_state["_preson_stderr"] = str(e)
                    st.session_state["_preson_out"] = None

        _preson_rc = st.session_state.pop("_preson_rc", None)
        if _preson_rc is not None:
            _preson_out = st.session_state.pop("_preson_out", None)
            _preson_stdout = st.session_state.pop("_preson_stdout", "")
            _preson_stderr = st.session_state.pop("_preson_stderr", "")
            if _preson_rc == 0 and _preson_out:
                st.success(f"✅ Presenton finished — output ready!")
                st.session_state["_preson_last_out"] = _preson_out
                with open(_preson_out, "rb") as _pf:
                    st.download_button(
                        f"💾 Download {os.path.basename(_preson_out)}",
                        data=_pf.read(),
                        file_name=os.path.basename(_preson_out),
                        mime="application/octet-stream",
                        key="preson_dl")
            else:
                st.error(f"Presenton exited with code {_preson_rc}")
                if _preson_stderr:
                    with st.expander("Stderr output"):
                        st.code(_preson_stderr, language="text")
                if _preson_stdout:
                    with st.expander("Stdout output"):
                        st.code(_preson_stdout, language="text")

        # Show output from previous successful runs
        _preson_last_out = st.session_state.get("_preson_last_out")
        if _preson_last_out and os.path.isfile(_preson_last_out):
            with open(_preson_last_out, "rb") as _pf:
                st.download_button(
                    f"💾 Re-download {os.path.basename(_preson_last_out)}",
                    data=_pf.read(),
                    file_name=os.path.basename(_preson_last_out),
                    mime="application/octet-stream",
                    key="preson_dl_prev")
        st.stop()

    # ========================= LLM GENERATE MODE (existing) ===================
    _wp_projects = list_wp_projects()
    _sv_projects = list_survey_projects()
    _crew_projects = list_crew_projects()
    _agent_sessions = list_coding_sessions()
    _cc_sessions = list_cc_sessions()
    _pres_conv = st.session_state.get("conversations", {}).get(st.session_state.get("current_conv_id"))

    pres_categories = ["✍️ Type/paste an idea"]
    if _wp_projects:
        pres_categories.append("✍️ Write Paper project")
    if _sv_projects:
        pres_categories.append("📚 Survey project")
    if _crew_projects:
        pres_categories.append("🔬 Research Crew project")
    if st.session_state.get("sm_summary"):
        pres_categories.append("📖 Summarize Paper tab — last summary")
    if _agent_sessions:
        pres_categories.append("🤖 Coding Agent session")
    if _cc_sessions:
        pres_categories.append("🧭 Claude Code session")
    if st.session_state.get("coding_code", "").strip():
        pres_categories.append("💻 Coding tab — current Quick Edit code")
    if _pres_conv and _pres_conv.get("messages"):
        pres_categories.append("💬 Chat tab — current conversation")
    pres_categories.append("🌐 Search the web")
    pres_categories.append("📚 Search academic papers")
    pres_categories.append("📤 Upload file(s)")

    # Category -> the specific saved items inside it, for the second selector.
    _pres_items_by_category = {
        "✍️ Write Paper project": _wp_projects,
        "📚 Survey project": _sv_projects,
        "🔬 Research Crew project": _crew_projects,
        "🤖 Coding Agent session": _agent_sessions,
        "🧭 Claude Code session": _cc_sessions,
    }

    scol1, scol2 = st.columns(2)
    with scol1:
        pres_category = st.selectbox("Source", pres_categories, key="pres_category")
    pres_item = None
    if pres_category in _pres_items_by_category:
        with scol2:
            pres_item = st.selectbox("Which one", _pres_items_by_category[pres_category], key="pres_item")

    pres_source_text, pres_default_title = "", ""
    if pres_category == "✍️ Type/paste an idea":
        pres_source_text = st.text_area(
            "Describe the idea for the presentation:", height=160, key="pres_manual_idea",
            placeholder="e.g. 'A wearable ECG patch for continuous heart monitoring using human "
                        "body communication instead of Bluetooth, aimed at hospital patients — "
                        "cover the motivation, how it works, and results so far.'")

    elif pres_category == "✍️ Write Paper project" and pres_item:
        payload, err = load_wp_project(pres_item)
        wp_data = (payload or {}).get("wp_paper") or {}
        pres_source_text = wp_data.get("md") or (payload or {}).get("wp_idea", "")
        pres_default_title = wp_data.get("title") or pres_item
        st.caption(f"Using **{pres_item}** ({len(pres_source_text):,} characters).")

    elif pres_category == "📚 Survey project" and pres_item:
        payload, err = load_survey_project(pres_item)
        sv_res = (payload or {}).get("survey_result") or {}
        pres_source_text = sv_res.get("md", "")
        pres_default_title = pres_item
        st.caption(f"Using **{pres_item}** ({len(pres_source_text):,} characters).")

    elif pres_category == "🔬 Research Crew project" and pres_item:
        payload, err = load_crew_project(pres_item)
        res_data = (payload or {}).get("research_result") or {}
        pres_source_text = res_data.get("result", "")
        pres_default_title = pres_item
        st.caption(f"Using **{pres_item}** ({len(pres_source_text):,} characters).")

    elif pres_category == "📖 Summarize Paper tab — last summary":
        pres_source_text = st.session_state["sm_summary"]
        pres_default_title = re.sub(r"\.pdf$", "", st.session_state.get("sm_file_name", "Paper"), flags=re.I)
        st.caption(f"Using the summary of **{pres_default_title}** ({len(pres_source_text):,} characters).")

    elif pres_category == "🤖 Coding Agent session" and pres_item:
        payload, err = load_coding_session(pres_item)
        transcript = (payload or {}).get("transcript", [])
        pres_source_text = "\n\n".join(f"{t.get('type', '').upper()}: {t.get('text', '')}"
                                       for t in transcript if t.get("text"))
        pres_default_title = pres_item
        st.caption(f"Using **{pres_item}** ({len(transcript)} transcript entries).")

    elif pres_category == "🧭 Claude Code session" and pres_item:
        payload, err = load_cc_session(pres_item)
        turns = (payload or {}).get("turns", [])
        pres_source_text = "\n\n".join(f"TASK: {t['task']}\nRESULT:\n{t['output']}" for t in turns)
        pres_default_title = pres_item
        st.caption(f"Using **{pres_item}** ({len(turns)} turn(s)).")

    elif pres_category == "💻 Coding tab — current Quick Edit code":
        pres_source_text = f"```{st.session_state.get('coding_lang', '')}\n{st.session_state['coding_code']}\n```"
        st.caption(f"Using the current code ({len(st.session_state['coding_code']):,} characters).")

    elif pres_category == "💬 Chat tab — current conversation":
        pres_source_text = "\n\n".join(f"{m['role'].upper()}: {m['content']}" for m in _pres_conv["messages"])
        st.caption(f"Using {len(_pres_conv['messages'])} message(s) from the current chat.")

    elif pres_category == "🌐 Search the web":
        pres_web_query = st.text_input("Search query", key="pres_web_query")
        if pres_web_query.strip():
            with st.spinner("Searching the web…"):
                _wres, _werr = web_search(pres_web_query.strip(), max_results=6, read_pages=True, chars_per_page=2500)
            if _wres:
                pres_source_text = results_to_context(_wres, use_pages=True)
                pres_default_title = pres_web_query.strip()
                st.caption(f"Found {len(_wres)} web result(s), {len(pres_source_text):,} characters.")
            else:
                st.warning(_werr or "No results found.")

    elif pres_category == "📚 Search academic papers":
        pres_acad_query = st.text_input("Search query", key="pres_academic_query")
        if pres_acad_query.strip():
            with st.spinner("Searching arXiv and OpenAlex…"):
                _found = []
                try:
                    _found += search_arxiv_meta(pres_acad_query.strip(), n=6, sort="relevance")
                except Exception:
                    pass
                try:
                    _found += search_openalex(pres_acad_query.strip(), n=6)
                except Exception:
                    pass
            if _found:
                pres_source_text = papers_to_context_block(_found)
                pres_default_title = pres_acad_query.strip()
                st.caption(f"Found {len(_found)} paper(s), {len(pres_source_text):,} characters.")
            else:
                st.warning("No papers found for that query.")

    elif pres_category == "📤 Upload file(s)":
        # accepts multiple files, and .zip as a stand-in for a folder (browsers don't
        # expose a folder picker to a plain file input — zip the folder instead)
        pres_ups = st.file_uploader("Upload file(s), or a .zip of a folder", type=["txt", "md", "pdf", "zip"],
                                    accept_multiple_files=True, key="pres_upload")
        if pres_ups:
            _parts = []
            for _u in pres_ups:
                _lname = _u.name.lower()
                try:
                    if _lname.endswith(".pdf"):
                        _txt, _, _ = extract_pdf_text(_u, max_chars=24000)
                        _parts.append(f"--- {_u.name} ---\n{_txt}")
                    elif _lname.endswith(".zip"):
                        with zipfile.ZipFile(_u) as _z:
                            for _zn in _z.namelist():
                                _zln = _zn.lower()
                                if _zln.endswith((".txt", ".md")):
                                    _parts.append(f"--- {_zn} ---\n" + _z.read(_zn).decode("utf-8", errors="replace"))
                                elif _zln.endswith(".pdf"):
                                    try:
                                        _reader = PdfReader(io.BytesIO(_z.read(_zn)))
                                        _ptxt = "\n".join((p.extract_text() or "") for p in _reader.pages[:30])
                                        _parts.append(f"--- {_zn} ---\n{_ptxt[:24000]}")
                                    except Exception:
                                        pass
                    else:
                        _parts.append(f"--- {_u.name} ---\n" + _u.read().decode("utf-8", errors="replace"))
                except Exception as _e:
                    st.warning(f"Couldn't read {_u.name}: {_e}")
            pres_source_text = "\n\n".join(_parts)
            pres_default_title = (re.sub(r"\.(zip|pdf|txt|md)$", "", pres_ups[0].name, flags=re.I)
                                  if len(pres_ups) == 1 else "Uploaded files")
            st.caption(f"Read **{len(pres_source_text):,}** characters from {len(pres_ups)} file(s).")

    st.markdown("##### 2. Shape")
    pc1, pc2 = st.columns([1, 2])
    with pc1:
        pres_n_slides = st.slider("Content slides", 3, 25, 8, key="pres_n_slides")
    with pc2:
        pres_guidance = st.text_input(
            "Extra guidance (optional)", key="pres_guidance",
            placeholder="e.g. 'technical audience', 'keep it high-level', 'emphasize results'")

    if st.button("🧠 Generate Outline", type="primary", disabled=not pres_source_text.strip()):
        pres_source_choice = f"{pres_category} / {pres_item}" if pres_item else pres_category
        log_event(_CURRENT_USER, "presentation_outline", detail=pres_source_choice)
        with st.spinner("Thinking through how to structure this…"):
            pres_outline, pres_raw = generate_slide_outline(
                pres_source_text, pres_n_slides, global_model_string,
                api_base=selected_api_base, api_key=selected_api_key, num_ctx=num_ctx,
                guidance=pres_guidance)
        if pres_outline and pres_outline.get("slides"):
            if not pres_outline.get("title") and pres_default_title:
                pres_outline["title"] = pres_default_title
            st.session_state.pres_outline = pres_outline
            st.session_state.pop("pres_pptx_bytes", None)
            add_memory_fact(_CURRENT_USER, f"Made a presentation about: {pres_outline.get('title', 'a topic')}")
        else:
            st.error("Couldn't parse an outline from the model's reply — try again, or a different model.")
            with st.expander("Raw reply"):
                st.code(pres_raw or "(empty)")

    outline = st.session_state.get("pres_outline")
    if outline:
        st.divider()
        st.markdown("##### 3. Review & edit the outline")
        outline["title"] = st.text_input("Deck title", value=outline.get("title", ""), key="pres_title_edit")
        outline["subtitle"] = st.text_input("Subtitle (optional)", value=outline.get("subtitle", ""),
                                            key="pres_subtitle_edit")

        pres_remove_idx = None
        for i, sl in enumerate(outline.get("slides", [])):
            with st.expander(f"Slide {i + 1}: {sl.get('title', '') or '(untitled)'}", expanded=False):
                sl["title"] = st.text_input("Title", value=sl.get("title", ""), key=f"pres_slide_title_{i}")
                bullets_text = st.text_area(
                    "Bullets (one per line)", value="\n".join(sl.get("bullets", []) or []),
                    height=120, key=f"pres_slide_bullets_{i}")
                sl["bullets"] = [b.strip() for b in bullets_text.split("\n") if b.strip()]
                sl["notes"] = st.text_area("Speaker notes (optional)", value=sl.get("notes", "") or "",
                                           height=60, key=f"pres_slide_notes_{i}")
                if st.button("🗑️ Remove slide", key=f"pres_slide_remove_{i}"):
                    pres_remove_idx = i
        if pres_remove_idx is not None:
            outline["slides"].pop(pres_remove_idx)
            st.rerun()

        if st.button("➕ Add blank slide"):
            outline.setdefault("slides", []).append({"title": "New slide", "bullets": [""], "notes": ""})
            st.rerun()

        st.divider()
        st.markdown("##### 4. Build the file")
        pres_color = st.color_picker("Accent color", value="#1F4E79", key="pres_color")
        if st.button("🎬 Build .pptx", type="primary"):
            log_event(_CURRENT_USER, "presentation_build", detail=f"{len(outline.get('slides', []))} slides")
            pres_pptx_bytes = build_pptx_from_outline(outline, theme_color=pres_color.lstrip("#"))
            if pres_pptx_bytes is None:
                st.error("python-pptx isn't installed. Run: `pip install python-pptx`, then try again.")
            else:
                st.session_state["pres_pptx_bytes"] = pres_pptx_bytes
                st.session_state["pres_pptx_name"] = re.sub(
                    r"[^A-Za-z0-9_ -]+", "", outline.get("title", "presentation"))[:60] or "presentation"
                st.success(f"Built a {len(outline.get('slides', [])) + 1}-slide deck.")

        if st.session_state.get("pres_pptx_bytes"):
            st.download_button(
                "💾 Download Presentation (.pptx)",
                data=st.session_state["pres_pptx_bytes"],
                file_name=f"{st.session_state.get('pres_pptx_name', 'presentation')}.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")

# =============================== ADMIN TAB (admin accounts only) =============
if _CURRENT_IS_ADMIN:
    with tab_admin:
        st.subheader("🛡️ Admin Dashboard")
        st.caption("Only visible to admin accounts. Admin accounts also bypass the per-user "
                  "folder sandbox in the Agent/Claude Code tabs — they can point at any folder "
                  "on this machine.")

        st.markdown("##### 👥 Users")
        all_users = list_all_users()
        user_rows = []
        for uname_key, rec in sorted(all_users.items()):
            folder = os.path.join(BASE_DIR, "user_data", _safe_username_for_path(rec["username"]))
            size_mb = dir_size_bytes(folder) / (1024 * 1024) if os.path.isdir(folder) else 0.0
            user_rows.append({
                "Username": rec["username"], "Email": rec["email"],
                "Admin": "🛡️" if rec.get("is_admin") else "",
                "Created": rec.get("created_at", ""),
                "Last login": rec.get("last_login") or "never",
                "Data used (MB)": round(size_mb, 2),
            })
        st.dataframe(user_rows, use_container_width=True, hide_index=True)

        st.markdown("###### Promote / demote admin")
        pa1, pa2, pa3 = st.columns([2, 1, 1])
        with pa1:
            admin_target_user = st.selectbox(
                "User", [r["username"] for r in all_users.values()], key="admin_target_user")
        with pa2:
            if st.button("⬆️ Make admin", key="admin_promote", use_container_width=True):
                set_admin_flag(admin_target_user, True)
                st.rerun()
        with pa3:
            if st.button("⬇️ Remove admin", key="admin_demote", use_container_width=True,
                        disabled=(admin_target_user == _CURRENT_USER)):
                set_admin_flag(admin_target_user, False)
                st.rerun()
        if admin_target_user == _CURRENT_USER:
            st.caption("(Can't remove your own admin access from here — have another admin do it.)")

        st.divider()
        st.markdown("##### 📊 Tool usage")
        admin_events = read_events(limit=10000)
        if not admin_events:
            st.caption("No activity logged yet.")
        else:
            from collections import Counter
            by_tool = Counter(e.get("event", "?") for e in admin_events)
            by_user = Counter(e.get("user", "?") for e in admin_events)

            ec1, ec2 = st.columns(2)
            with ec1:
                st.markdown("**By tool**")
                st.dataframe([{"Tool": k, "Uses": v} for k, v in by_tool.most_common()],
                            use_container_width=True, hide_index=True)
            with ec2:
                st.markdown("**By user**")
                st.dataframe([{"User": k, "Uses": v} for k, v in by_user.most_common()],
                            use_container_width=True, hide_index=True)

            st.markdown("##### 🕒 Recent activity")
            admin_user_filter = st.selectbox(
                "Filter by user", ["(all)"] + sorted(by_user.keys()), key="admin_activity_filter")
            recent = list(reversed(admin_events))
            if admin_user_filter != "(all)":
                recent = [e for e in recent if e.get("user") == admin_user_filter]
            recent = recent[:300]
            st.dataframe(
                [{"Time": e.get("ts", ""), "User": e.get("user", ""), "Event": e.get("event", ""),
                  "Detail": e.get("detail", "")} for e in recent],
                use_container_width=True, hide_index=True)

# =============================== CHAT TAB (3-panel: conversations | chat | context) ===
with tab_chat:
    # ---- Panel toggle state --------------------------------------------------
    if "chat_show_convos" not in st.session_state:
        st.session_state.chat_show_convos = True
    if "chat_show_context" not in st.session_state:
        st.session_state.chat_show_context = True

    # ---- Conversation store: supports multiple named chats, like ChatGPT's
    # sidebar chat list. Lives only for this browser session (not persisted to
    # disk). -------------------------------------------------------------------
    # ---- Load conversations from disk on first access (survives browser close) --
    if "conversations" not in st.session_state:
        _loaded, _active = load_all_chat_conversations()
        if _loaded:
            st.session_state.conversations = _loaded
            st.session_state.current_conv_id = _active or list(_loaded.keys())[0]
        else:
            # No saved conversations — start fresh
            _first_id = f"conv_{datetime.now().strftime('%H%M%S%f')}"
            st.session_state.conversations = {_first_id: {"title": "New chat", "messages": []}}
            st.session_state.current_conv_id = _first_id
            save_chat_conversation(_first_id, st.session_state.conversations[_first_id])

    conv_ids = list(st.session_state.conversations.keys())
    if st.session_state.current_conv_id not in st.session_state.conversations:
        st.session_state.current_conv_id = conv_ids[0]

    # ---- Callbacks for conversation management (runs before next script body,
    # so state is already updated when the body renders — avoids the double-rerun
    # dance that inline st.rerun() inside nested st.columns requires). ----------
    def _switch_conv(cid):
        st.session_state.current_conv_id = cid
        _save_chat_index(cid)

    def _new_conv():
        new_id = f"conv_{datetime.now().strftime('%H%M%S%f')}"
        st.session_state.conversations[new_id] = {"title": "New chat", "messages": []}
        st.session_state.current_conv_id = new_id
        save_chat_conversation(new_id, st.session_state.conversations[new_id])

    def _delete_conv(cid):
        convs = st.session_state.conversations
        del convs[cid]
        delete_chat_conversation(cid)
        if st.session_state.current_conv_id == cid:
            st.session_state.current_conv_id = next(iter(convs.keys()), None)
            if st.session_state.current_conv_id:
                _save_chat_index(st.session_state.current_conv_id)

    def _rename_conv():
        new_name = st.session_state.get("chat_rename_input", "").strip()
        cid = st.session_state.current_conv_id
        if new_name and new_name != st.session_state.conversations[cid]["title"]:
            st.session_state.conversations[cid]["title"] = new_name
            save_chat_conversation(cid, st.session_state.conversations[cid])

    # ---- Top toolbar — compact picker + panel toggles ------------------------
    with st.container():
        tb1, tb2, tb3, tb4 = st.columns([5, 0.7, 0.7, 1.2])
        with tb1:
            chosen = st.selectbox(
                "Conversation", conv_ids,
                index=conv_ids.index(st.session_state.current_conv_id),
                format_func=lambda cid: st.session_state.conversations[cid]["title"],
                label_visibility="collapsed", key="conv_picker")
            if chosen != st.session_state.current_conv_id:
                st.session_state.current_conv_id = chosen
                st.rerun()
        with tb2:
            st.session_state.chat_show_convos = st.toggle(
                "💬", value=st.session_state.chat_show_convos, key="chat_toggle_convos",
                help="Show/hide conversation list panel")
        with tb3:
            st.session_state.chat_show_context = st.toggle(
                "📋", value=st.session_state.chat_show_context, key="chat_toggle_context",
                help="Show/hide context & info panel")
        with tb4:
            st.button("➕ New", key="chat_new_top_btn", use_container_width=True,
                      on_click=_new_conv)

    conv = st.session_state.conversations[st.session_state.current_conv_id]
    messages = conv["messages"]  # mutated in place; persists in session_state automatically

    # ---- CSS: sticky side panels + scrollable chat (matches Claude Code tab) --
    st.markdown("""
    <style>
    .chat-sticky-left, .chat-sticky-right {
        position: sticky;
        top: 0.5rem;
        max-height: 85vh;
        overflow-y: auto;
        padding-right: 0.3rem;
    }
    .chat-scroll {
        max-height: 55vh;
        overflow-y: auto;
        padding: 0.5rem;
        margin-bottom: 0.5rem;
        border: 1px solid rgba(128,128,128,0.12);
        border-radius: 8px;
    }
    .chat-input-area {
        position: sticky;
        bottom: 0;
        padding-top: 0.3rem;
        z-index: 10;
    }
    </style>
    """, unsafe_allow_html=True)

    st.divider()

    # ---- 3-panel layout: conversations | chat | context ---------------------
    _show_left = st.session_state.chat_show_convos
    _show_right = st.session_state.chat_show_context

    if _show_left and _show_right:
        left_col, mid_col, right_col = st.columns([1.8, 4, 2.2])
    elif _show_left:
        left_col, mid_col, right_col = st.columns([2, 5, 0.02])
    elif _show_right:
        left_col, mid_col, right_col = st.columns([0.02, 5, 2.2])
    else:
        left_col, mid_col = st.columns([0.02, 5])
        right_col = None

    # ── LEFT PANEL: Conversations ─────────────────────────────────────────
    if _show_left:
        with left_col:
            st.markdown('<div class="chat-sticky-left">', unsafe_allow_html=True)
            st.caption("##### 💬 Conversations")

            # New conversation button
            st.button("➕ New chat", key="chat_new_btn", use_container_width=True,
                      on_click=_new_conv)

            # Gather current state for the loop (avoids reading stale session_state
            # mid-render when on_click callbacks have already mutated it)
            _active_id = st.session_state.current_conv_id
            _convos = dict(st.session_state.conversations)  # shallow snapshot

            if not _convos:
                st.caption("No conversations yet.")
            else:
                for cid, cdata in _convos.items():
                    is_active = (cid == _active_id)
                    ctitle = cdata["title"]
                    _btn_label = ("🟢 " if is_active else "   ") + ctitle[:35]
                    crow1, crow2 = st.columns([5, 1])
                    with crow1:
                        st.button(
                            _btn_label, key=f"chat_conv_{cid}",
                            use_container_width=True,
                            help=f"Switch to '{ctitle}'",
                            type="primary" if is_active else "secondary",
                            on_click=_switch_conv, args=(cid,),
                            disabled=is_active)
                    with crow2:
                        st.button(
                            "🗑", key=f"chat_del_{cid}",
                            help=f"Delete '{ctitle}'",
                            disabled=len(conv_ids) <= 1,
                            on_click=_delete_conv, args=(cid,))

                # Rename active conversation
                if _active_id and _active_id in _convos:
                    st.divider()
                    st.text_input(
                        "Rename", value=_convos[_active_id]["title"],
                        key="chat_rename_input", label_visibility="collapsed")
                    _cur_title = _convos[_active_id]["title"]
                    _new_name = st.session_state.get("chat_rename_input", "").strip()
                    st.button("✏️ Rename", key="chat_rename_btn", use_container_width=True,
                              disabled=(_new_name == _cur_title or not _new_name),
                              on_click=_rename_conv)

            st.markdown('</div>', unsafe_allow_html=True)  # close chat-sticky-left

    # ── CENTER: Chat ──────────────────────────────────────────────────────
    with mid_col:
        # Settings expander at top
        with st.expander("⚙️ Settings & Attachments", expanded=False):
            c1, c2 = st.columns([1, 1])
            with c1:
                chat_temperature = st.slider("Temperature", 0.0, 1.5, 0.7, 0.05, key="chat_temp")
                tts_on = st.toggle("🔊 Speak replies", value=False)
                tts_voice_label = None
                if tts_on:
                    tts_voice_label = st.selectbox("Voice", list(TTS_VOICES.keys()), index=0, key="chat_tts_voice")
            with c2:
                chat_files = st.file_uploader("Attach files (PDF, TXT)", type=["pdf", "txt"], accept_multiple_files=True)
                if chat_files:
                    st.caption(f"📎 {len(chat_files)} file(s) attached")

            system_prompt = st.text_area("System prompt", "You are a helpful research assistant.", height=68)

            web_on = st.toggle("🌐 Enable Web Search", value=False)
            if web_on:
                web_n = st.number_input("Search Results", 3, 10, 5)
                web_read = st.checkbox("Read full pages")
                web_max_rounds = st.slider(
                    "Max search rounds", 1, 5, 3,
                    help="If the first search isn't enough, the model can search again with a "
                         "different query — up to this many times — before it answers.")

            st.markdown("##### 🧠 Memory")
            mem_enabled = st.toggle(
                "Remember things about me across sessions", value=True, key="chat_mem_enabled",
                help="Like ChatGPT's memory — the model notices durable facts (your name, role, "
                     "preferences, ongoing projects) as you chat and recalls them in future "
                     "sessions, even after the app restarts. This is the same memory every other "
                     "tool in this app draws on and adds to — it's tied to your account, not just "
                     "this chat.")
            mem_profile = _CURRENT_USER
            if mem_enabled:
                mem_facts = load_memory(mem_profile)
                if mem_facts:
                    st.caption(f"{len(mem_facts)} thing(s) remembered:")
                    for i, f in enumerate(mem_facts):
                        mc1, mc2 = st.columns([6, 1])
                        with mc1:
                            st.caption(f"• {f['fact']}")
                        with mc2:
                            if st.button("🗑️", key=f"mem_del_{i}"):
                                delete_memory_fact(mem_profile, i)
                                st.rerun()
                    if st.button("🧹 Forget everything", key="mem_clear_all"):
                        clear_memory(mem_profile)
                        st.rerun()
                else:
                    st.caption("Nothing remembered yet — it builds up automatically as you use "
                              "any tool in this app, not just chat.")

        # Scrollable message history
        st.markdown('<div class="chat-scroll">', unsafe_allow_html=True)
        for m in messages:
            avatar = "🧑" if m["role"] == "user" else "✨"
            with st.chat_message(m["role"], avatar=avatar):
                st.markdown(m["content"])
                if m.get("sources"):
                    render_source_cards(m["sources"])

        if not messages:
            st.caption("👋 Ask me anything — or turn on web search in ⚙️ Settings for current info.")
        st.markdown('</div>', unsafe_allow_html=True)  # close chat-scroll

        # Follow-up question chips for the last answer, Perplexity-style
        followups = conv.get("followups") or []
        clicked_followup = None
        if followups and messages and messages[-1]["role"] == "assistant":
            st.markdown("**Related**")
            fcols = st.columns(len(followups))
            for i, q in enumerate(followups):
                with fcols[i]:
                    if st.button(q, key=f"followup_{st.session_state.current_conv_id}_{len(messages)}_{i}",
                                use_container_width=True):
                        clicked_followup = q

        # Input area
        st.markdown('<div class="chat-input-area">', unsafe_allow_html=True)

        # ---- Sending a message — shared by both the text box and follow-up chips --
        def send_message(um):
            # Auto-title the conversation from the first message, like ChatGPT does
            if not messages:
                conv["title"] = (um.strip()[:40] + "…") if len(um.strip()) > 40 else um.strip()
            conv["followups"] = []  # clear stale suggestions now that a new turn is starting

            # Process file context
            file_context = ""
            if chat_files:
                file_context = "\n\n[Attached Context]\n"
                for f in chat_files:
                    try:
                        if f.type == "application/pdf":
                            reader = PdfReader(f)
                            text = "\n".join([page.extract_text() or "" for page in reader.pages])
                            file_context += f"File: {f.name}\n{text}\n"
                        else:
                            file_context += f"File: {f.name}\n{f.read().decode('utf-8', errors='ignore')}\n"
                    except Exception as e:
                        file_context += f"File: {f.name} [Error: {e}]\n"

            # Display user message
            messages.append({"role": "user", "content": um})
            with st.chat_message("user", avatar="🧑"):
                st.markdown(um)

            # Process Web Search — plan → search → judge → (maybe) search again, then flatten
            # every round into ONE deduped, globally-numbered source list.
            flat_sources = []
            if web_on:
                with st.status("🌐 Researching…", expanded=True) as status_box:
                    def _status(msg):
                        status_box.write(msg)

                    _raw_context, search_rounds = run_iterative_search(
                        um, messages[:-1], global_model_string,
                        api_base=selected_api_base, api_key=selected_api_key, num_ctx=num_ctx,
                        max_results=int(web_n), read_pages=web_read, max_rounds=int(web_max_rounds),
                        status_cb=_status)

                    flat_sources = flatten_search_results(search_rounds)
                    n_searches = sum(1 for r in search_rounds if r["type"] == "search")
                    if n_searches:
                        status_box.update(label=f"🌐 Searched {n_searches} round(s) — {len(flat_sources)} source(s)",
                                         state="complete")
                    else:
                        status_box.update(label="🌐 Web search skipped — not needed for this question", state="complete")

            # Prepare messages — instruct inline [n] citations tied to the numbered source list
            # (Memory is NOT manually appended here — stream_chat() injects the account's
            # USER_MEMORY_CONTEXT into every call automatically.)
            sys_txt = system_prompt.strip()
            if flat_sources:
                cited_context = build_cited_context(flat_sources, use_pages=web_read)
                sys_txt += (
                    "\n\nYou have the following numbered web sources available. Cite them inline "
                    "using [n] immediately after any claim they support. Only cite when a source "
                    "directly supports the specific claim — don't force citations onto general "
                    "knowledge.\n\n" + cited_context
                )
            msgs = [{"role": "system", "content": sys_txt}] + messages[:-1]
            msgs.append({"role": "user", "content": um + file_context})

            # Stream response
            with st.chat_message("assistant", avatar="✨"):
                try:
                    response = st.write_stream(stream_chat(msgs, global_model_string,
                                                         api_base=selected_api_base,
                                                         api_key=selected_api_key,
                                                         temperature=chat_temperature))
                    if flat_sources:
                        render_source_cards(flat_sources)
                    messages.append({"role": "assistant", "content": response, "sources": flat_sources})

                    # Optional TTS
                    if tts_on and tts_voice_label:
                        audio, _ = tts_to_mp3(response, voice=TTS_VOICES[tts_voice_label])
                        if audio:
                            st.audio(audio, format="audio/mpeg", autoplay=True)

                    # Suggest related follow-up questions, Perplexity-style
                    with st.spinner("Thinking of related questions…"):
                        conv["followups"] = suggest_followups(
                            um, response, global_model_string,
                            api_base=selected_api_base, api_key=selected_api_key, num_ctx=num_ctx)

                    # Update long-term memory, ChatGPT-style
                    if mem_enabled:
                        existing = [f["fact"] for f in load_memory(mem_profile)]
                        new_fact = extract_memory_fact(
                            um, response, existing, global_model_string,
                            api_base=selected_api_base, api_key=selected_api_key, num_ctx=num_ctx)
                        if new_fact:
                            add_memory_fact(mem_profile, new_fact)
                            st.caption(f"🧠 Remembered: {new_fact}")

                    # Persist the conversation to disk so it survives browser close
                    save_chat_conversation(st.session_state.current_conv_id, conv)
                except Exception as e:
                    st.error(f"Error: {e}")
                    messages.pop()  # drop the empty/failed turn so retrying doesn't duplicate it

        # ---- Input -----------------------------------------------------------
        um_typed = st.chat_input("Message the AI...")
        um = clicked_followup or um_typed
        if um:
            log_event(_CURRENT_USER, "chat_message")
            send_message(um)
            st.rerun()

        st.markdown('</div>', unsafe_allow_html=True)  # close chat-input-area

    # ── RIGHT PANEL: Context & Info ───────────────────────────────────────
    if _show_right and right_col is not None:
        with right_col:
            st.markdown('<div class="chat-sticky-right">', unsafe_allow_html=True)
            st.caption("##### 📋 Context")

            # Attached files summary
            if chat_files:
                st.caption(f"**📎 {len(chat_files)} attached file(s):**")
                for f in chat_files:
                    st.caption(f"• {f.name}")
            else:
                st.caption("No files attached.")

            st.divider()

            # Current conversation stats
            st.caption(f"**💬 Messages:** {len(messages)}")
            if messages:
                _last_role = messages[-1]["role"]
                st.caption(f"**Last:** {_last_role}")
                _total_chars = sum(len(m.get("content", "")) for m in messages)
                st.caption(f"**Length:** {_total_chars:,} chars")
                # Source count across all messages
                _total_sources = sum(len(m.get("sources", [])) for m in messages)
                if _total_sources:
                    st.caption(f"**🔗 Sources:** {_total_sources} total")

            st.divider()

            # Active settings at a glance
            st.caption(f"**🌡️ Temperature:** {chat_temperature}")
            if tts_on:
                st.caption(f"**🔊 TTS:** {tts_voice_label or 'on'}")
            st.caption(f"**🌐 Web search:** {'on' if web_on else 'off'}")
            if web_on:
                st.caption(f"**🔍 Results/round:** {web_n}")
                st.caption(f"**📖 Read pages:** {'yes' if web_read else 'no'}")
                st.caption(f"**🔄 Max rounds:** {web_max_rounds}")

            st.divider()

            # Memory stats
            if mem_enabled:
                mem_facts = load_memory(mem_profile)
                st.caption(f"**🧠 Memory:** {len(mem_facts) or 'no'} fact(s) stored")
            else:
                st.caption("**🧠 Memory:** off")

            # Quick actions
            st.divider()
            st.button("🗑️ Delete this chat", key="chat_del_right_btn", use_container_width=True,
                      disabled=len(conv_ids) <= 1,
                      on_click=_delete_conv, args=(st.session_state.current_conv_id,))

            st.markdown('</div>', unsafe_allow_html=True)  # close chat-sticky-right
